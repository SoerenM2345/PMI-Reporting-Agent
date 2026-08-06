"""Populated PPTX: SteerCo deck, two workstream one-pagers, integration roadmap."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import case as C

L_TITLE = 1
L_ONLY = 27
L_SUB = 28
L_1COL = 29
L_DIV = 19


def rgb(h):
    return RGBColor.from_string(h)


CASECODE = (f"{C.NEWCO}  |  {C.OFFICE} ({C.OFFICE_ABBR})\n"
            f"{C.en(C.TODAY)}  |  Day {C.DAYS_AFTER_DAY1} after Day 1  |  "
            f"Strictly confidential")


def deck():
    """Open the template, drop its sample slide, and replace the master's
    CaseCode placeholder so no template boilerplate reaches the output."""
    p = Presentation(C.TEMPLATE)
    lst = p.slides._sldIdLst
    for s in list(lst):
        p.part.drop_rel(s.rId)
        lst.remove(s)
    for sh in p.slide_masters[0].shapes:
        if sh.name == "Copyright" and sh.has_text_frame:
            for r in sh.text_frame.paragraphs[0].runs:
                r.text = r.text.replace("2026", "2016")
        if sh.name == "CaseCode" and sh.has_text_frame:
            tf = sh.text_frame
            for para in list(tf.paragraphs)[1:]:
                para._p.getparent().remove(para._p)
            p0 = tf.paragraphs[0]
            for r in list(p0.runs):
                r._r.getparent().remove(r._r)
            for i, line in enumerate(CASECODE.split("\n")):
                para = p0 if i == 0 else tf.add_paragraph()
                para.alignment = PP_ALIGN.RIGHT
                run = para.add_run()
                run.text = line
                run.font.size = Pt(7.5)
                run.font.color.rgb = rgb(C.D_GREY)
    return p


def add(p, i):
    return p.slides.add_slide(p.slide_masters[0].slide_layouts[i])


def ph_text(slide, idx, text, size=None, bold=None):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            r = ph.text_frame.paragraphs[0].add_run()
            r.text = text
            if size:
                r.font.size = Pt(size)
            if bold is not None:
                r.font.bold = bold
            return ph


def title(slide, t, sub=None):
    ph_text(slide, 0, t, bold=True)
    if sub:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 13:
                ph.text_frame.word_wrap = True
                r = ph.text_frame.paragraphs[0].add_run()
                r.text = sub
                r.font.size = Pt(13)
    return slide


def tb(slide, l, t, w, h, text, size=12, bold=False, color=C.D_BLACK,
       align=PP_ALIGN.LEFT, fill=None):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    f = box.text_frame
    f.word_wrap = True
    f.margin_left = f.margin_right = Inches(0.06)
    f.margin_top = f.margin_bottom = Inches(0.03)
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        para = f.paragraphs[0] if i == 0 else f.add_paragraph()
        para.alignment = align
        r = para.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = rgb(color)
    if fill:
        box.fill.solid()
        box.fill.fore_color.rgb = rgb(fill)
        box.line.fill.background()
    return box


def bullets(slide, idx, items, size=12):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            f = ph.text_frame
            f.word_wrap = True
            for i, (lvl, txt, bold) in enumerate(items):
                para = f.paragraphs[0] if i == 0 else f.add_paragraph()
                para.level = lvl
                r = para.add_run()
                r.text = txt
                r.font.size = Pt(size)
                r.font.bold = bold
            return ph


def table(slide, l, t, w, h, headers, rows, widths, hdr=9, body=8.5, colcolor=None):
    gf = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(l), Inches(t),
                                Inches(w), Inches(h))
    tbl = gf.table
    total = sum(widths)
    for i, ww in enumerate(widths):
        tbl.columns[i].width = Inches(w * ww / total)
    for j, htxt in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = ""
        c.fill.solid()
        c.fill.fore_color.rgb = rgb(C.D_DARK)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = c.margin_right = Inches(0.04)
        r = c.text_frame.paragraphs[0].add_run()
        r.text = htxt
        r.font.size = Pt(hdr)
        r.font.bold = True
        r.font.color.rgb = rgb("FFFFFF")
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = ""
            c.fill.solid()
            c.fill.fore_color.rgb = rgb("FFFFFF" if i % 2 else "F7F9F3")
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = c.margin_right = Inches(0.04)
            para = c.text_frame.paragraphs[0]
            r = para.add_run()
            r.text = str(val)
            r.font.size = Pt(body)
            col = C.D_BLACK
            if colcolor and (i - 1, j) in colcolor:
                col = colcolor[(i - 1, j)]
                r.font.bold = True
            r.font.color.rgb = rgb(col)
    return tbl


def footer(slide, extra=""):
    """The master already carries the standing footer; this adds only a source line."""
    if extra:
        tb(slide, 0.5, 6.92, 8.0, 0.25, extra, size=7.5, color=C.D_GREY)


# =====================================================================
# 1. SteerCo deck, session 02
# =====================================================================
def steerco():
    p = deck()

    s = add(p, L_TITLE)
    ph_text(s, 0, f"{C.OFFICE}\nSteering Committee, Session 02", size=26, bold=True)
    ph_text(s, 10, f"{C.ACQUIRER} and {C.TARGET}   |   {C.WEEK_LABEL}   |   "
                   f"Day {C.DAYS_AFTER_DAY1} after Day 1   |   {C.en(C.STEERCO_02)}", size=10)

    # --- executive summary
    s = add(p, L_1COL)
    title(s, "Executive Summary",
          f"Key message: the Day 100 commitment of {C.en(C.DAY100)} still holds, but three "
          f"workstreams now carry a delayed gate milestone and the value case is "
          f"{C.SYN_SECURED_PCT} percent secured against register target.")
    reds = [c for c in C.WS_CODES if C.ws_rag(c) == "Red"]
    bullets(s, 1, [
        (0, f"Overall status {C.OVERALL_RAG}   |   Milestones on track {C.MS_ON_TRACK} of "
            f"{C.MS_TOTAL}   |   Open actions {C.T_OPEN + C.T_INPROG}, of which "
            f"{C.T_OVERDUE} overdue and {C.T_BLOCKED} blocked   |   "
            f"Secured synergy USD {C.SYN_SECURED_DECK} m of USD {C.SYN_TARGET} m", True),
        (0, "", False),
        (0, "What moved this week", True),
        (1, f"Value prioritisation framework scoring published and adopted as the sole basis "
            f"for sequencing, per decision B-01", False),
        (1, f"Interim account ownership list published, closing the duplicate coverage issue "
            f"raised in the first fortnight", False),
        (1, f"Portfolio cross-reference released ahead of the flagship event on "
            f"{C.en(C.DELL_EMC_WORLD)}", False),
        (0, "", False),
        (0, "What is off track", True),
        (1, f"{', '.join(C.ws_full(c) for c in reds)} each carry a delayed gate milestone", False),
        (1, f"M-07 ERP consolidation blueprint has slipped {(C.M07_MAIL - C.M07_ROADMAP).days} "
            f"days; the legacy upgrade at the target and the integration design are competing "
            f"for the same architecture decisions", False),
        (1, f"M-13 works council consultation in DACH has slipped 14 days and is a legal gate, "
            f"not a schedule preference", False),
        (0, "", False),
        (0, "What the Steering Committee is asked to decide today", True),
        (1, "B-07 hold the Day 100 review date and re-phase scope instead of moving the "
            "commitment, page 7", False),
        (1, "Endorse the ERP blueprint recovery plan and the associated scope freeze, page 7", False),
    ], size=11)
    footer(s)

    # --- workstream status
    s = add(p, L_SUB)
    title(s, "Workstream Status Overview",
          f"Key message: the two workstreams that gate everything downstream, "
          f"{C.ws_full('WS2')} and {C.ws_full('WS3')}, are the ones in difficulty.")
    rows, colors = [], {}
    for i, code in enumerate(C.WS_CODES):
        rag = C.ws_rag(code)
        ms = C.ws_milestones(code)
        openms = sum(1 for m in ms if m[8] != "done")
        keyrisk = max(C.ws_risks(code), key=C.sev, default=None)
        rows.append([code, C.WS_NAME[code], C.nm(code).split()[-1] if code in C.PEOPLE else "",
                     rag, C.RAG_LAST_WEEK[code], f"{C.ws_progress(code)}%",
                     openms, C.ws_overdue(code),
                     keyrisk[0] if keyrisk else "-",
                     "Yes" if any(r[11] for r in C.ws_risks(code)) else "No"])
        colors[(i, 3)] = C.RAG_COLOR[rag]
    table(s, 0.5, 1.95, 12.33, 3.2,
          ["WS", "Workstream", "Lead", "RAG now", "RAG last", "Progress",
           "Open milestones", "Overdue or blocked", "Top risk", "At SteerCo"],
          rows, [0.6, 1.9, 1.1, 0.9, 0.9, 0.9, 1.2, 1.4, 0.8, 1.0], colcolor=colors)
    tb(s, 0.5, 5.35, 12.33, 1.5,
       f"Reading: RAG is derived, not asserted. A workstream shows Red when a gate milestone "
       f"is already delayed, Amber when a gate milestone is at risk or a severity {C.ESCALATION_THRESHOLD} "
       f"or higher risk is open, Green otherwise. Three workstreams moved to Red this week; none "
       f"improved. Progress percentages come from the integration tracker, which decision B-02 "
       f"makes the single source of truth. Where a workstream deck disagrees with this table, "
       f"the tracker governs.", size=10, fill=C.D_PALE)
    footer(s)

    # --- milestones
    s = add(p, L_SUB)
    title(s, f"Milestone Status: Day 1 to Day 100",
          f"Key message: {C.MS_DELAYED} of {C.MS_TOTAL} milestones are delayed and the largest "
          f"single slip is {C.MS_MAX_SLIP} days, which the Day 100 date can still absorb only "
          f"if scope is re-phased.")
    rows, colors = [], {}
    slipped = [m for m in C.MILESTONES if m[7] != m[6] or m[8] in ("at_risk", "delayed")]
    others = [m for m in C.MILESTONES if m not in slipped][:5]
    for i, m in enumerate(slipped + others):
        delta = (m[7] - m[6]).days
        rows.append([m[0], m[1][:62], m[3], C.de(m[6]).replace(".2016", ""),
                     C.de(m[7]).replace(".2016", ""),
                     f"+{delta}" if delta else "0",
                     {"done": "Done", "on_track": "On track", "at_risk": "At risk",
                      "delayed": "Delayed"}[m[8]],
                     "Yes" if m[9] else "No"])
        if m[8] == "delayed":
            colors[(i, 6)] = C.RAG_COLOR["Red"]
        elif m[8] == "at_risk":
            colors[(i, 6)] = C.RAG_COLOR["Amber"]
    table(s, 0.5, 1.95, 12.33, 4.3,
          ["ID", "Milestone", "WS", "Baseline", "Forecast", "Delta (d)", "Status", "Gate"],
          rows, [0.6, 5.6, 0.7, 1.0, 1.0, 0.9, 1.0, 0.7], colcolor=colors)
    tb(s, 0.5, 6.35, 12.33, 0.5,
       f"{C.MS_GATE} of {C.MS_TOTAL} milestones are gate relevant. {C.en(C.DELL_EMC_WORLD)} "
       f"is an external fixed date and cannot move.", size=9, color=C.D_GREY)
    footer(s)

    # --- risks
    s = add(p, L_SUB)
    title(s, "Top Risks and Issues",
          f"Key message: {C.R_HIGH} risks sit at or above the escalation threshold of "
          f"severity {C.ESCALATION_THRESHOLD}, and two of them share the same root cause.")
    top = sorted(C.RISKS, key=C.sev, reverse=True)[:7]
    rows, colors = [], {}
    for i, r in enumerate(top):
        srv = C.sev(r)
        rows.append([r[0], r[1][:78], r[3], f"{r[5]}x{r[6]}={srv}", C.band(srv),
                     r[7][:64], C.de(r[8]).replace(".2016", ""), r[10]])
        colors[(i, 4)] = C.RAG_COLOR["Red"] if srv >= 15 else C.RAG_COLOR["Amber"]
    table(s, 0.5, 1.95, 12.33, 3.3,
          ["ID", "Risk", "WS", "L x I", "Band", "Mitigation", "Due", "Trend"],
          rows, [0.6, 4.3, 0.6, 0.9, 0.8, 3.9, 0.8, 0.9], colcolor=colors)
    tb(s, 0.5, 5.45, 6.0, 1.4,
       f"Movement since session 01\n"
       f"Increasing: {', '.join(r[0] for r in C.RISKS if r[10] == 'Increasing')}\n"
       f"Decreasing: {', '.join(r[0] for r in C.RISKS if r[10] == 'Decreasing')}\n"
       f"New this period: none", size=10, fill=C.D_PALE)
    tb(s, 6.84, 5.45, 6.0, 1.4,
       f"Register health\n"
       f"{C.R_TOTAL} risks open, {C.R_HIGH} at or above threshold, "
       f"{C.R_ESCALATED} escalated beyond the workstream.\n"
       f"{C.DEP_UNCONFIRMED_CRITICAL} critical or high dependencies are still not confirmed by "
       f"both leads, contrary to decision B-03.", size=10, fill=C.D_PALE)
    footer(s)

    # --- synergies
    s = add(p, L_SUB)
    title(s, "Value Realisation Status",
          f"Key message: USD {C.SYN_SECURED_DECK} m of USD {C.SYN_TARGET} m is secured, but only "
          f"USD {C.SYN_VALIDATED} m has passed Finance validation, which is the figure decision "
          f"B-06 says we report.")
    rows = []
    for b in C.SYN_BUCKETS:
        items = [x for x in C.SYNERGIES if x[3] == b]
        tgt = sum(x[7] for x in items)
        sec = sum(x[8] for x in items)
        rows.append([b, len(items), f"{tgt:,}", f"{sec:,}",
                     f"{round(100 * sec / tgt)}%", f"{sum(x[10] for x in items):,}",
                     f"{sum(x[11] for x in items):,}"])
    rows.append(["Total", len(C.SYNERGIES), f"{C.SYN_TARGET:,}", f"{C.SYN_SECURED:,}",
                 f"{C.SYN_SECURED_PCT}%", f"{C.SYN_FY17:,}", f"{C.SYN_CTA:,}"])
    table(s, 0.5, 1.95, 7.3, 3.1,
          ["Synergy bucket", "Init.", "Target", "Secured", "%", "FY17 in-year", "Cost to achieve"],
          rows, [2.2, 0.6, 1.0, 1.0, 0.6, 1.1, 1.2])
    tb(s, 8.1, 1.95, 4.73, 0.3, "Bridge to the deal model, USD m run-rate",
       size=11, bold=True, color=C.D_DARK)
    bridge = [["Deal model target", f"{C.DEAL_MODEL_TARGET:,}"],
              ["In register", f"{C.SYN_TARGET:,}"],
              ["Gap to deal model", f"{C.SYN_GAP:,}"],
              ["Secured, all initiatives", f"{C.SYN_SECURED:,}"],
              ["Secured, Finance validated", f"{C.SYN_VALIDATED:,}"],
              ["Realised to date", f"{C.SYN_REALISED:,}"],
              ["Cost synergy target", f"{C.COST_TARGET:,}"],
              ["Revenue synergy target", f"{C.REV_TARGET:,}"]]
    table(s, 8.1, 2.3, 4.73, 2.75, ["Line item", "USD m"], bridge, [3.0, 1.4])
    tb(s, 0.5, 5.25, 12.33, 1.5,
       f"Revenue synergies are designed at roughly {C.REV_COST_RATIO} times cost synergies, "
       f"consistent with the ratio management stated publicly at announcement. Only "
       f"USD {C.SYN_VALIDATED} m has passed Finance validation, all of it in cost buckets; no "
       f"revenue initiative has yet been validated, which is why the secured figure and the "
       f"reportable figure differ by USD {C.SYN_SECURED - C.SYN_VALIDATED} m. Issue I-03 records "
       f"that the session 01 pack overstated this, and decision B-06 fixed the rule.",
       size=10, fill=C.D_PALE)
    footer(s)

    # --- decisions
    s = add(p, L_SUB)
    title(s, "Decisions Requested and Actions Carried Forward",
          "Key message: one decision today, and the two actions still open from session 01 "
          "both sit on the critical path.")
    dec_rows = [[d[0], d[1][:86], d[3][:74], C.nm(d[7]), d[8]]
                for d in C.DECISIONS if d[5] == C.STEERCO_02 or d[8] == "Open"]
    table(s, 0.5, 1.95, 12.33, 1.9,
          ["ID", "Decision requested", "Rationale", "Proposed by", "Status"],
          dec_rows, [0.6, 5.0, 4.6, 1.3, 0.9])
    tb(s, 0.5, 4.05, 12.33, 0.3, "Actions carried forward from session 01",
       size=11, bold=True, color=C.D_DARK)
    act_rows = []
    for a in C.ACTIONS:
        if "session 01" in a[3] and a[9] != "Done":
            act_rows.append([a[0], a[1][:80], C.nm(a[4]), a[5],
                             C.de(a[6]).replace(".2016", ""),
                             C.de(a[7]).replace(".2016", ""), a[8], a[9]])
    table(s, 0.5, 4.4, 12.33, 1.5,
          ["ID", "Action", "Owner", "WS", "Original due", "Revised due", "Shifts", "Status"],
          act_rows, [0.7, 5.4, 1.4, 0.6, 1.2, 1.2, 0.7, 1.1])
    tb(s, 0.5, 6.05, 12.33, 0.8,
       f"{C.ACT_DONE} of {C.ACT_TOTAL} tracked actions are closed and {C.ACT_SHIFTED} have been "
       f"rescheduled at least once. Action OP-01 is recorded in the signed session 01 minutes "
       f"against {C.nm(C.OP01_OWNER_MINUTES)}; it was reassigned by e-mail after the meeting and "
       f"the minutes were never amended.", size=10, fill=C.D_PALE)
    footer(s)

    # --- appendix
    s = add(p, L_DIV)
    for ph in s.placeholders:
        r = ph.text_frame.paragraphs[0].add_run()
        r.text = "Appendix"
        r.font.bold = True
        break

    s = add(p, L_SUB)
    title(s, "Appendix: Governance and Deal Reference",
          "Source: company filings and the published advisory case study.")
    gov = [[n, r, m] for n, r, m in C.PUBLIC_GOVERNANCE]
    table(s, 0.5, 1.95, 6.0, 1.6, ["Name", "Role", "Integration mandate"], gov, [1.6, 2.4, 2.0])
    facts = [
        ["Announcement", C.en(C.ANNOUNCE)],
        ["Close, Day 1", C.en(C.DAY1)],
        ["Transaction value", f"{C.DEAL_VALUE} (headline {C.DEAL_VALUE_HEADLINE})"],
        ["Consideration", f"{C.CASH_PER_SHARE} cash plus {C.TRACKING_RATIO} tracking shares "
                          f"per target share"],
        ["Committed debt financing", C.DEBT_COMMITMENT],
        ["Combined revenue", C.COMBINED_REVENUE],
        ["Employees", f"{C.EMPLOYEES:,}"],
        ["Countries served", str(C.COUNTRIES)],
        ["Workstreams in the programme", C.WORKSTREAM_COUNT_TOTAL],
        ["Announcement to close", f"{C.JOURNEY_MONTHS} months"],
    ]
    table(s, 6.84, 1.95, 6.0, 2.8, ["Reference point", "Value"], facts, [2.4, 3.6])
    tb(s, 0.5, 5.0, 12.33, 1.8,
       "Sources\n" + "\n".join(f"  {i+1}. {s_}" for i, s_ in enumerate(C.SOURCES)) +
       "\n\nAll operational names, figures, risks, actions and synergy values in this pack are "
       "synthetic and were constructed for testing purposes. They are internally consistent but "
       "are not a representation of what any party actually did.",
       size=9, color=C.D_GREY)
    footer(s)

    p.save(C.OUT / "DellEMC_VCIO_SteerCo_Update_Session02_2016-09-29.pptx")


# =====================================================================
# 2 + 3. Workstream one-pagers
# =====================================================================
def one_pager(code, fname, progress_override=None):
    p = deck()
    s = add(p, L_ONLY)
    prog = progress_override if progress_override is not None else C.ws_progress(code)
    rag = C.ws_rag(code)
    title(s, f"{code} {C.WS_NAME[code]}: Weekly Status")
    tb(s, 0.5, 0.78, 12.33, 0.3,
       f"{C.OFFICE} ({C.OFFICE_ABBR})  |  Reporting week {C.WEEK_LABEL}  |  "
       f"Workstream lead {C.nm(code)}  |  Prepared {C.en(C.D_WS_ONEPAGER)}",
       size=10, color=C.D_GREY)

    ms = C.ws_milestones(code)
    rk = C.ws_risks(code)
    ts = C.ws_tasks(code)
    tiles = [
        ("Overall RAG", rag, C.RAG_COLOR[rag]),
        ("Progress vs. plan", f"{prog}%", C.D_BLACK),
        ("Milestones open", str(sum(1 for m in ms if m[8] != "done")), C.D_BLACK),
        ("Overdue or blocked", str(C.ws_overdue(code)), C.D_BLACK),
        ("Open risks", str(len(rk)), C.D_BLACK),
        ("Top risk severity", str(max((C.sev(r) for r in rk), default=0)), C.D_BLACK),
    ]
    for i, (lab, val, col) in enumerate(tiles):
        x = 0.5 + i * 2.06
        tb(s, x, 1.2, 1.95, 0.28, lab, size=8.5, bold=True, color="FFFFFF", fill=C.D_DARK)
        tb(s, x, 1.48, 1.95, 0.42, val, size=15, bold=True, color=col, fill=C.D_PALE)

    tb(s, 0.5, 2.1, 6.0, 0.28, "Progress this week", size=11, bold=True, color=C.D_DARK)
    rows = [[t[0], t[1][:46], C.nm(t[4]).split()[-1], C.de(t[5]).replace(".2016", ""),
             t[6], f"{t[7]}%"] for t in ts[:5]]
    table(s, 0.5, 2.45, 6.0, 1.85,
          ["ID", "Task", "Owner", "Due", "Status", "%"], rows, [0.7, 2.5, 0.9, 0.7, 0.8, 0.4],
          hdr=8, body=7.5)

    tb(s, 6.84, 2.1, 6.0, 0.28, "Milestones and dependencies", size=11, bold=True, color=C.D_DARK)
    mrows = [[m[0], m[1][:40], C.de(m[6]).replace(".2016", ""),
              C.de(m[7]).replace(".2016", ""),
              {"done": "Done", "on_track": "On track", "at_risk": "At risk",
               "delayed": "Delayed"}[m[8]]] for m in ms]
    deps = [d for d in C.DEPENDENCIES if d[2] == code or d[3] == code][:2]
    for d in deps:
        mrows.append([d[0], d[1][:40], "-", C.de(d[4]).replace(".2016", ""), d[7]])
    table(s, 6.84, 2.45, 6.0, 1.85,
          ["ID", "Milestone or dependency", "Baseline", "Forecast", "Status"],
          mrows, [0.7, 2.6, 0.9, 0.9, 0.9], hdr=8, body=7.5)

    tb(s, 0.5, 4.5, 6.0, 0.28, "Risks and blockers", size=11, bold=True, color=C.D_DARK)
    rrows = [[r[0], r[1][:52], f"{C.sev(r)}", r[7][:44], C.nm(r[4]).split()[-1]]
             for r in sorted(rk, key=C.sev, reverse=True)]
    table(s, 0.5, 4.85, 6.0, 1.5,
          ["ID", "Risk", "Sev", "Mitigation", "Owner"], rrows, [0.6, 2.4, 0.5, 1.8, 0.7],
          hdr=8, body=7.5)

    tb(s, 6.84, 4.5, 6.0, 0.28, "Asks and escalations", size=11, bold=True, color=C.D_DARK)
    asks = [a for a in C.ACTIONS if a[5] == code][:3]
    arows = [[a[1][:56], C.nm(a[4]).split()[-1], C.de(a[7]).replace(".2016", ""), a[9]]
             for a in asks]
    if not arows:
        arows = [["No open escalation from this workstream this week", "-", "-", "-"]]
    table(s, 6.84, 4.85, 6.0, 1.5,
          ["Ask", "Owner", "Needed by", "Status"], arows, [3.4, 0.9, 0.9, 0.8],
          hdr=8, body=7.5)

    note = ""
    if progress_override is not None:
        note = (f"  |  Progress stated here is the workstream's own roll-up and is "
                f"{progress_override - C.ws_progress(code):+d} points against the integration "
                f"tracker")
    tb(s, 0.5, 6.5, 12.33, 0.45,
       f"Source: integration tracker as at {C.en(C.D_TRACKER)}, decision B-02 makes it the "
       f"single source of truth{note}.", size=8.5, color=C.D_GREY)
    footer(s)
    p.save(C.OUT / fname)


# =====================================================================
# 4. Integration roadmap
# =====================================================================
def roadmap():
    p = deck()
    s = add(p, L_TITLE)
    ph_text(s, 0, f"{C.OFFICE}\nIntegration Roadmap, Day 1 to Day 100", size=24, bold=True)
    ph_text(s, 10, f"{C.ACQUIRER} and {C.TARGET}   |   Version 2.0   |   "
                   f"Last refreshed {C.en(C.D_ROADMAP)}", size=10)

    s = add(p, L_SUB)
    title(s, "Integration Phases and Gates",
          f"Key message: the programme has {C.DAYS_TO_DAY100} days from today to Day 100 and two "
          f"of the four remaining gates depend on decisions that are not yet taken.")
    phases = [
        ("Pre-Close", f"{C.en(C.ANNOUNCE)} to Day 1",
         "Regulatory clearance, clean team protocol, Day 1 readiness",
         "Value prioritisation framework agreed, e-guide and runbook built"),
        ("Day 1", C.en(C.DAY1),
         "Legal close, brand launch, employee and sales enablement live",
         f"{C.MS_DONE} of {C.MS_DONE} Day 1 milestones achieved on the day"),
        ("Day 1 to Day 30", f"to {C.en(C.DAY30)}",
         "Opening balance sheet, entity plan, ERP blueprint, organisation design",
         "Two gate milestones delayed, one at risk"),
        ("Day 30 to Day 100", f"to {C.en(C.DAY100)}",
         "Coverage model live, consultation concluded, synergy wave 1 contracted",
         f"Flagship event {C.en(C.DELL_EMC_WORLD)} is a fixed external date"),
        ("Post Day 100", "Handover to run",
         "Integration organisation dissolved into functional ownership",
         "Open items transferred with named owners"),
    ]
    for i, (nme, sub, act, stat) in enumerate(phases):
        x = 0.5 + i * 2.47
        tb(s, x, 1.9, 2.35, 0.3, nme, size=11, bold=True, color="FFFFFF", fill=C.D_DARK)
        tb(s, x, 2.22, 2.35, 0.3, sub, size=8, color=C.D_GREY, fill=C.D_LIGHT)
        tb(s, x, 2.58, 2.35, 1.5, "Key activities\n" + act, size=8, fill=C.D_PALE)
        tb(s, x, 4.15, 2.35, 1.3, "Status\n" + stat, size=8)
    footer(s)

    s = add(p, L_SUB)
    title(s, "Master Milestone List",
          f"Key message: {C.MS_GATE} of {C.MS_TOTAL} milestones are gate relevant, and the "
          f"critical path now runs through Tax and IT.")
    rows, colors = [], {}
    for i, m in enumerate(C.MILESTONES):
        d = (m[7] - m[6]).days
        rows.append([m[0], m[1][:66], m[4] if m[4] in C.WS_CODES else m[3], m[3],
                     C.nm(m[4]) if m[4] in C.PEOPLE else "",
                     C.de(m[6]).replace(".2016", ""), f"+{d}" if d else "0",
                     "Yes" if m[9] else "No",
                     {"done": "Done", "on_track": "On track", "at_risk": "At risk",
                      "delayed": "Delayed"}[m[8]]])
        if m[8] == "delayed":
            colors[(i, 8)] = C.RAG_COLOR["Red"]
        elif m[8] == "at_risk":
            colors[(i, 8)] = C.RAG_COLOR["Amber"]
    # roadmap still carries the pre-slip forecast for M-07: planted conflict C2
    for i, m in enumerate(C.MILESTONES):
        if m[0] == "M-07":
            rows[i][5] = C.de(C.M07_ROADMAP).replace(".2016", "")
            rows[i][6] = "0"
    table(s, 0.5, 1.9, 12.33, 4.6,
          ["ID", "Milestone", "WS", "Phase WS", "Owner", "Baseline", "Delta", "Gate", "Status"],
          [[r[0], r[1], r[3], m[5], r[4], r[5], r[6], r[7], r[8]]
           for r, m in zip(rows, C.MILESTONES)],
          [0.6, 5.0, 0.6, 1.5, 1.3, 1.0, 0.7, 0.6, 1.0], hdr=8, body=7.5, colcolor=colors)
    tb(s, 0.5, 6.6, 12.33, 0.35,
       f"This roadmap was last refreshed on {C.en(C.D_ROADMAP)}. Milestone forecasts move faster "
       f"than this document; the integration tracker carries the current forecast.",
       size=8, color=C.D_GREY)
    footer(s)

    s = add(p, L_SUB)
    title(s, "Cross-Workstream Dependencies",
          f"Key message: {C.DEP_UNCONFIRMED_CRITICAL} critical or high dependencies are still "
          f"unconfirmed by both leads, which decision B-03 was meant to eliminate.")
    rows, colors = [], {}
    for i, d in enumerate(C.DEPENDENCIES):
        rows.append([d[0], d[1][:72], d[2], d[3], C.de(d[4]).replace(".2016", ""),
                     d[5], d[6], d[7]])
        if d[7] in ("At risk", "Delayed"):
            colors[(i, 7)] = C.RAG_COLOR["Red" if d[7] == "Delayed" else "Amber"]
        if d[6] == "No":
            colors[(i, 6)] = C.RAG_COLOR["Red"]
    table(s, 0.5, 1.9, 12.33, 3.2,
          ["ID", "Deliverable handed over", "From", "To", "Needed by", "Criticality",
           "Both leads confirmed", "Status"],
          rows, [0.7, 5.6, 0.7, 0.7, 1.1, 1.1, 1.5, 1.0], colcolor=colors)
    tb(s, 0.5, 5.3, 12.33, 1.5,
       "The advisory team identified interdependencies between workstreams as the defining "
       "challenge of a programme of this size: individuals in each workstream understand their "
       "own responsibilities but may not see how their activities land on others. This register "
       "exists to make that visible, and decision B-03 requires both sides of every handover to "
       "confirm the date. Two entries still show one-sided confirmation.",
       size=10, fill=C.D_PALE)
    footer(s)

    p.save(C.OUT / "DellEMC_VCIO_Integration_Roadmap_Day1_to_Day100_2016-09-16.pptx")


if __name__ == "__main__":
    steerco()
    one_pager("WS3", "DellEMC_VCIO_Workstream_Status_IT_W3_2016-09-28.pptx",
              progress_override=C.IT_PROGRESS_ONEPAGER)
    one_pager("WS4", "DellEMC_VCIO_Workstream_Status_HumanCapital_W3_2016-09-28.pptx")
    roadmap()
    print("pptx done")
