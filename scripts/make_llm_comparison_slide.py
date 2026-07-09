"""Build the one-slide LLM provider scorecard: Claude Sonnet 5 vs GPT-5.5.

Content source: docs/LLM_Provider_Comparison.md §6.
Output: ../../LLM_Provider_Comparison_Slide.pptx (shared PMI folder).
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

GREEN = RGBColor(0x2E, 0x7D, 0x32)
GREEN_DARK = RGBColor(0x1B, 0x5E, 0x20)
TINT = RGBColor(0xE8, 0xF5, 0xE9)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x69, 0x69, 0x69)
LIGHT = RGBColor(0xF4, 0xF4, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xB5, 0x6A, 0x00)
FONT = "Verdana"

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
slide = prs.slides.add_slide(prs.slide_layouts[6])
shapes = slide.shapes


def textbox(x, y, w, h, lines, margin=0.03):
    box = shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    for i, (text, size, bold, color, space) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = text
        r.font.name, r.font.size, r.font.bold = FONT, Pt(size), bold
        r.font.color.rgb = color
        p.space_after = Pt(space)
    return box


def card(x, y, w, h, fill=WHITE, line=RGBColor(0xD6, 0xD6, 0xD6)):
    from pptx.enum.shapes import MSO_SHAPE
    shp = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


# ---------------------------------------------------------------- header + action title
textbox(0.35, 0.12, 9.0, 0.25,
        [("H2 Automated Reporting · LLM Provider Decision", 9, False, GREY, 0)])
textbox(10.6, 0.10, 2.4, 0.3, [("Deloitte. | TUM", 11, True, DARK, 0)])
textbox(0.35, 0.38, 12.63, 0.80,
        [("Evidence-backed parameters lean Claude Sonnet 5 (cost, calibration, governance) — "
          "the four open cells are exactly what the evaluation harness measures, so run the A/B before committing",
          17, True, DARK, 0)])

# ---------------------------------------------------------------- scope note
card(0.35, 1.24, 12.63, 0.38, fill=LIGHT, line=LIGHT)
textbox(0.45, 1.30, 12.43, 0.28, [
    ("Workload basis: 2 LLM calls per report — request classification (~0.3K tokens) + summary bullets from the "
     "extracted data model (~10–20K in / 1–2K out). All numbers are computed deterministically upstream — the LLM never invents figures.",
     8.5, False, GREY, 0),
])

# ---------------------------------------------------------------- table
TX, TY = 0.35, 1.76
COL = [2.05, 4.00, 3.65, 2.93]  # parameter | sonnet | gpt | edge
XS = [TX]
for c in COL[:-1]:
    XS.append(XS[-1] + c)
HDR_H = 0.30

hdrs = ["Parameter", "Claude Sonnet 5", "OpenAI GPT-5.5 (standard)", "Edge"]
for j, (hx, hw, ht) in enumerate(zip(XS, COL, hdrs)):
    card(hx, TY, hw, HDR_H, fill=GREEN, line=GREEN)
    textbox(hx + 0.04, TY + 0.035, hw - 0.08, 0.24, [(ht, 10, True, WHITE, 0)])

rows = [
    ("Cost per report",
     "$0.03–0.06 (intro $2/$10 per 1M; $3/$15 from Sep) — 20 EUR ≈ 350–1,000 reports",
     "$0.09–0.16 ($5/$30 per 1M); caching −90% on repeated prompt — 20 EUR ≈ 130–250 reports",
     "Sonnet 5 — ~3× cheaper [volume ≤200 reports/mo assumed]", "sonnet", 0.52),
    ("Structured JSON reliability",
     "Structured outputs / tool-use JSON schema supported",
     "Native json_object mode — already implemented in app/agent/llm.py",
     "OPEN — % valid parses over 100 harness runs [both ≥99% assumed: two small schemas]", "open", 0.50),
    ("Numerical faithfulness",
     "Better-calibrated model family in 2026 hallucination studies",
     "Higher fabrication tendency in the same studies",
     "OPEN — Stage-3 bullet-number precision test [both ≥95% assumed: figures handed in-prompt]", "open", 0.50),
    ("Audience fit (SteerCo register)",
     "Opus 4.8 leads OfficeQA Pro 66.2 vs 54.1 [Sonnet 5 inheritance unverified]",
     "54.1 on OfficeQA Pro (flagship measurement)",
     "OPEN — blinded SM rating on the 15–30 held-out PMI examples is the only valid test", "open", 0.50),
    ("Calibration on missing data",
     "Abstains on ~64% of unknown-answer questions [family behavior assumed for Sonnet 5]",
     "Attempts ~86% of unknowns — mostly confidently wrong",
     "Claude — missing owners/dates must surface as \"unknown\" behind the SM gate", "sonnet", 0.50),
    ("Latency (target <15 s)",
     "Anthropic fast tier; short calls",
     "~2× faster on agentic evals; batch/flex modes async (not UI-suitable)",
     "OPEN [both <10 s assumed for ≤20K-token calls — not a differentiator]", "open", 0.44),
    ("Governance / hosting",
     "On project's governance-track shortlist; AWS Bedrock / GCP Vertex [Sonnet 5 EU regions unverified]",
     "Prototype-only clearance per prior research; no Deloitte alliance; Azure EU hosting exists",
     "Sonnet 5 for anything beyond prototype; tie for current MVP", "sonnet", 0.56),
    ("Engineering fit",
     "langchain-anthropic first-class in LangGraph; swap = one client class [~1 h incl. tests]",
     "Already implemented — the PDF spec's named choice; zero switching cost",
     "GPT-5.5 today — delta trivial by design (env-switched provider layer)", "gpt", 0.50),
]

y = TY + HDR_H
for i, (param, sonnet, gpt, edge, tag, rh) in enumerate(rows):
    fill = WHITE if i % 2 == 0 else LIGHT
    for hx, hw in zip(XS, COL):
        card(hx, y, hw, rh, fill=fill, line=RGBColor(0xE3, 0xE3, 0xE3))
    edge_color = {"sonnet": GREEN_DARK, "gpt": DARK, "open": AMBER}[tag]
    textbox(XS[0] + 0.04, y + 0.04, COL[0] - 0.08, rh - 0.08, [(param, 8.5, True, DARK, 0)])
    textbox(XS[1] + 0.04, y + 0.04, COL[1] - 0.08, rh - 0.08, [(sonnet, 8, False, DARK, 0)])
    textbox(XS[2] + 0.04, y + 0.04, COL[2] - 0.08, rh - 0.08, [(gpt, 8, False, DARK, 0)])
    textbox(XS[3] + 0.04, y + 0.04, COL[3] - 0.08, rh - 0.08, [(edge, 8, True, edge_color, 0)])
    y += rh

# ---------------------------------------------------------------- bottom line
card(0.35, y + 0.10, 12.63, 0.52, fill=TINT, line=GREEN)
textbox(0.45, y + 0.16, 12.43, 0.42, [
    ("Bottom line:  Sonnet 5 is the default if cost matters or a production path opens; GPT-5.5 stands for zero-touch "
     "continuity with the spec. Both fit the 20 EUR/month budget. Decide the four OPEN cells via the Stage 1–4 A/B on our "
     "own harness — the provider switch is a one-line environment change.", 9.5, True, GREEN_DARK, 0),
])

# ---------------------------------------------------------------- footer
textbox(0.35, 7.12, 12.63, 0.30, [
    ("Sources: OpenAI & Anthropic pricing pages (checked 09.07.2026); OfficeQA Pro & agentic evals via CodingFleet, "
     "DataCamp, BenchLM; calibration studies via DigitalApplied, CometAPI, MindStudio — third-party, not peer-reviewed; "
     "governance per PMI_Deep_Research_Open_Decisions_v2 / UC2_V2 §3.   |   TUM Project Study x Deloitte 2026 · 09.07.2026",
     7.5, False, GREY, 0),
])

out = Path(__file__).resolve().parents[2] / "LLM_Provider_Comparison_Slide.pptx"
prs.save(str(out))
print(f"Saved {out}")
