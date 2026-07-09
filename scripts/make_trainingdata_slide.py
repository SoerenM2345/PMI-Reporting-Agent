"""Build the one-slide training-data decision summary (Deloitte style), v2.

Layout per team feedback 2026-07-09:
- Key decision (synthetic vs proxy) alone in the decision box
- "Further decision questions" section below, each with a context column
- Compact, project-tailored evaluation plan

Output: ../../TrainingData_Decision_Slide.pptx (shared PMI folder).
Content source: docs/TrainingData_Decision.md
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

GREEN = RGBColor(0x2E, 0x7D, 0x32)
GREEN_DARK = RGBColor(0x1B, 0x5E, 0x20)
TINT = RGBColor(0xE8, 0xF5, 0xE9)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x69, 0x69, 0x69)
LIGHT = RGBColor(0xF4, 0xF4, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xB7, 0x1C, 0x1C)
FONT = "Verdana"

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
slide = prs.slides.add_slide(prs.slide_layouts[6])
shapes = slide.shapes


def textbox(x, y, w, h, lines, default_size=9, margin=0.03):
    box = shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, spec in enumerate(lines):
        if isinstance(spec, str):
            spec = (spec, default_size, False, DARK, 2)
        text, size, bold, color, space = spec
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = text
        r.font.name, r.font.size, r.font.bold = FONT, Pt(size), bold
        r.font.color.rgb = color
        p.space_after = Pt(space)
    return box


def card(x, y, w, h, fill=WHITE, line=RGBColor(0xD6, 0xD6, 0xD6)):
    from pptx.enum.shapes import MSO_SHAPE
    shp = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def header_bar(x, y, w, title, h=0.28):
    bar = card(x, y, w, h, fill=GREEN, line=GREEN)
    tf = bar.text_frame
    tf.margin_left, tf.margin_top, tf.margin_bottom = Inches(0.06), Inches(0.01), Inches(0.01)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name, r.font.size, r.font.bold, r.font.color.rgb = FONT, Pt(11), True, WHITE


# ------------------------------------------------------------- header + action title
textbox(0.35, 0.12, 9.0, 0.25,
        [("H2 Automated Reporting · Training Data Concept & Decision", 9, False, GREY, 0)])
textbox(10.6, 0.10, 2.4, 0.3, [("Deloitte. | TUM", 11, True, DARK, 0)])
textbox(0.35, 0.38, 12.63, 0.80,
        [("Proxy corpora cover report generation only — closing the two remaining data "
          "gaps requires a synthetic vs. SEC EDGAR decision, validated on a held-out PMI test set",
          17, True, DARK, 0)])

# ------------------------------------------------------------- row 1: bottleneck strip
strip = [
    ("BOTTLENECK — UNCHANGED", "0 PMI-native reporting datasets exist (confirmed by search); every candidate is a proxy domain", True),
    ("VERIFIED PROXIES", "ECTSum 2,425 pairs in repo (1681/249/495); QMSum + AMI gated by open transcript-scope decision", False),
    ("GAP 1 — NOT COVERED", "Multi-format extraction (xlsx/pptx/docx/pdf/html → one PMI schema): no corpus teaches it", True),
    ("GAP 2 — NOT COVERED", "Cross-source conflicts (Excel 82% vs. PPT 75%): no corpus contains them by design", True),
]
sw = 3.10
for i, (label, body, is_red) in enumerate(strip):
    x = 0.35 + i * (sw + 0.075)
    card(x, 1.26, sw, 0.80, fill=LIGHT, line=LIGHT)
    textbox(x + 0.04, 1.30, sw - 0.08, 0.72, [
        (label, 8.5, True, RED if is_red else GREEN_DARK, 2),
        (body, 8.5, False, DARK, 0),
    ])

# ------------------------------------------------------------- row 2: THE key decision
card(0.35, 2.20, 12.63, 0.42, fill=GREEN_DARK, line=GREEN_DARK)
textbox(0.45, 2.27, 12.43, 0.30, [
    ("KEY DECISION:  Close data gaps 1–2 with synthetic data, a SEC EDGAR proxy corpus — or a hybrid of both?",
     12, True, WHITE, 0),
])

# ------------------------------------------------------------- row 3: option boxes
row_y, row_h = 2.74, 2.02
cx, cw = 0.35, 6.25
header_bar(cx, row_y, cw, "OPTION A — SYNTHETIC DATA")
card(cx, row_y + 0.28, cw, row_h - 0.28)
textbox(cx + 0.07, row_y + 0.35, cw - 0.14, row_h - 0.42, [
    ("Context: gap 2 is narrow and mechanical, on a schema we own; seeded-conflict generator + tests already exist in the repo", 8, False, GREY, 4),
    ("+  Exact fit to PMI schema and priority rule; gold labels by construction", 8.5, False, GREEN_DARK, 2),
    ("+  Only practical source of guaranteed conflicts; unlimited volume, no license risk; near-zero cost", 8.5, False, GREEN_DARK, 4),
    ("–  Circular: tests what we generated; weak external validity as thesis evidence", 8.5, False, RED, 2),
    ("–  LLM-generated data risks bias and model collapse if it replaces real data (Long et al. 2024, ACL Findings)", 8.5, False, RED, 2),
    ("–  Lacks real-world formatting noise", 8.5, False, RED, 0),
])

px, pw = 6.73, 6.25
header_bar(px, row_y, pw, "OPTION B — PROXY: SEC EDGAR")
card(px, row_y + 0.28, pw, row_h - 0.28)
textbox(px + 0.07, row_y + 0.35, pw - 0.14, row_h - 0.42, [
    ("Context: free SEC filings APIs (no key, ~10 req/s); same fact appears per company-quarter in XBRL tags, 10-K/10-Q HTML, 8-K Ex-99.1 press releases, PDFs", 8, False, GREY, 4),
    ("+  XBRL = free ground-truth labels for gap-1 extraction from real documents", 8.5, False, GREEN_DARK, 2),
    ("+  Real formatting noise; auditable and citable; M&A-adjacent forms (8-K, S-4, DEFM14A)", 8.5, False, GREEN_DARK, 4),
    ("–  Register mismatch: regulated filings ≠ internal PMI status reports", 8.5, False, RED, 2),
    ("–  True conflicts rare — gap 2 still needs synthetic injection (exception: non-GAAP vs. GAAP deltas)", 8.5, False, RED, 2),
    ("–  Prep effort (alignment, parsing); US-listed firms only", 8.5, False, RED, 0),
])

# ------------------------------------------------------------- row 4: recommendation
card(0.35, 4.88, 12.63, 0.44, fill=TINT, line=GREEN)
textbox(0.45, 4.94, 12.45, 0.34, [
    ("Recommendation — hybrid per sub-skill:  ECTSum (+QMSum if transcripts stay in scope) for generation  ·  "
     "EDGAR XBRL↔HTML pairs for extraction (gap 1)  ·  synthetic conflict injection (gap 2)  ·  "
     "15–30 real PMI examples always held out", 9.5, True, GREEN_DARK, 0),
])

# ------------------------------------------------------------- row 5: further questions | eval plan
fy, fh = 5.46, 1.52
# left: further decision questions with context column
fqx, fqw = 0.35, 6.95
header_bar(fqx, fy, fqw, "FURTHER DECISION QUESTIONS", h=0.26)
card(fqx, fy + 0.26, fqw, fh - 0.26)
questions = [
    ("Transcripts in V2 scope?",
     "Slide 5 lists documents only; Interview 7 demands transcripts — gates QMSum + AMI entirely"),
    ("Full 8-step review for this choice?",
     "Xiao & Watson protocol caught real errors before; costs ~1 research day"),
    ("ECTSum re-split by ticker?",
     "310 tickers overlap train/test; overlap inflates our generation scores"),
    ("License gate (GPL-3.0)?",
     "Internal fine-tuning/eval OK; re-check before sharing weights or data outside the project"),
]
qy = fy + 0.33
for q, ctx in questions:
    textbox(fqx + 0.06, qy, 2.35, 0.30, [(q, 8, True, DARK, 0)])
    textbox(fqx + 2.48, qy, fqw - 2.56, 0.30, [(ctx, 8, False, GREY, 0)])
    qy += 0.295

# right: compact evaluation plan
epx, epw = 7.43, 5.55
header_bar(epx, fy, epw, "EVALUATION PLAN — 5 STEPS, RUN ON OUR REPO", h=0.26)
card(epx, fy + 0.26, epw, fh - 0.26)
steps = [
    ("0 Freeze", "15–30 real PMI examples + ticker-disjoint ECTSum, never in training"),
    ("1 Extract", "F1 per PMI-schema field; numbers ≥ 95% correct"),
    ("2 Conflict", "Recall on seeded 82/75-type conflicts; resolution vs. SM gold"),
    ("3 Generate", "ROUGE/BERTScore + bullet-number precision; blinded SM rating"),
    ("4 Gate", "PMI score within 10–15% of proxy score — else revisit key decision"),
]
sy = fy + 0.32
for label, body in steps:
    textbox(epx + 0.06, sy, 0.85, 0.24, [(label, 8, True, GREEN_DARK, 0)])
    textbox(epx + 0.95, sy, epw - 1.05, 0.24, [(body, 8, False, DARK, 0)])
    sy += 0.235

# ------------------------------------------------------------- footer
textbox(0.35, 7.10, 12.63, 0.32, [
    ("Sources: Mukherjee et al. 2022 (EMNLP); Zhong et al. 2021 (NAACL); Carletta et al. 2005 (MLMI); "
     "Long et al. 2024 (ACL Findings); SEC.gov EDGAR developer resources; H2 Deep Dive doc (2026-07-08); "
     "UC2_V2_SingleAgent_Definition §5.   |   TUM Project Study x Deloitte 2026 · 09.07.2026", 7.5, False, GREY, 0),
])

out = Path(__file__).resolve().parents[2] / "TrainingData_Decision_Slide.pptx"
prs.save(str(out))
print(f"Saved {out}")
