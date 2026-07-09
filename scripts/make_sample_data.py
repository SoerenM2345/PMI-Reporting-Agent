"""Generate realistic sample PMI input files into data/samples/.

Intentionally includes a cross-source conflict: the Excel tracker reports
overall progress = 82%, the PowerPoint update says 75% — exactly the
consistency-check example from the spec (slide 5, step 5).

Run:  python scripts/make_sample_data.py
"""
from __future__ import annotations

from pathlib import Path

SAMPLES = Path(__file__).resolve().parents[1] / "data" / "samples"
SAMPLES.mkdir(parents=True, exist_ok=True)


def make_excel():
    import xlsxwriter

    wb = xlsxwriter.Workbook(str(SAMPLES / "integration_tracker.xlsx"))
    hdr = wb.add_format({"bold": True, "bg_color": "#2E7D32", "font_color": "white"})

    ws = wb.add_worksheet("Workplan")
    ws.write_row(0, 0, ["Task", "Owner", "Status", "Progress %", "Due Date",
                        "Workstream", "Priority"], hdr)
    rows = [
        ["Day-1 readiness checklist", "Anna Schmidt", "Done", 100, "2026-06-15", "PMO", "High"],
        ["Migrate ERP chart of accounts", "Jonas Weber", "In Progress", 60, "2026-07-31", "Finance", "High"],
        ["Harmonize payroll systems", "Lisa Chen", "In Progress", 45, "2026-08-15", "HR", "Medium"],
        ["Consolidate CRM instances", "Marco Rossi", "Blocked", 30, "2026-07-20", "IT", "High"],
        ["Define target operating model", "Anna Schmidt", "In Progress", 70, "2026-07-25", "PMO", "High"],
        ["Customer communication plan", "Sofia Ivanova", "Not Started", 0, "2026-08-01", "Sales", "Medium"],
        ["Rebrand office signage", "Marco Rossi", "In Progress", 50, "2026-09-01", "Facilities", "Low"],
    ]
    for i, r in enumerate(rows, 1):
        ws.write_row(i, 0, r)
    ws.write(len(rows) + 2, 0, "Overall Progress: 82%")

    ws2 = wb.add_worksheet("Risk Log")
    ws2.write_row(0, 0, ["Risk", "Severity", "Likelihood", "Owner", "Mitigation", "Status"], hdr)
    risks = [
        ["Key engineer attrition in target company", "High", "Medium", "Lisa Chen",
         "Retention bonus program approved by SteerCo", "Open"],
        ["ERP cutover slips past Q3", "Critical", "Medium", "Jonas Weber",
         "Weekly cutover war-room; fallback to parallel run", "Open"],
        ["Customer churn during rebranding", "Medium", "Low", "Sofia Ivanova",
         "Proactive top-50 account outreach", "Open"],
    ]
    for i, r in enumerate(risks, 1):
        ws2.write_row(i, 0, r)

    ws3 = wb.add_worksheet("Budget")
    ws3.write_row(0, 0, ["Category", "Planned", "Actual", "Workstream"], hdr)
    budget = [
        ["IT integration", 1200000, 950000, "IT"],
        ["Severance & retention", 800000, 610000, "HR"],
        ["Rebranding", 300000, 340000, "Sales"],
        ["External advisors", 500000, 420000, "PMO"],
    ]
    for i, r in enumerate(budget, 1):
        ws3.write_row(i, 0, r)
    wb.close()


def make_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Weekly Integration Update — Workstream Leads"
    box = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
    tf = box.text_frame
    tf.text = "Overall progress: 75% (per workstream self-reporting)"
    for line in [
        "Action: finalize TSA exit plan for IT (Marco Rossi)",
        "Action: prepare SteerCo pre-read for budget reforecast (Jonas Weber)",
        "Action: schedule culture-integration workshops (Lisa Chen)",
        "Milestone Day-100 review moved to 2026-08-20",
    ]:
        p = tf.add_paragraph(); p.text = line; p.font.size = Pt(14)

    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = "Milestones"
    rows, cols = 4, 4
    tbl = s2.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9),
                              Inches(2.4)).table
    for j, h in enumerate(["Milestone", "Due Date", "Status", "Owner"]):
        tbl.cell(0, j).text = h
    data = [
        ["Day-1 legal close", "2026-06-01", "Done", "Anna Schmidt"],
        ["Day-100 review", "2026-08-20", "On Track", "Anna Schmidt"],
        ["ERP go-live", "2026-10-01", "At Risk", "Jonas Weber"],
    ]
    for i, row in enumerate(data, 1):
        for j, v in enumerate(row):
            tbl.cell(i, j).text = v
    prs.save(str(SAMPLES / "weekly_update.pptx"))


def make_docx():
    from docx import Document

    doc = Document()
    doc.add_heading("Integration SteerCo — Meeting Notes, 6 July 2026", level=1)
    doc.add_paragraph("Attendees: A. Schmidt (PMO lead), J. Weber (Finance), "
                      "L. Chen (HR), M. Rossi (IT), S. Ivanova (Sales)")
    doc.add_paragraph("The CRM consolidation is blocked pending the data-privacy "
                      "assessment; legal review expected by 14 July.")
    doc.add_paragraph("Action: deliver data-privacy assessment for CRM merge -> Marco Rossi")
    doc.add_paragraph("Action: update synergy tracking baseline with Q2 actuals -> Jonas Weber")
    doc.add_paragraph("Action: present retention-program uptake figures at next SteerCo -> Lisa Chen")
    doc.add_paragraph("Decision: rebranding budget overrun of 40k EUR approved.")
    doc.add_heading("KPIs discussed", level=2)
    t = doc.add_table(rows=4, cols=4)
    t.style = "Table Grid"
    hdr = ["KPI", "Value", "Target", "Unit"]
    for j, h in enumerate(hdr):
        t.rows[0].cells[j].text = h
    kpis = [
        ["Synergies captured", "12.4", "30", "MEUR"],
        ["Employee attrition (target co.)", "8", "5", "%"],
        ["Systems migrated", "14", "40", "count"],
    ]
    for i, row in enumerate(kpis, 1):
        for j, v in enumerate(row):
            t.rows[i].cells[j].text = v
    doc.save(str(SAMPLES / "steerco_meeting_notes.docx"))


if __name__ == "__main__":
    make_excel()
    make_pptx()
    make_docx()
    print(f"Sample files written to {SAMPLES}")
