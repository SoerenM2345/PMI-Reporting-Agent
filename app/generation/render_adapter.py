"""Adapt GeneratedContent to existing rendering infrastructure.

Bridge between the new LLM-generated content and the rendering pipeline.
"""
from __future__ import annotations

import logging
from pathlib import Path

from app.generation.content_schema import ContentSection, GeneratedContent

log = logging.getLogger("pmi.generation.render_adapter")


def generate_outputs(
    content: GeneratedContent,
    output_formats: list[str],
    output_dir: Path,
) -> dict:
    """Generate output files for approved content.

    Args:
        content: The approved GeneratedContent
        output_formats: List of formats (powerpoint, excel, pdf, word, html, chart)
        output_dir: Directory to write output files

    Returns:
        {format: filepath} mapping
    """
    output_dir.mkdir(parents=True, exist_ok=True)
    outputs = {}

    for fmt in output_formats:
        fmt_lower = fmt.lower()

        if fmt_lower == "powerpoint":
            path = _render_powerpoint(content, output_dir)
        elif fmt_lower == "excel":
            path = _render_excel(content, output_dir)
        elif fmt_lower == "pdf":
            path = _render_pdf(content, output_dir)
        elif fmt_lower == "word":
            path = _render_word(content, output_dir)
        elif fmt_lower == "html":
            path = _render_html(content, output_dir)
        elif fmt_lower in ("chart", "png", "image"):
            path = _render_image(content, output_dir)
        else:
            log.warning("Unknown output format: %s", fmt)
            continue

        if path:
            outputs[fmt_lower] = str(path)

    return outputs


def _render_powerpoint(content: GeneratedContent, output_dir: Path) -> Path:
    """Render to PowerPoint."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = content.title
    subtitle.text = content.subtitle or ""

    # Content slides
    blank_layout = prs.slide_layouts[6]
    for section in content.sections:
        slide = prs.slides.add_slide(blank_layout)
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)

        # Title
        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.75))
        title_frame = title_box.text_frame
        title_frame.text = section.title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True

        # Content
        content_top = top + Inches(1)
        content_height = Inches(6)

        if section.type == "text":
            text_box = slide.shapes.add_textbox(left, content_top, width, content_height)
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.text = str(section.content)

        elif section.type == "bullets":
            text_box = slide.shapes.add_textbox(left, content_top, width, content_height)
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            items = section.content if isinstance(section.content, list) else [section.content]
            for i, item in enumerate(items):
                if i == 0:
                    text_frame.text = str(item)
                else:
                    p = text_frame.add_paragraph()
                    p.text = str(item)
                    p.level = 0

        elif section.type == "table":
            rows = section.content if isinstance(section.content, list) else []
            if rows:
                table_shape = slide.shapes.add_table(
                    len(rows), len(rows[0]) if rows[0] else 1, left, content_top, width
                ).table
                for row_idx, row in enumerate(rows):
                    for col_idx, cell_data in enumerate(row if isinstance(row, list) else [row]):
                        table_shape.cell(row_idx, col_idx).text = str(cell_data)

    output_file = output_dir / f"{_safe_name(content.title)}.pptx"
    prs.save(output_file)
    log.info("rendered powerpoint: %s", output_file)
    return output_file


def _render_excel(content: GeneratedContent, output_dir: Path) -> Path:
    """Render to Excel."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill

    wb = Workbook()
    wb.remove(wb.active)

    # Summary sheet
    ws = wb.create_sheet("Summary", 0)
    ws["A1"] = content.title
    ws["A1"].font = Font(bold=True, size=14)
    if content.subtitle:
        ws["A2"] = content.subtitle

    # Content sheets
    for section in content.sections:
        ws = wb.create_sheet(section.title[:31])  # Excel limit is 31 chars

        # Title
        ws["A1"] = section.title
        ws["A1"].font = Font(bold=True, size=12)

        if section.type == "text":
            ws["A2"] = section.content

        elif section.type == "bullets":
            items = section.content if isinstance(section.content, list) else [section.content]
            for i, item in enumerate(items, start=2):
                ws[f"A{i}"] = str(item)

        elif section.type == "table":
            rows = section.content if isinstance(section.content, list) else []
            for row_idx, row in enumerate(rows, start=2):
                for col_idx, cell_data in enumerate(row if isinstance(row, list) else [row]):
                    ws.cell(row=row_idx, column=col_idx + 1, value=cell_data)

    output_file = output_dir / f"{_safe_name(content.title)}.xlsx"
    wb.save(output_file)
    log.info("rendered excel: %s", output_file)
    return output_file


def _render_pdf(content: GeneratedContent, output_dir: Path) -> Path:
    """Render to PDF."""
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    from reportlab.lib import colors

    output_file = output_dir / f"{_safe_name(content.title)}.pdf"
    doc = SimpleDocTemplate(str(output_file), pagesize=letter)
    styles = getSampleStyleSheet()
    story = []

    # Title
    title_style = ParagraphStyle("CustomTitle", parent=styles["Heading1"], fontSize=24, textColor=colors.HexColor("#1a1a1a"))
    story.append(Paragraph(content.title, title_style))

    if content.subtitle:
        story.append(Paragraph(content.subtitle, styles["Normal"]))

    story.append(Spacer(1, 0.3 * inch))

    # Sections
    for section in content.sections:
        story.append(Paragraph(section.title, styles["Heading2"]))
        story.append(Spacer(1, 0.1 * inch))

        if section.type == "text":
            story.append(Paragraph(str(section.content), styles["Normal"]))

        elif section.type == "bullets":
            items = section.content if isinstance(section.content, list) else [section.content]
            for item in items:
                story.append(Paragraph(f"• {item}", styles["Normal"]))

        elif section.type == "table":
            rows = section.content if isinstance(section.content, list) else []
            if rows:
                table_data = [[str(cell) for cell in row] if isinstance(row, list) else [str(row)] for row in rows]
                tbl = Table(table_data)
                tbl.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#f5f5f5")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.black),
                    ("ALIGN", (0, 0), (-1, -1), "LEFT"),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, 0), 10),
                    ("BOTTOMPADDING", (0, 0), (-1, 0), 12),
                    ("BACKGROUND", (0, 1), (-1, -1), colors.beige),
                    ("GRID", (0, 0), (-1, -1), 1, colors.black),
                ]))
                story.append(tbl)

        story.append(Spacer(1, 0.2 * inch))

    doc.build(story)
    log.info("rendered pdf: %s", output_file)
    return output_file


def _render_word(content: GeneratedContent, output_dir: Path) -> Path:
    """Render to Word."""
    from docx import Document
    from docx.shared import Pt, RGBColor

    doc = Document()
    doc.add_heading(content.title, level=0)

    if content.subtitle:
        doc.add_paragraph(content.subtitle, style="Subtitle")

    for section in content.sections:
        doc.add_heading(section.title, level=1)

        if section.type == "text":
            doc.add_paragraph(str(section.content))

        elif section.type == "bullets":
            items = section.content if isinstance(section.content, list) else [section.content]
            for item in items:
                doc.add_paragraph(str(item), style="List Bullet")

        elif section.type == "table":
            rows = section.content if isinstance(section.content, list) else []
            if rows:
                table = doc.add_table(rows=len(rows), cols=len(rows[0]) if rows[0] else 1)
                for row_idx, row in enumerate(rows):
                    for col_idx, cell_data in enumerate(row if isinstance(row, list) else [row]):
                        table.rows[row_idx].cells[col_idx].text = str(cell_data)

    output_file = output_dir / f"{_safe_name(content.title)}.docx"
    doc.save(output_file)
    log.info("rendered word: %s", output_file)
    return output_file


def _render_html(content: GeneratedContent, output_dir: Path) -> Path:
    """Render to HTML."""
    from app.generation.preview import content_to_html

    html = content_to_html(content)
    output_file = output_dir / f"{_safe_name(content.title)}.html"
    output_file.write_text(html)
    log.info("rendered html: %s", output_file)
    return output_file


def _render_image(content: GeneratedContent, output_dir: Path) -> Path:
    """Render to PNG (placeholder)."""
    log.info("image rendering not yet implemented for generated content")
    # TODO: Implement chart/image rendering
    return None


def _safe_name(text: str) -> str:
    """Make text safe for filename."""
    import re
    return re.sub(r"[^A-Za-z0-9_-]+", "_", text).strip("_") or "document"
