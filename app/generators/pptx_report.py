"""PowerPoint generation (spec §12). python-pptx; the output is an editable .pptx.

Four audience templates, because §12.1-12.4 are genuinely different documents:

* **Steering Committee** (§12.1) — decision-oriented. Is it on track, what changed,
  what needs you, what could hurt value. Concise.
* **IMO / PMO** (§12.2) — operational. Scorecard, overdue work, dependencies, and
  crucially "missing updates and data-quality gaps": who has not reported.
* **Finance** (§12.3) — budget vs actual vs forecast, synergy target vs realized.
* **Workstream** (§12.4) — one stream's objectives, blockers, and what it needs from
  the others.

§12.5's design rules are applied throughout, and two of them do real work:

**"Use clear management-message titles."** A slide titled "Risks" tells a reader
nothing. A slide titled "Two critical risks are unmitigated; both need an owner today"
tells them what to do. Titles are written from the data, not from the section name.

**"Mark data-quality limitations."** Every deck ends with a slide saying what this
report could not do — files that failed, figures read off a screenshot, conflicts
nobody resolved. That slide exists because the alternative is a deck that looks
equally confident whether or not it is.
"""
from __future__ import annotations

import logging
from pathlib import Path
from typing import Optional

from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from app.models.pmi import (
    Audience,
    DataQualityReport,
    PMIDataModel,
    Severity,
    Status,
)

log = logging.getLogger("pmi.pptx")

# The Deloitte brand palette (matches `Deloitte_Master.pptx`'s theme). Kept as
# module constants because the drawn body content (tables, tiles, bullets) sets
# colours explicitly; the master's own graphics, backgrounds and logos are
# inherited by every slide automatically.
# From `app/report/brand.py` — the one place these hexes live, so the deck, the
# workbook, the charts and the HTML dashboard cannot disagree about what the
# brand green is. `brand.rgb()` hands the same hex to python-pptx.
from app.report import brand

GREEN = brand.rgb(brand.GREEN)     # Deloitte dark green — readable on white
BRIGHT = brand.rgb(brand.BRIGHT)   # Deloitte signature green (accent1)
DARK = brand.rgb(brand.DARK)       # theme dk2
GREY = brand.rgb(brand.GREY)       # theme grey
RED = brand.rgb(brand.RED)
AMBER = brand.rgb(brand.AMBER)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

#: Master layouts used when a Deloitte template is loaded. Falls back to the
#: python-pptx blank layout (index 6) when the template is absent.
COVER_LAYOUT = "Title slide - White"
CONTENT_LAYOUT = "Title Only"

RAG = {
    Severity.LOW: GREEN, Severity.MEDIUM: AMBER,
    Severity.HIGH: RGBColor(0xEF, 0x6C, 0x00), Severity.CRITICAL: RED,
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)


def generate(
    model: PMIDataModel,
    audience: Audience,
    bullets: list[str],
    out_dir: Path,
    quality: Optional[DataQualityReport] = None,
) -> Path:
    """Plan the content, then render it.

    The signature is unchanged so every existing caller and test keeps working,
    but the deck is no longer built by walking `PMIDataModel` — it is built from
    a `ReportContent`, the same structure the text preview and the Word, PDF and
    HTML renderers read. That is what lets a user see what a deck will say
    before paying to generate it, and revise it without re-extracting anything.
    """
    from app.report.planner import plan
    from app.report.render import pptx as renderer

    content = plan(model, audience, bullets=bullets, quality=quality)
    return renderer.render(content, out_dir, model)


def generate_from_content(content, out_dir: Path, model=None) -> Path:
    """Render an already-planned (and possibly user-revised) report."""
    from app.report.render import pptx as renderer

    return renderer.render(content, out_dir, model)


# ============================================================ slide primitives
# Geometry, the RAG palette and the §12.5 footer. These are the only things left
# here: what a slide *says* is decided in `app/report/planner.py` and assembled
# by `app/report/render/pptx.py`, so that the deck, the document, the PDF and the
# text preview cannot drift apart. Management-message titles moved to
# `app/report/messages.py` for the same reason.


def _table(prs, title: str, subtitle: str, headers: list[str],
           rows: list[list[str]], note: str = "", footer: str = "") -> None:
    if not rows:
        return

    slide = _blank(prs)
    _header(slide, title, subtitle)

    shape = slide.shapes.add_table(
        len(rows) + 1, len(headers),
        MARGIN, Inches(1.85),
        SLIDE_W - Inches(1.2), Inches(0.32) * (len(rows) + 1),
    )
    table = shape.table

    for column, heading in enumerate(headers):
        cell = table.cell(0, column)
        cell.text = heading
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            warn = str(value).startswith("⚠")
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
                    run.font.color.rgb = RED if warn else DARK
                    run.font.bold = warn

    if note:
        _text(slide, note, MARGIN, SLIDE_H - Inches(0.95), size=10, colour=GREY)

    _footer(slide, footer)


def _tile(slide, label: str, value: str, colour, left, top) -> None:
    box = slide.shapes.add_textbox(left, top, Inches(1.95), Inches(1.6))
    frame = box.text_frame
    frame.word_wrap = True

    run = frame.paragraphs[0].add_run()
    run.text = label
    run.font.size = Pt(11)
    run.font.color.rgb = GREY

    paragraph = frame.add_paragraph()
    run = paragraph.add_run()
    run.text = value
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = colour


def _layout(prs, name: str):
    """A master layout by name, or the blank fallback when no template is loaded."""
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    # python-pptx's default template: index 6 is the blank layout.
    return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]


def _clear_placeholders(slide):
    """Drop the layout's empty placeholders so no "Click to add…" prompt is left
    behind. The master's background graphics, logo and theme are on the *master*,
    not these placeholders, so the brand design survives their removal."""
    for placeholder in list(slide.placeholders):
        placeholder._element.getparent().remove(placeholder._element)


def _blank(prs, layout_name: str = CONTENT_LAYOUT):
    """A branded content slide: the master's background and logo, nothing else —
    titles and body are drawn on top so the deck's structure is decided by the
    planner, not by which placeholders a layout happens to carry."""
    slide = prs.slides.add_slide(_layout(prs, layout_name))
    _clear_placeholders(slide)
    return slide


def _header(slide, title: str, subtitle: str) -> None:
    _text(slide, title, MARGIN, Inches(0.45), size=22, bold=True, colour=DARK,
          width=SLIDE_W - Inches(1.2))
    if subtitle:
        _text(slide, subtitle, MARGIN, Inches(1.15), size=13, colour=GREEN)


def _footer(slide, text: str = "") -> None:
    """The footer the *plan* asked for, or nothing at all.

    This used to stamp a hard-coded prototype disclaimer on every slide,
    independently of `ReportContent.footer` — so what the deck said about its own
    status was not something the planner or the user could decide. Renderers draw;
    they never decide.
    """
    if not text:
        return
    _text(slide, text, MARGIN, SLIDE_H - Inches(0.5), size=9, colour=GREY,
          width=SLIDE_W - Inches(1.2))


def _text(slide, text: str, left, top, *, size: int, bold: bool = False,
          colour=DARK, width=None):
    box = slide.shapes.add_textbox(
        left, top, width or (SLIDE_W - Inches(1.2)), Inches(0.5)
    )
    frame = box.text_frame
    frame.word_wrap = True
    run = frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = colour
    return box
