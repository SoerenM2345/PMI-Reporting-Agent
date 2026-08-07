"""Simple file text extraction for LLM input."""
from __future__ import annotations

import logging
from pathlib import Path

log = logging.getLogger("pmi.storage.file_storage")


def get_file_text(session_id: str, filename: str) -> str:
    """Extract text from an uploaded file for LLM input.

    Supports: PDF, XLSX, CSV, DOCX, PPTX, HTML, PNG/JPG, TXT
    """
    from app.storage import json_store

    file_path = json_store.uploads_dir(session_id) / filename

    if not file_path.exists():
        raise FileNotFoundError(f"File not found: {filename}")

    suffix = file_path.suffix.lower()

    # Text files
    if suffix in {".txt", ".md"}:
        return file_path.read_text(encoding="utf-8", errors="ignore")

    # PDF
    if suffix == ".pdf":
        return _extract_pdf(file_path)

    # Excel
    if suffix in {".xlsx", ".xls"}:
        return _extract_excel(file_path)

    # CSV
    if suffix == ".csv":
        return file_path.read_text(encoding="utf-8", errors="ignore")

    # Word
    if suffix == ".docx":
        return _extract_docx(file_path)

    # PowerPoint
    if suffix == ".pptx":
        return _extract_pptx(file_path)

    # HTML
    if suffix == ".html":
        return _extract_html(file_path)

    # Images
    if suffix in {".png", ".jpg", ".jpeg", ".gif", ".bmp"}:
        return _extract_image(file_path)

    # Unknown format
    log.warning("Unknown file format: %s", suffix)
    return f"[File: {filename} - format not supported for text extraction]"


def _extract_pdf(path: Path) -> str:
    """Extract text from PDF."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(path)
        text_parts = []
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            if text.strip():
                text_parts.append(text)
        return "\n\n".join(text_parts)
    except Exception as e:
        log.error("Failed to extract PDF: %s", e)
        return f"[PDF extraction failed: {str(e)}]"


def _extract_excel(path: Path) -> str:
    """Extract text from Excel file."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"\n## Sheet: {sheet_name}")
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(v) if v is not None else "" for v in row)
                if row_text.strip():
                    parts.append(row_text)
        return "\n".join(parts)
    except Exception as e:
        log.error("Failed to extract Excel: %s", e)
        return f"[Excel extraction failed: {str(e)}]"


def _extract_docx(path: Path) -> str:
    """Extract text from Word document."""
    try:
        from docx import Document
        doc = Document(path)
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        log.error("Failed to extract DOCX: %s", e)
        return f"[Word extraction failed: {str(e)}]"


def _extract_pptx(path: Path) -> str:
    """Extract text from PowerPoint."""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            parts.append(f"\n## Slide {slide_num}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
        return "\n".join(parts)
    except Exception as e:
        log.error("Failed to extract PPTX: %s", e)
        return f"[PowerPoint extraction failed: {str(e)}]"


def _extract_html(path: Path) -> str:
    """Extract text from HTML."""
    try:
        from html.parser import HTMLParser

        class TextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.text = []
                self.skip = False

            def handle_starttag(self, tag, attrs):
                if tag in {"script", "style"}:
                    self.skip = True

            def handle_endtag(self, tag):
                if tag in {"script", "style"}:
                    self.skip = False
                elif tag in {"p", "div", "br"}:
                    self.text.append("\n")

            def handle_data(self, data):
                if not self.skip:
                    text = data.strip()
                    if text:
                        self.text.append(text)

        extractor = TextExtractor()
        extractor.feed(path.read_text(encoding="utf-8", errors="ignore"))
        return " ".join(extractor.text)
    except Exception as e:
        log.error("Failed to extract HTML: %s", e)
        return f"[HTML extraction failed: {str(e)}]"


def _extract_image(path: Path) -> str:
    """Extract text from image using vision."""
    try:
        # Use Claude Vision API to extract text from image
        from app.llm import get_client

        client = get_client()
        if not client or not client.supports_vision:
            return f"[Image: {path.name} - vision not available]"

        with open(path, "rb") as f:
            image_data = f.read()

        # For now, just indicate that the image was uploaded
        # In a real implementation, you'd call Vision API here
        return f"[Image file: {path.name}]"
    except Exception as e:
        log.error("Failed to extract image: %s", e)
        return f"[Image extraction failed: {str(e)}]"
