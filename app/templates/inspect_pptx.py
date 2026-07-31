"""Read a `.pptx` template's structure. No brand knowledge, no opinions.

This is the bottom of the template stack: it opens the file, walks the masters
and layouts, and reports what is actually in there — placeholder types, resolved
geometry, and the raw theme/master XML for `extract_theme` to parse. Everything
above it (`extract_layouts`, `layout_catalog`, `brand_system`) works from this
report rather than re-opening the file.

Two things here exist because python-pptx does not do them:

* **`idx` is not a slot identity.** `Content Placeholder 3` is `idx=10` on one
  layout, `idx=15` on another and `idx=30` on a third — in the *same* template,
  for the same visual column. Callers get `left`/`top` and must order by
  geometry. `ph_idx` is reported for round-tripping only.

* **`idx == 4294967295` is python-pptx's "this placeholder has no idx attribute"
  sentinel**, not a real index. Four layouts in the Deloitte master carry two
  such placeholders each. They are reported with `ph_idx=None` so nothing
  downstream tries `slide.placeholders[4294967295]`.
"""
from __future__ import annotations

import logging
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Emu

log = logging.getLogger("pmi.templates.inspect")

#: python-pptx reports a missing `idx` attribute as 2**32 - 1.
NO_IDX = 4294967295


@dataclass(frozen=True)
class RawPlaceholder:
    """One placeholder on a layout, with geometry resolved through inheritance."""

    ph_idx: Optional[int]
    ph_type: str
    name: str
    left_in: float
    top_in: float
    width_in: float
    height_in: float
    #: The layout's own prompt ("Click to add title"). Useful as a hint and as
    #: the string a design critic must never find in a *rendered* slide.
    prompt: str = ""

    @property
    def is_picture(self) -> bool:
        return self.ph_type in ("PICTURE", "MEDIA_CLIP")

    @property
    def is_title(self) -> bool:
        return self.ph_type in ("TITLE", "CENTER_TITLE")

    @property
    def is_addressable(self) -> bool:
        """False for the no-idx sentinel: there is no safe way to target it."""
        return self.ph_idx is not None


@dataclass(frozen=True)
class RawLayout:
    index: int
    raw_name: str
    placeholders: tuple[RawPlaceholder, ...]
    #: Shapes on the layout that are *not* placeholders — the think-cell OLE
    #: object, the `CaseCode`/`Copyright` textboxes, background pictures. These
    #: are inherited visually and are never cloned onto a slide, so they are
    #: reported for auditing, not for filling.
    decorations: tuple[str, ...] = ()

    @property
    def normalized_name(self) -> str:
        """Casefolded, whitespace-collapsed. Seven layouts have trailing spaces
        and one differs from its siblings only by the case of "Only"."""
        return " ".join(self.raw_name.split()).casefold()


@dataclass(frozen=True)
class RawTemplate:
    path: Path
    slide_w_in: float
    slide_h_in: float
    master_count: int
    layouts: tuple[RawLayout, ...]
    slide_count: int
    theme_xml: bytes = b""
    master_xml: bytes = b""
    table_style_xml: bytes = b""
    layout_xml: dict[int, bytes] = field(default_factory=dict)

    def by_index(self, index: int) -> RawLayout:
        return self.layouts[index]

    @property
    def defines_table_styles(self) -> bool:
        """The Deloitte master ships a 182-byte stub with no `<a:tblStyle>`.

        When this is False every table property must be set explicitly by the
        renderer; there is nothing to inherit.

        The stub still names the built-in "Table Grid" GUID on its
        `<a:tblStyleLst def=...>` wrapper, so a substring test for "tblStyle"
        matches an empty file. Only a real `<a:tblStyle>` child counts.
        """
        xml = self.table_style_xml
        return b"<a:tblStyle " in xml or b"<a:tblStyle>" in xml


def inspect(path: Path | str) -> RawTemplate:
    """Open `path` and report its structure. Read-only; the file is not modified."""
    path = Path(path)
    prs = Presentation(str(path))

    masters = list(prs.slide_masters)
    layouts: list[RawLayout] = []
    for index, layout in enumerate(masters[0].slide_layouts if masters else []):
        layouts.append(
            RawLayout(
                index=index,
                raw_name=layout.name,
                placeholders=tuple(_placeholder(ph) for ph in layout.placeholders),
                decorations=tuple(
                    shape.name for shape in layout.shapes if not shape.is_placeholder
                ),
            )
        )

    parts = _zip_parts(path)
    template = RawTemplate(
        path=path,
        slide_w_in=_inches(prs.slide_width),
        slide_h_in=_inches(prs.slide_height),
        master_count=len(masters),
        layouts=tuple(layouts),
        slide_count=len(prs.slides),
        theme_xml=parts.get("ppt/theme/theme1.xml", b""),
        master_xml=parts.get("ppt/slideMasters/slideMaster1.xml", b""),
        table_style_xml=parts.get("ppt/tableStyles.xml", b""),
        layout_xml={
            i: parts.get(f"ppt/slideLayouts/slideLayout{i + 1}.xml", b"")
            for i in range(len(layouts))
        },
    )
    log.info(
        "inspected %s: %d master(s), %d layouts, %d demo slides, %.3f x %.3f in",
        path.name, template.master_count, len(template.layouts),
        template.slide_count, template.slide_w_in, template.slide_h_in,
    )
    return template


def _placeholder(ph) -> RawPlaceholder:
    fmt = ph.placeholder_format
    idx = fmt.idx
    return RawPlaceholder(
        ph_idx=None if idx is None or idx == NO_IDX else int(idx),
        ph_type=str(fmt.type).split(" ")[0] if fmt.type is not None else "UNKNOWN",
        name=ph.name,
        left_in=_inches(ph.left),
        top_in=_inches(ph.top),
        width_in=_inches(ph.width),
        height_in=_inches(ph.height),
        prompt=_prompt(ph),
    )


def _prompt(ph) -> str:
    try:
        return " ".join(ph.text_frame.text.split()) if ph.has_text_frame else ""
    except (AttributeError, ValueError):                       # noqa: BLE001
        return ""


def _inches(value) -> float:
    # Placeholder geometry on a layout resolves through inheritance, but a shape
    # with no explicit extent anywhere still reports None.
    return round(Emu(value).inches, 4) if value is not None else 0.0


def _zip_parts(path: Path) -> dict[str, bytes]:
    """The XML parts `extract_theme` needs, read once.

    A template that is a valid pptx but missing a part (some generators omit
    `tableStyles.xml`) yields an empty entry rather than raising — a missing
    part means "nothing is defined here", which is exactly what the caller
    needs to know.
    """
    wanted = {"ppt/theme/theme1.xml", "ppt/slideMasters/slideMaster1.xml",
              "ppt/tableStyles.xml"}
    out: dict[str, bytes] = {}
    try:
        with zipfile.ZipFile(path) as archive:
            names = set(archive.namelist())
            for name in sorted(names):
                if name in wanted or name.startswith("ppt/slideLayouts/slideLayout"):
                    out[name] = archive.read(name)
    except (OSError, zipfile.BadZipFile) as exc:      # noqa: BLE001
        log.warning("could not read XML parts from %s (%s)", path, exc)
    return out
