"""Export a saved draft to PowerPoint, Word, PDF, Excel, HTML or Markdown (Phase 4).

Every exporter consumes the `NormalizedDoc` parsed from the *saved draft*, never
the knowledge base, so an export always matches the text the user approved and
edited — it cannot "independently regenerate a different narrative" (Scenario 6).
Markdown is the draft's own text, verbatim; the binary formats are built from the
same parsed structure with the libraries the report renderers already use, so no
new dependency is introduced.

Exporting is an **audit** action (correction #1): it logs a `draft_export` event
and never changes knowledge or the draft.
"""
from __future__ import annotations

import logging
import re
import uuid
from datetime import date
from pathlib import Path
from typing import Optional

from app.project import docmodel, paths
from app.project.docmodel import DocSection, NormalizedDoc
from app.project.json_repositories import Repositories, default_repositories
from app.project.models import AuditEvent, DraftRecord

log = logging.getLogger("pmi.project.export")

FORMATS = {
    "markdown": ".md", "md": ".md", "html": ".html",
    "word": ".docx", "docx": ".docx", "powerpoint": ".pptx", "pptx": ".pptx",
    "pdf": ".pdf", "excel": ".xlsx", "xlsx": ".xlsx",
}


class ExportError(RuntimeError):
    """The draft could not be exported (unknown draft or format)."""


def export_draft(project_id: str, draft_id: str, fmt: str, *,
                 repos: Optional[Repositories] = None,
                 chat_id: Optional[str] = None) -> Path:
    """Export the latest saved version of a draft. Returns the written file path."""
    repos = repos or default_repositories()
    draft = repos.drafts.get(project_id, draft_id)
    if draft is None:
        raise ExportError(f"No such draft: {draft_id}")

    key = (fmt or "").strip().lower()
    if key not in FORMATS:
        raise ExportError(f"Unknown export format: {fmt}")

    out_dir = paths.exports_dir(project_id)
    out_dir.mkdir(parents=True, exist_ok=True)
    stem = f"{_slug(draft.title)}_{draft.draft_id}_v{draft.version}"
    path = out_dir / f"{stem}{FORMATS[key]}"

    if key == "markdown" or key == "md":
        # The draft's own Markdown, byte for byte — the most faithful export there is.
        path.write_text(draft.content, encoding="utf-8")
    else:
        doc = docmodel.to_document(draft)
        _BUILDERS[FORMATS[key]](doc, path)

    repos.audit.append(AuditEvent(
        event_id=f"aud_{uuid.uuid4().hex[:10]}", project_id=project_id,
        chat_id=chat_id, type="draft_export",
        content=f"Exported “{draft.title}” as {key}",
        metadata={"draft_id": draft_id, "version": draft.version,
                  "file": path.name}))
    log.info("project %s: exported draft %s v%d as %s", project_id, draft_id,
             draft.version, key)
    return path


# --------------------------------------------------------------- HTML / MD
def _build_html(doc: NormalizedDoc, path: Path) -> None:
    from html import escape

    parts = ["<!doctype html><meta charset='utf-8'>",
             f"<title>{escape(doc.title)}</title>",
             "<style>body{font-family:system-ui,sans-serif;max-width:52rem;"
             "margin:2rem auto;padding:0 1rem;line-height:1.5}"
             "table{border-collapse:collapse;width:100%;margin:1rem 0}"
             "th,td{border:1px solid #ccc;padding:.4rem .6rem;text-align:left}"
             "h2{margin-top:2rem}</style>",
             f"<h1>{escape(doc.title)}</h1>"]
    for section in doc.sections:
        parts.append(f"<h2>{escape(section.heading)}</h2>")
        for block in section.blocks:
            if block.type == "paragraph" and block.text:
                parts.append(f"<p>{escape(block.text)}</p>")
            elif block.type == "note" or block.type == "chart":
                parts.append(f"<blockquote>{escape(block.text)}</blockquote>")
            elif block.type == "bullets":
                parts.append("<ul>" + "".join(
                    f"<li>{escape(i)}</li>" for i in block.items) + "</ul>")
            elif block.type == "table":
                head = "".join(f"<th>{escape(c)}</th>" for c in block.columns)
                body = "".join(
                    "<tr>" + "".join(f"<td>{escape(c)}</td>" for c in row) + "</tr>"
                    for row in block.rows)
                parts.append(f"<table><tr>{head}</tr>{body}</table>")
    path.write_text("\n".join(parts), encoding="utf-8")


# ------------------------------------------------------------------- DOCX
def _build_docx(doc: NormalizedDoc, path: Path) -> None:
    from docx import Document

    document = Document()
    document.add_heading(doc.title, level=0)
    for section in doc.sections:
        document.add_heading(section.heading or "", level=1)
        for block in section.blocks:
            if block.type == "paragraph" and block.text:
                document.add_paragraph(block.text)
            elif block.type in ("note", "chart") and block.text:
                document.add_paragraph(block.text, style="Intense Quote")
            elif block.type == "bullets":
                for item in block.items:
                    document.add_paragraph(item, style="List Bullet")
            elif block.type == "table" and block.columns:
                table = document.add_table(rows=1, cols=len(block.columns))
                table.style = "Light Grid Accent 1"
                for cell, header in zip(table.rows[0].cells, block.columns):
                    cell.text = header
                for row in block.rows:
                    cells = table.add_row().cells
                    for cell, value in zip(cells, row):
                        cell.text = value
    document.save(str(path))


# ------------------------------------------------------------------- PPTX
def _build_pptx(doc: NormalizedDoc, path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = doc.title
    if doc.subtitle:
        title_slide.placeholders[1].text = doc.subtitle

    for section in doc.sections:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
        slide.shapes.title.text = section.heading or ""
        table_block = next((b for b in section.blocks if b.type == "table"
                            and b.columns), None)
        if table_block is not None:
            _pptx_table(slide, table_block, Inches, Pt)
        else:
            _pptx_bullets(slide, section, Inches, Pt)
    prs.save(str(path))


def _pptx_bullets(slide, section: DocSection, Inches, Pt) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.6),
                                   Inches(8.8), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for block in section.blocks:
        lines = ([block.text] if block.type in ("paragraph", "note", "chart")
                 and block.text else block.items if block.type == "bullets" else [])
        for line in lines:
            para = tf.paragraphs[0] if first else tf.add_paragraph()
            para.text = line
            para.font.size = Pt(16)
            first = False
    if first:
        tf.paragraphs[0].text = "—"


def _pptx_table(slide, block, Inches, Pt) -> None:
    rows = min(len(block.rows) + 1, 12)
    cols = len(block.columns)
    shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.6),
                                   Inches(9.0), Inches(0.4 * rows))
    table = shape.table
    for c, header in enumerate(block.columns):
        table.cell(0, c).text = header
    for r, row in enumerate(block.rows[:rows - 1], start=1):
        for c in range(cols):
            table.cell(r, c).text = row[c] if c < len(row) else ""


# ------------------------------------------------------------------- XLSX
def _build_xlsx(doc: NormalizedDoc, path: Path) -> None:
    import xlsxwriter

    book = xlsxwriter.Workbook(str(path))
    bold = book.add_format({"bold": True})
    wrap = book.add_format({"text_wrap": True, "valign": "top"})

    overview = book.add_worksheet("Report")
    overview.set_column(0, 0, 100)
    overview.write(0, 0, doc.title, bold)
    row = 2
    used_names: set[str] = set()

    for section in doc.sections:
        overview.write(row, 0, section.heading, bold)
        row += 1
        for block in section.blocks:
            if block.type == "paragraph" and block.text:
                overview.write(row, 0, block.text, wrap)
                row += 1
            elif block.type == "bullets":
                for item in block.items:
                    overview.write(row, 0, f"• {item}", wrap)
                    row += 1
            elif block.type == "table" and block.columns:
                _xlsx_table_sheet(book, bold, section.heading, block, used_names)
        row += 1
    book.close()


def _xlsx_table_sheet(book, bold, heading: str, block, used: set[str]) -> None:
    name = _sheet_name(heading or "Table", used)
    sheet = book.add_worksheet(name)
    for c, header in enumerate(block.columns):
        sheet.write(0, c, header, bold)
        sheet.set_column(c, c, max(12, len(header) + 2))
    for r, row in enumerate(block.rows, start=1):
        for c, value in enumerate(row):
            sheet.write(r, c, value)


def _sheet_name(heading: str, used: set[str]) -> str:
    # Excel sheet names: <=31 chars, no []:*?/\ , unique.
    name = re.sub(r"[\[\]:*?/\\]", " ", heading).strip()[:31] or "Sheet"
    base, n = name, 2
    while name.lower() in used:
        name = f"{base[:28]} {n}"
        n += 1
    used.add(name.lower())
    return name


# -------------------------------------------------------------------- PDF
def _build_pdf(doc: NormalizedDoc, path: Path) -> None:
    from fpdf import FPDF
    from fpdf.enums import XPos, YPos

    def line(pdf, height, text):
        # Return the cursor to the left margin after each block, or the next
        # `multi_cell` sees zero remaining width and raises.
        pdf.multi_cell(0, height, text, new_x=XPos.LMARGIN, new_y=YPos.NEXT)

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 18)
    line(pdf, 9, _ascii(doc.title))
    pdf.ln(2)

    for section in doc.sections:
        pdf.set_font("Helvetica", "B", 14)
        line(pdf, 8, _ascii(section.heading))
        pdf.set_font("Helvetica", size=11)
        for block in section.blocks:
            if block.type == "paragraph" and block.text:
                line(pdf, 6, _ascii(block.text))
            elif block.type in ("note", "chart") and block.text:
                pdf.set_font("Helvetica", "I", 11)
                line(pdf, 6, _ascii(block.text))
                pdf.set_font("Helvetica", size=11)
            elif block.type == "bullets":
                for item in block.items:
                    line(pdf, 6, _ascii(f"- {item}"))
            elif block.type == "table" and block.columns:
                _pdf_table(pdf, block)
        pdf.ln(3)
    pdf.output(str(path))


def _pdf_table(pdf, block) -> None:
    epw = pdf.w - 2 * pdf.l_margin
    ncols = max(len(block.columns), 1)
    width = epw / ncols
    pdf.set_font("Helvetica", "B", 9)
    for header in block.columns:
        pdf.cell(width, 6, _ascii(header)[:22], border=1)
    pdf.ln()
    pdf.set_font("Helvetica", size=9)
    for row in block.rows[:30]:
        for c in range(ncols):
            pdf.cell(width, 6, _ascii(row[c] if c < len(row) else "")[:22], border=1)
        pdf.ln()


def _ascii(text: str) -> str:
    """FPDF's core fonts are Latin-1, and it cannot wrap a word wider than the page.

    So we drop characters the font can't encode *and* break any over-long token (a
    long citation or file path) with a space, or `multi_cell` raises "not enough
    horizontal space to render a single character".
    """
    s = (text or "").encode("latin-1", "replace").decode("latin-1")
    return re.sub(r"(\S{35})(?=\S)", r"\1 ", s)


def _slug(text: str) -> str:
    return re.sub(r"[^A-Za-z0-9]+", "_", text or "report").strip("_")[:40] or "report"


_BUILDERS = {
    ".html": _build_html, ".docx": _build_docx, ".pptx": _build_pptx,
    ".xlsx": _build_xlsx, ".pdf": _build_pdf,
}
