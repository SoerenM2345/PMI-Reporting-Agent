"""Render a `ReportContent` as a PowerPoint deck (§12).

This draws; it does not decide. Which sections exist, what they are titled and
which rows are worth showing were all settled by the planner, so this module can
be read as a straight mapping from block kind to slide furniture:

    tiles   -> the 2x3 grid of figures
    table   -> a titled table, truncated with a "showing N of M" note
    bullets -> a stack of text boxes
    chart   -> a full-bleed picture

The slide primitives (`_table`, `_tile`, `_header`, `_footer`) still live in
`app/generators/pptx_report.py` and are imported rather than copied — they carry
the geometry, the RAG palette and the §12.5 footer, none of which changes when
the content does.
"""
from __future__ import annotations

import logging
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt

from app.models.pmi import PMIDataModel
from app.report import charts_registry
from app.report.content import ReportContent, Section

log = logging.getLogger("pmi.render.pptx")


def render(
    content: ReportContent,
    out_dir: Path,
    model: Optional[PMIDataModel] = None,
) -> Path:
    """Write the deck. `model` is needed only to build chart images."""
    from app.generators import pptx_report as deck

    prs = Presentation()
    prs.slide_width, prs.slide_height = deck.SLIDE_W, deck.SLIDE_H

    _title_slide(prs, content, deck)

    for section in content.narrative():
        _section_slide(prs, section, content, out_dir, model, deck)

    out_dir.mkdir(parents=True, exist_ok=True)
    path = out_dir / (
        f"PMI_Report_{content.audience.value}_{_today()}.pptx"
    )
    prs.save(str(path))
    log.info("wrote %s (%d slides)", path.name, len(prs.slides))
    return path


def _today() -> str:
    from datetime import date

    return date.today().isoformat()


def _title_slide(prs, content: ReportContent, deck) -> None:
    slide = deck._blank(prs)
    deck._text(slide, content.title, deck.MARGIN, Inches(2.4), size=40, bold=True)
    deck._text(slide, content.subtitle, deck.MARGIN, Inches(3.4), size=20,
               colour=deck.GREY)
    if content.meta_line:
        deck._text(slide, "   ·   ".join(content.meta_line), deck.MARGIN,
                   Inches(4.1), size=13, colour=deck.GREY)
    deck._footer(slide)


def _section_slide(
    prs, section: Section, content: ReportContent, out_dir: Path,
    model: Optional[PMIDataModel], deck,
) -> None:
    """One section, one slide — except charts and tables, which own theirs."""
    # A chart section is a picture slide; nothing else goes on it.
    chart = next((b for b in section.blocks if b.kind == "chart"), None)
    if chart is not None:
        _chart_slide(prs, section, chart, out_dir, model, deck)
        return

    tables = [b for b in section.blocks if b.kind == "table" and b.rows]
    if tables:
        for block in tables:
            deck._table(
                prs, section.headline, section.label,
                [c.header for c in block.columns],
                [[cell.text for cell in row] for row in
                 block.rows[:(block.row_limit or len(block.rows))]],
                note=block.note or "",
            )
        return

    slide = deck._blank(prs)
    deck._header(slide, section.headline, section.label)

    tiles = next((b for b in section.blocks if b.kind == "tiles"), None)
    if tiles is not None:
        _tiles(slide, tiles, content, deck)
        deck._footer(slide)
        return

    top = Inches(1.9)
    rendered = False
    for block in section.blocks:
        if block.kind != "bullets" or not block.items:
            continue
        rendered = True
        if block.title:
            deck._text(slide, block.title, deck.MARGIN, top, size=15, bold=True)
            top += Inches(0.5)
        for item in block.items[:9]:
            colour = deck.RED if item.emphasis == "bad" else deck.DARK
            box = slide.shapes.add_textbox(
                deck.MARGIN, top, deck.SLIDE_W - Inches(1.2), Inches(0.5)
            )
            frame = box.text_frame
            frame.word_wrap = True
            run = frame.paragraphs[0].add_run()
            run.text = f"•  {item.text}"
            run.font.size = Pt(13)
            run.font.color.rgb = colour
            top += Inches(0.5)

    # "No critical risks" and "nobody reported the risks" are different claims;
    # an empty slide must say which one it is rather than showing white space.
    if not rendered and section.empty_explanation:
        deck._text(slide, section.empty_explanation, deck.MARGIN, Inches(2.2),
                   size=14, colour=deck.GREY)

    deck._footer(slide)


def _tiles(slide, block, content: ReportContent, deck) -> None:
    palette = {
        "good": deck.GREEN, "bad": deck.RED, "warn": deck.AMBER,
        "muted": deck.GREY, "none": deck.DARK,
    }
    for index, tile in enumerate(block.tiles):
        value = (content.facts.display(tile.fact_key)
                 if tile.fact_key else (tile.value or "Not Reported"))
        left = deck.MARGIN + Inches(2.1) * (index % 3)
        top = Inches(2.1) + Inches(1.9) * (index // 3)
        deck._tile(slide, tile.label, value, palette.get(tile.emphasis, deck.DARK),
                   left, top)


def _chart_slide(
    prs, section: Section, block, out_dir: Path,
    model: Optional[PMIDataModel], deck,
) -> None:
    builder = charts_registry.resolve(block.builder)
    if builder is None or model is None:
        # §21.17: an unknown builder is recorded, never guessed at and never
        # silently dropped without trace.
        log.warning("no chart builder named %r — slide skipped", block.builder)
        return

    try:
        path = builder(model, out_dir)
    except Exception as exc:
        log.warning("chart %s failed: %s", block.builder, exc)
        return
    if path is None:
        return

    slide = deck._blank(prs)
    deck._header(slide, section.headline, "")
    slide.shapes.add_picture(str(path), Inches(1.9), Inches(1.7),
                             height=Inches(5.2))
    deck._footer(slide)
