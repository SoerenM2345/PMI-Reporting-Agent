"""PowerPoint primitives that respect the template instead of overwriting it.

The old renderer did this, on every slide:

    slide = prs.slides.add_slide(layout)
    for placeholder in list(slide.placeholders):
        placeholder._element.getparent().remove(placeholder._element)

Deleting the placeholders is why nothing in the deck inherited the theme. A
placeholder carries the layout's typography, colour, bullet characters and
geometry by reference; remove it and you are drawing a textbox on white, which is
what every generated slide had been for as long as the feature existed.

`fill` writes *into* the placeholder and leaves the formatting alone —
deliberately setting no font name and no size, so the run inherits Aptos at
21pt from the master. `release_unused` then removes only the content
placeholders nothing was written to, because an unfilled placeholder shows
"Click to add text" in edit mode.

`consume` is the one honest exception. python-pptx cannot put a table or a chart
*inside* an `OBJECT` placeholder, so it reads that placeholder's resolved
geometry, hands it back, and deletes it. The shape is then positioned exactly
where the template said the column goes. That is "positioned by the layout" — not
the old sin of a magic coordinate.
"""
from __future__ import annotations

import logging
from typing import Iterable, Optional, Sequence

from pptx.util import Emu, Inches, Pt

from app.templates.brand_system import BrandSystem
from app.templates.extract_layouts import LayoutSlot, TemplateLayout

log = logging.getLogger("pmi.renderers.pptx_base")

Box = tuple[int, int, int, int]                # left, top, width, height in EMU


# ================================================================ the deck
def open_deck(reference):
    """The template with its demo slides removed and its masters intact."""
    from pptx import Presentation

    path = getattr(reference, "template_path", "") or ""
    if path:
        presentation = Presentation(path)
        _strip_slides(presentation)
        return presentation
    log.warning("no template path on the reference; using a blank deck")
    return Presentation()


def _strip_slides(presentation) -> None:
    """Remove the template's own example slides, keeping masters and layouts.

    The Deloitte master ships 26 empty layout swatches. They are not content and
    must not appear in a deliverable — but the layouts they demonstrate are
    exactly what this renderer uses.
    """
    slides = presentation.slides._sldIdLst
    for slide_id in list(slides):
        relationship_id = slide_id.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        if relationship_id:
            presentation.part.drop_rel(relationship_id)
        slides.remove(slide_id)


def add_page(presentation, layout: Optional[TemplateLayout]):
    """A new slide on `layout`, or on the blankest layout available."""
    native = _resolve_layout(presentation, layout)
    return presentation.slides.add_slide(native)


def _resolve_layout(presentation, layout: Optional[TemplateLayout]):
    layouts = presentation.slide_masters[0].slide_layouts
    if layout is not None:
        # Index first: it is stable and unique. Names in this template have
        # trailing spaces and inconsistent casing.
        if 0 <= layout.index < len(layouts):
            candidate = layouts[layout.index]
            if _normalize(candidate.name) == layout.normalized_name:
                return candidate
        for candidate in layouts:
            if _normalize(candidate.name) == layout.normalized_name:
                return candidate
        log.warning("layout %r is not in this template; using the first",
                    layout.raw_name)
    return layouts[6] if len(layouts) > 6 else layouts[0]


def _normalize(name: str) -> str:
    return " ".join((name or "").split()).casefold()


# ============================================================= placeholders
def slot_shape(slide, slot: Optional[LayoutSlot]):
    """The cloned placeholder for `slot`, or `None`.

    Matched by `idx`, which does survive cloning — unlike `shapes.title`, which
    returns `None` on the fourteen layouts in this template that expose only
    `BODY` placeholders.
    """
    if slot is None or slot.ph_idx is None:
        return None
    for placeholder in slide.placeholders:
        if placeholder.placeholder_format.idx == slot.ph_idx:
            return placeholder
    return None


def geometry(slide, slot: Optional[LayoutSlot], *,
             brand: Optional[BrandSystem] = None) -> Optional[Box]:
    """The slot's box in EMU, resolved through inheritance.

    A freshly cloned placeholder reports its geometry correctly even though the
    slide XML carries no `<a:off>` — python-pptx walks to the layout. Falling
    back to the catalog's measured inches covers the case where the placeholder
    was already consumed.
    """
    shape = slot_shape(slide, slot)
    if shape is not None and shape.left is not None:
        return (shape.left, shape.top, shape.width, shape.height)
    if slot is None:
        return None
    return (Inches(slot.left_in), Inches(slot.top_in),
            Inches(slot.width_in), Inches(slot.height_in))


def consume(slide, slot: Optional[LayoutSlot]) -> Optional[Box]:
    """Take the slot's geometry and remove its placeholder.

    For a table, chart or grouped diagram: python-pptx can only add those as
    free shapes, so the placeholder is replaced by something occupying exactly
    its box. Leaving it behind would show "Click to add text" underneath.
    """
    box = geometry(slide, slot)
    shape = slot_shape(slide, slot)
    if shape is not None:
        shape._element.getparent().remove(shape._element)
    return box


def fill(slide, slot: Optional[LayoutSlot], text: str, *,
         brand: Optional[BrandSystem] = None,
         size_pt: Optional[float] = None,
         color: Optional[str] = None,
         bold: Optional[bool] = None) -> bool:
    """Write `text` into the slot's placeholder. Formatting is inherited.

    Size, colour and weight are left unset unless a caller insists, so the run
    picks up the layout's own typography. A test asserts `font.size is None` on
    a filled run for exactly this reason: `None` here means "inheriting", which
    is the whole point of using the template.
    """
    shape = slot_shape(slide, slot)
    if shape is None or not shape.has_text_frame:
        return False

    from pptx.enum.text import MSO_AUTO_SIZE

    frame = shape.text_frame
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.text = text
    paragraph = frame.paragraphs[0]
    if size_pt is not None:
        paragraph.font.size = Pt(size_pt)
    if bold is not None:
        paragraph.font.bold = bold
    if color is not None and brand is not None:
        paragraph.font.color.rgb = brand.pptx_rgb(color)
    return True


def fill_bullets(slide, slot: Optional[LayoutSlot], items: Sequence[str], *,
                 brand: Optional[BrandSystem] = None,
                 size_pt: Optional[float] = None) -> bool:
    """Write bullets into a placeholder, using the master's own bullet levels."""
    shape = slot_shape(slide, slot)
    if shape is None or not shape.has_text_frame or not items:
        return False

    from pptx.enum.text import MSO_AUTO_SIZE

    frame = shape.text_frame
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.clear()
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        # Level 1 in this master carries a bullet character; level 0 does not.
        paragraph.level = 1
        if size_pt is not None:
            paragraph.font.size = Pt(size_pt)
    return True


def release_unused(slide, used_slot_ids: Iterable[str],
                   layout: Optional[TemplateLayout]) -> int:
    """Remove content placeholders nothing was written to.

    An empty placeholder renders blank in a show but shows its prompt text in
    edit mode, and a deck full of "Click to add text" is not deliverable. Title
    and subtitle slots the composition named are left alone even when empty:
    those are part of the layout's design, and removing one changes the slide's
    proportions.
    """
    if layout is None:
        return 0
    used = set(used_slot_ids)
    protected = {"title", "subtitle"}
    removed = 0

    for slot in layout.slots:
        if slot.slot_id in used or slot.slot_id in protected:
            continue
        shape = slot_shape(slide, slot)
        if shape is None:
            continue
        if shape.has_text_frame and shape.text_frame.text.strip():
            continue                       # something wrote here after all
        shape._element.getparent().remove(shape._element)
        removed += 1
    return removed


# ================================================================ furniture
def draw_footer(slide, brand: BrandSystem, *, left_text: str = "",
                page_number: Optional[int] = None,
                right_text: str = "") -> None:
    """Draw the footer band.

    This template defines **no** `FOOTER`, `SLIDE_NUMBER` or `DATE` placeholder
    on any of its 59 layouts, so there is nothing to inherit and nothing to
    switch on. The geometry mirrors the hand-drawn `CaseCode`/`Copyright` boxes
    the template's own divider layouts use, so a generated footer reads as
    native rather than bolted on.
    """
    from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN

    width = brand.slide_w_in - 2 * brand.grid.margin_x_in
    box = slide.shapes.add_textbox(
        Inches(brand.grid.margin_x_in), Inches(brand.footer_top_in),
        Inches(width), Inches(brand.footer_height_in))
    box.name = "pmi:footer"
    frame = box.text_frame
    frame.word_wrap = False
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0

    right = right_text or (str(page_number) if page_number is not None else "")
    paragraph = frame.paragraphs[0]
    paragraph.text = left_text
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.font.size = Pt(brand.footer_pt)
    paragraph.font.color.rgb = brand.pptx_rgb("muted")

    if right:
        # A second right-aligned box: one paragraph cannot hold two alignments.
        number = slide.shapes.add_textbox(
            Inches(brand.slide_w_in - brand.grid.margin_x_in - 1.0),
            Inches(brand.footer_top_in), Inches(1.0),
            Inches(brand.footer_height_in))
        number.name = "pmi:page-number"
        number.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        number_paragraph = number.text_frame.paragraphs[0]
        number_paragraph.text = right
        number_paragraph.alignment = PP_ALIGN.RIGHT
        number_paragraph.font.size = Pt(brand.footer_pt)
        number_paragraph.font.color.rgb = brand.pptx_rgb("muted")


def draw_source_note(slide, brand: BrandSystem, text: str, *,
                     top_in: Optional[float] = None) -> None:
    """The page's provenance, under the content and above the footer."""
    if not text:
        return
    top = top_in if top_in is not None else brand.footer_top_in - 0.34
    box = slide.shapes.add_textbox(
        Inches(brand.grid.margin_x_in), Inches(top),
        Inches(brand.slide_w_in - 2 * brand.grid.margin_x_in), Inches(0.3))
    box.name = "pmi:source-note"
    from pptx.enum.text import MSO_AUTO_SIZE

    frame = box.text_frame
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(6)
    paragraph.font.italic = True
    from pptx.dml.color import RGBColor
    paragraph.font.color.rgb = RGBColor(0xA6, 0xA6, 0xA6)


def draw_kpi_row(slide, brand: BrandSystem, tiles: Sequence, box: Box) -> None:
    """A row of figures, as shapes. Values were resolved from evidence already."""
    from pptx.enum.text import MSO_AUTO_SIZE

    if not tiles:
        return
    left, top, width, height = box
    gap = Inches(brand.spacing_in["md"])
    count = min(len(tiles), 6)
    tile_w = int((width - gap * (count - 1)) / count)

    for index, tile in enumerate(tiles[:count]):
        x = left + (tile_w + gap) * index
        panel = slide.shapes.add_shape(_rounded(), x, top, tile_w,
                                       min(height, Inches(1.5)))
        panel.name = f"pmi:kpi:{tile.evidence_id}"
        panel.fill.solid()
        panel.fill.fore_color.rgb = brand.pptx_rgb("surface_alt")
        panel.line.fill.background()
        panel.shadow.inherit = False

        frame = panel.text_frame
        frame.word_wrap = True
        frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        frame.text = tile.display or "Not Reported"
        value = frame.paragraphs[0]
        value.font.size = Pt(brand.font("kpi").size_pt)
        value.font.bold = True
        value.font.color.rgb = brand.pptx_rgb(_tile_color(tile, brand))

        label = frame.add_paragraph()
        label.text = tile.label
        label.font.size = Pt(brand.font("label").size_pt)
        label.font.color.rgb = brand.pptx_rgb("muted")
        if tile.note:
            note = frame.add_paragraph()
            note.text = tile.note
            note.font.size = Pt(brand.font("label").size_pt)
            note.font.italic = True
            note.font.color.rgb = brand.pptx_rgb("rag_amber")


def _tile_color(tile, brand: BrandSystem) -> str:
    return {"bad": "rag_red", "warn": "rag_amber", "good": "rag_green",
            "muted": "muted"}.get(tile.emphasis, "primary")


def _rounded():
    from pptx.enum.shapes import MSO_SHAPE

    return MSO_SHAPE.ROUNDED_RECTANGLE


# =================================================================== tables
def draw_table(slide, brand: BrandSystem, spec, box: Box,
               needs_explicit_styles: bool = True):
    """Draw a table with every property set explicitly.

    `ppt/tableStyles.xml` in this template is a 182-byte stub with no
    `<a:tblStyle>` entries, so there is no style to inherit and no banding to
    switch on. Every fill, rule and alignment below has to be stated.
    """
    from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN

    left, top, width, height = box
    rows = len(spec.rows) + 1
    columns = len(spec.columns)
    if not columns or rows < 2:
        return None

    # Give wrapped cells enough vertical room while never exceeding the slot.
    # The former fixed 0.32in rows made a two-line cell paint over every row
    # below it even though the table shape itself technically stayed on-slide.
    row_height = min(Inches(0.62), max(Inches(0.32), int(height / rows)))
    graphic = slide.shapes.add_table(rows, columns, left, top, width,
                                     row_height * rows)
    graphic.name = f"pmi:table:{spec.spec_id}"
    table = graphic.table

    for index, column in enumerate(spec.columns):
        cell = table.cell(0, index)
        cell.text = column.header
        _style_cell(cell, brand, size_pt=brand.font("small").size_pt,
                    bold=True, color="text_inverse", fill="primary")
        if column.kind in ("number", "currency", "percent"):
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    for row_index, row in enumerate(spec.rows, start=1):
        emphasised = (row_index - 1) in spec.emphasis_rows
        for column_index, cell_value in enumerate(row):
            cell = table.cell(row_index, column_index)
            cell.text = _shorten_table_cell(cell_value.text)
            column = spec.columns[column_index] if column_index < columns else None
            _style_cell(
                cell, brand, size_pt=brand.font("small").size_pt,
                bold=emphasised,
                color=_cell_color(cell_value.emphasis),
                fill="surface_alt" if row_index % 2 == 0 else "surface")
            if column is not None and column.kind in ("number", "currency",
                                                      "percent"):
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            cell.text_frame.word_wrap = True
            cell.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    return graphic


def _shorten_table_cell(text: str, limit: int = 110) -> str:
    """Keep a PowerPoint cell scannable; the full value remains in the spec."""
    compact = " ".join((text or "").replace("_x000D_", " ").split())
    if len(compact) <= limit:
        return compact
    cut = compact[:limit]
    boundary = max(cut.rfind(" "), cut.rfind(";"), cut.rfind("."))
    if boundary >= int(limit * 0.65):
        cut = cut[:boundary]
    return cut.rstrip(" ,;.") + "…"


def _cell_color(emphasis: str) -> str:
    return {"bad": "rag_red", "warn": "rag_amber", "good": "rag_green",
            "muted": "muted"}.get(emphasis, "text")


def _style_cell(cell, brand: BrandSystem, *, size_pt: float, bold: bool,
                color: str, fill: str) -> None:
    from pptx.enum.text import MSO_AUTO_SIZE

    cell.fill.solid()
    cell.fill.fore_color.rgb = brand.pptx_rgb(fill)
    cell.margin_left = cell.margin_right = Inches(0.06)
    cell.margin_top = cell.margin_bottom = Inches(0.02)
    cell.text_frame.word_wrap = True
    cell.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(size_pt)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = brand.pptx_rgb(color)


# ================================================================= textboxes
def draw_text(slide, brand: BrandSystem, text: str, box: Box, *,
              token: str = "body", color: Optional[str] = None,
              name: str = "pmi:text", italic: bool = False,
              bold: Optional[bool] = None):
    """A free textbox, for content whose placeholder was consumed."""
    left, top, width, height = box
    shape = slide.shapes.add_textbox(left, top, width, height)
    shape.name = name
    from pptx.enum.text import MSO_AUTO_SIZE

    frame = shape.text_frame
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.text = text

    style = brand.font(token)
    paragraph = frame.paragraphs[0]
    paragraph.font.size = Pt(style.size_pt)
    paragraph.font.bold = style.bold if bold is None else bold
    paragraph.font.italic = italic
    paragraph.font.color.rgb = brand.pptx_rgb(color or style.color)
    return shape


def draw_callout(slide, brand: BrandSystem, text: str, box: Box, *,
                 emphasis: str = "warn"):
    """A panel with a rule down its leading edge. Used for a caveat or an ask."""
    left, top, width, height = box
    panel = slide.shapes.add_shape(_rect(), left, top, width, height)
    panel.name = "pmi:callout"
    panel.fill.solid()
    panel.fill.fore_color.rgb = brand.pptx_rgb("surface_alt")
    panel.line.fill.background()
    panel.shadow.inherit = False

    rule = slide.shapes.add_shape(_rect(), left, top, Inches(0.05), height)
    rule.name = "pmi:callout-rule"
    rule.fill.solid()
    rule.fill.fore_color.rgb = brand.pptx_rgb(
        {"warn": "rag_amber", "bad": "rag_red",
         "good": "rag_green"}.get(emphasis, "primary"))
    rule.line.fill.background()
    rule.shadow.inherit = False

    from pptx.enum.text import MSO_AUTO_SIZE

    frame = panel.text_frame
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = Inches(0.16)
    frame.text = text
    paragraph = frame.paragraphs[0]
    paragraph.font.size = Pt(brand.font("small").size_pt)
    paragraph.font.color.rgb = brand.pptx_rgb("text")
    return panel


def _rect():
    from pptx.enum.shapes import MSO_SHAPE

    return MSO_SHAPE.RECTANGLE


def inset(box: Box, dx_in: float = 0.0, dy_in: float = 0.0) -> Box:
    left, top, width, height = box
    return (left + Inches(dx_in), top + Inches(dy_in),
            max(width - Inches(2 * dx_in), Inches(0.2)),
            max(height - Inches(2 * dy_in), Inches(0.2)))


def split_vertical(box: Box, fraction: float, gap_in: float = 0.14
                   ) -> tuple[Box, Box]:
    """Divide a slot into a top and a bottom part — a chart over its conclusion."""
    left, top, width, height = box
    gap = Inches(gap_in)
    upper = int((height - gap) * fraction)
    return ((left, top, width, upper),
            (left, top + upper + gap, width, height - upper - gap))


def to_inches(box: Box) -> tuple[float, float, float, float]:
    return tuple(round(Emu(value).inches, 4) for value in box)   # type: ignore[return-value]
