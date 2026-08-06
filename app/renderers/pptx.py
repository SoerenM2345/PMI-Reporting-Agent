"""Render a `Deliverable` as a PowerPoint deck on the Deloitte master.

The difference from what this replaces, in one sentence: slides are built *on*
the template's own layouts with content written *into* its placeholders, instead
of being drawn as free textboxes on a blank white slide after every placeholder
had been deleted.

Concretely, per page:

* the layout the visual planner chose is instantiated by index (names in this
  template carry trailing spaces and inconsistent casing, so names are not keys);
* the title and subtitle are written into real placeholders and inherit Aptos at
  the layout's own size and colour;
* a chart becomes a **native, editable** PowerPoint chart in the column slot's
  own geometry;
* a diagram becomes grouped `AutoShape`s a partner can drag;
* unused content placeholders are released, so no slide shows "Click to add
  text";

Nothing in here decides what the deck says. Every string was written and every
figure resolved before this module ran.
"""
from __future__ import annotations

import logging
from pathlib import Path
from typing import Optional, Sequence

from pptx.util import Emu, Inches

from app.context.schemas import GenerationContext
from app.deliverable.model import (
    BulletsElement,
    ChartElement,
    Deliverable,
    DiagramElement,
    ImageElement,
    KpiRowElement,
    PageDesign,
    TableElement,
    TextElement,
)
from app.renderers import naming, pptx_base
from app.renderers.common import RenderResult, MeasuredBox
from app.templates.brand_system import BrandSystem
from app.templates.extract_layouts import LayoutSlot, TemplateLayout
from app.visualizations import charts as chart_render
from app.visualizations import diagrams as diagram_render

log = logging.getLogger("pmi.renderers.pptx")

#: How much of a column a primary visual takes when it shares the slot with a
#: conclusion beneath it.
VISUAL_FRACTION = 0.72


def render(deliverable: Deliverable, context: GenerationContext,
           out_dir: Path) -> RenderResult:
    """Write the deck and report what it measured."""
    reference = context.template_reference
    brand: BrandSystem = context.brand_system or _fallback_brand()

    presentation = pptx_base.open_deck(reference)
    catalog = getattr(reference, "catalog", None)
    boxes: list[MeasuredBox] = []
    warnings: list[str] = list(deliverable.warnings)
    charts_dir = out_dir / "assets"

    for page in deliverable.pages:
        layout = catalog.by_id(page.layout_id) if catalog else None
        slide = pptx_base.add_page(presentation, layout)
        used = _render_page(slide, page, deliverable, context, brand, layout,
                            charts_dir)

        pptx_base.release_unused(slide, used, layout)
        if page.source_note and page.purpose not in ("cover", "closing"):
            pptx_base.draw_source_note(slide, brand, page.source_note)
        if page.speaker_notes:
            slide.notes_slide.notes_text_frame.text = page.speaker_notes

        boxes.extend(_measure(slide, page))

    out_dir.mkdir(parents=True, exist_ok=True)
    path = out_dir / naming.output_name(deliverable, context, "pptx")
    presentation.save(str(path))
    log.info("rendered %s: %d slides, %d distinct layouts", path.name,
             len(deliverable.pages), len(deliverable.layouts_used))

    return RenderResult(path=path, page_count=len(deliverable.pages),
                        element_boxes=boxes, warnings=warnings)


def _slot(layout: Optional[TemplateLayout],
          slot_id: str) -> Optional[LayoutSlot]:
    return layout.slot(slot_id) if layout is not None else None


def _fallback_brand() -> BrandSystem:
    from app.templates import template_registry

    return template_registry.default().brand


# ================================================================== one page
def _render_page(slide, page: PageDesign, deliverable: Deliverable,
                 context: GenerationContext, brand: BrandSystem,
                 layout: Optional[TemplateLayout], charts_dir: Path) -> list[str]:
    """Fill one slide. Returns the slot ids that received content."""
    used: list[str] = []

    if pptx_base.fill(slide, _slot(layout, "title"), page.title or deliverable.title):
        used.append("title")
    subtitle = page.subtitle or (deliverable.subtitle
                                 if page.purpose == "cover" else "")
    if subtitle and pptx_base.fill(slide, _slot(layout, "subtitle"), subtitle):
        used.append("subtitle")

    if page.purpose == "cover":
        _render_cover(slide, page, deliverable, context, brand, layout, used)
        return used

    # Group elements by the slot they were assigned, so two elements sharing a
    # column stack inside it rather than overwriting each other.
    by_slot: dict[str, list] = {}
    for element in page.elements:
        if element.role in ("headline", "kicker"):
            continue                       # already written into the furniture
        by_slot.setdefault(element.slot or "col1", []).append(element)

    for slot_id, elements in by_slot.items():
        slot = _slot(layout, slot_id)
        if _render_slot(slide, slot, slot_id, elements, page, deliverable,
                        context, brand, charts_dir):
            used.append(slot_id)
    return used


def _render_cover(slide, page: PageDesign, deliverable: Deliverable,
                  context: GenerationContext, brand: BrandSystem,
                  layout: Optional[TemplateLayout], used: list[str]) -> None:
    """The title slide.

    The template's cover puts its title at the bottom left, not centred, and
    offers a large picture placeholder. An unfilled picture placeholder shows a
    prompt, so it is released rather than left. The cover title has already
    been decided by the planner and is filled into that native placeholder by
    ``_render_page``; no second free-positioned headline belongs near the logo.
    """
    for element in page.of_role("callout"):
        box = (Inches(brand.grid.margin_x_in),
               Inches(brand.slide_h_in - 1.15),
               Inches(brand.slide_w_in / 2), Inches(0.6))
        pptx_base.draw_callout(slide, brand, element.text, box,
                               emphasis=element.emphasis or "warn")


def _render_slot(slide, slot: Optional[LayoutSlot], slot_id: str,
                 elements: Sequence, page: PageDesign, deliverable: Deliverable,
                 context: GenerationContext, brand: BrandSystem,
                 charts_dir: Path) -> bool:
    """Render every element assigned to one slot, stacked top to bottom."""
    box = pptx_base.geometry(slide, slot)
    if box is None:
        box = _default_box(brand)
    box = _reserve_source_note(box, page, brand)

    visuals = [e for e in elements
               if isinstance(e, (ChartElement, DiagramElement, TableElement,
                                 KpiRowElement, ImageElement))]
    texts = [e for e in elements if isinstance(e, (TextElement, BulletsElement))]

    if visuals and texts:
        visual_box, text_box = pptx_base.split_vertical(box, VISUAL_FRACTION)
    elif visuals:
        visual_box, text_box = box, None
    else:
        visual_box, text_box = None, box

    wrote = False
    if visuals and visual_box is not None:
        # The placeholder cannot hold a chart or a table, so its geometry is
        # taken and the placeholder removed. See `pptx_base.consume`.
        pptx_base.consume(slide, slot)
        share = _stack(visual_box, len(visuals))
        for element, sub_box in zip(visuals, share):
            wrote = _render_visual(slide, element, sub_box, page, deliverable,
                                   brand, charts_dir) or wrote

    if texts:
        if visuals:
            for element, sub_box in zip(texts, _stack(text_box, len(texts))):
                wrote = _render_text_shape(slide, element, sub_box, brand) or wrote
        else:
            # No visual took the placeholder, so write into it and inherit.
            wrote = _render_text_placeholder(slide, slot, texts, brand,
                                            text_box) or wrote
    return wrote


def _stack(box, count: int) -> list:
    if count <= 1:
        return [box]
    left, top, width, height = box
    gap = Inches(0.12)
    each = int((height - gap * (count - 1)) / count)
    return [(left, top + (each + gap) * index, width, each)
            for index in range(count)]


def _default_box(brand: BrandSystem):
    grid = brand.grid
    return (Inches(grid.margin_x_in), Inches(grid.content_top_in),
            Inches(brand.slide_w_in - 2 * grid.margin_x_in),
            Inches(grid.content_height_in))


#: Height the source note occupies, plus breathing room above it.
SOURCE_NOTE_BAND_IN = 0.38


def _reserve_source_note(box, page: PageDesign, brand: BrandSystem):
    """Shorten a content box so the provenance note has somewhere to sit.

    The template's column slots run to y=6.96 and the note is drawn just above
    the footer, so without this the last third of an inch of every column is
    written over — which is exactly where a callout or the end of a paragraph
    lands.
    """
    if not page.source_note or page.purpose in ("cover", "closing"):
        return box
    left, top, width, height = box
    note_top = Inches(brand.footer_top_in - SOURCE_NOTE_BAND_IN)
    if top + height <= note_top:
        return box
    return (left, top, width, max(note_top - top, Inches(0.4)))


# ================================================================== visuals
def _render_visual(slide, element, box, page: PageDesign,
                   deliverable: Deliverable, brand: BrandSystem,
                   charts_dir: Path) -> bool:
    if isinstance(element, ChartElement):
        return _render_chart(slide, element, box, page, deliverable, brand,
                             charts_dir)
    if isinstance(element, DiagramElement):
        return _render_diagram(slide, element, box, deliverable, brand)
    if isinstance(element, TableElement):
        spec = deliverable.specs.tables.get(element.spec_id)
        if spec is None:
            return False
        shown = _table_view_for_slide(spec, box)
        table_box = box
        if shown.has_note:
            left, top, width, height = box
            table_box = (left, top, width, max(height - Inches(0.28), Inches(1)))
        pptx_base.draw_table(slide, brand, shown, table_box)
        _caption(slide, brand, shown.note(), table_box)
        return True
    if isinstance(element, KpiRowElement):
        pptx_base.draw_kpi_row(slide, brand, element.tiles, box)
        return bool(element.tiles)
    if isinstance(element, ImageElement) and element.image_ref:
        try:
            slide.shapes.add_picture(element.image_ref, box[0], box[1],
                                     width=box[2])
            return True
        except (OSError, ValueError) as exc:                    # noqa: BLE001
            page.warnings.append(f"An image could not be placed ({exc}).")
    return False


def _table_view_for_slide(spec, box, *, row_cap: int = 8):
    """A readable slide view; the stored spec keeps every row for other formats."""
    _left, _top, _width, height = box
    by_height = max(3, int(Emu(height).inches / 0.5) - 1)
    limit = min(row_cap, by_height, spec.displayed_row_count)
    if len(spec.rows) <= limit and not spec.is_truncated:
        return spec

    shown = spec.model_copy(deep=True)
    shown.rows = shown.rows[:limit]
    shown.row_evidence_ids = shown.row_evidence_ids[:limit]
    shown.evidence_ids = shown.evidence_ids[:limit]
    shown.emphasis_rows = [row for row in shown.emphasis_rows if row < limit]
    shown.total_rows = max(shown.total_rows, len(spec.rows))
    shown.row_limit = limit
    return shown


def _render_chart(slide, element: ChartElement, box, page: PageDesign,
                  deliverable: Deliverable, brand: BrandSystem,
                  charts_dir: Path) -> bool:
    spec = deliverable.specs.charts.get(element.spec_id)
    if spec is None:
        # Should be unreachable: an unbuildable chart was downgraded to a table
        # during planning. Loud rather than silent, because the visible symptom
        # would be a caption over empty space.
        page.warnings.append(
            f"A chart was expected on this page but no validated specification "
            f"exists for {element.spec_id!r}; nothing was drawn.")
        log.error("no chart spec for %s on page %s", element.spec_id, page.page_id)
        return False

    if spec.is_native_pptx:
        try:
            chart_render.to_pptx_native(spec, brand, slide, box)
            return True
        except Exception as exc:                               # noqa: BLE001
            log.warning("native chart %s failed (%s); falling back to an image",
                        spec.spec_id, exc)
            page.warnings.append(
                "This chart could not be built as an editable PowerPoint chart "
                "and was placed as an image instead.")

    path = chart_render.to_png(spec, brand, charts_dir,
                               size_in=_png_size(box))
    slide.shapes.add_picture(str(path), box[0], box[1], width=box[2])
    return True


def _render_diagram(slide, element: DiagramElement, box,
                    deliverable: Deliverable, brand: BrandSystem) -> bool:
    spec = deliverable.specs.diagrams.get(element.spec_id)
    if spec is None:
        return False
    inches = pptx_base.to_inches(box)
    diagram_render.to_pptx(spec, brand, slide,
                           diagram_render.Box(*inches))
    return True


def _png_size(box) -> tuple[float, float]:
    width = max(Emu(box[2]).inches, 2.0)
    height = max(Emu(box[3]).inches, 1.5)
    return (round(width, 2), round(height, 2))


def _caption(slide, brand: BrandSystem, text: str, box) -> None:
    if not text:
        return
    left, top, width, height = box
    pptx_base.draw_text(slide, brand, text,
                        (left, top + height, width, Inches(0.24)),
                        token="caption", color="muted", name="pmi:caption",
                        italic=True)


# ===================================================================== text
def _render_text_placeholder(slide, slot: Optional[LayoutSlot],
                             elements: Sequence, brand: BrandSystem,
                             box) -> bool:
    """Write text into the real placeholder, so it inherits the layout's type."""
    bullets = [e for e in elements if isinstance(e, BulletsElement)]
    prose = [e for e in elements if isinstance(e, TextElement)]

    if bullets and not prose:
        return pptx_base.fill_bullets(slide, slot, bullets[0].items)
    if prose and not bullets:
        first = prose[0]
        if first.role == "callout":
            pptx_base.consume(slide, slot)
            pptx_base.draw_callout(slide, brand, first.text, box,
                                   emphasis=first.emphasis or "warn")
            return True
        if first.role == "quote":
            return pptx_base.fill(slide, slot, first.text,
                                 brand=brand, size_pt=brand.font("display").size_pt,
                                 color="deep")
        return pptx_base.fill(slide, slot, first.text)

    # Mixed prose and bullets in one column: take the box and stack them, since
    # a single placeholder cannot hold both cleanly.
    pptx_base.consume(slide, slot)
    wrote = False
    for element, sub_box in zip(elements, _stack(box, len(elements))):
        wrote = _render_text_shape(slide, element, sub_box, brand) or wrote
    return wrote


def _render_text_shape(slide, element, box, brand: BrandSystem) -> bool:
    if isinstance(element, BulletsElement):
        if not element.items:
            return False
        text = "\n".join(f"•  {item}" for item in element.items)
        pptx_base.draw_text(slide, brand, text, box, token="body",
                            name="pmi:bullets")
        return True
    if isinstance(element, TextElement) and element.text:
        if element.role == "callout":
            pptx_base.draw_callout(slide, brand, element.text, box,
                                   emphasis=element.emphasis or "warn")
        elif element.role == "quote":
            pptx_base.draw_text(slide, brand, element.text, box,
                                token="display", color="deep", name="pmi:quote")
        else:
            pptx_base.draw_text(slide, brand, element.text, box, token="body",
                                name=f"pmi:{element.role}")
        return True
    return False


# ================================================================ measuring
def _measure(slide, page: PageDesign) -> list[MeasuredBox]:
    """Record every shape's box, for the overflow critic to check later."""
    boxes: list[MeasuredBox] = []
    for shape in slide.shapes:
        if shape.left is None:
            continue
        boxes.append(MeasuredBox(
            page_id=page.page_id,
            name=shape.name,
            left_in=round(Emu(shape.left).inches, 4),
            top_in=round(Emu(shape.top).inches, 4),
            width_in=round(Emu(shape.width).inches, 4),
            height_in=round(Emu(shape.height).inches, 4),
            text=(shape.text_frame.text if shape.has_text_frame else ""),
            is_placeholder=shape.is_placeholder,
        ))
    return boxes
