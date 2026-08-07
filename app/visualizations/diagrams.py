"""Consulting diagrams as real, editable shapes.

Most of what a PMI deck needs to communicate is not quantitative: a phase plan,
a governance model, a dependency chain, a risk grid. The old system had no way
to express any of it, so those pages became tables.

In PowerPoint these are built as genuine `AutoShape`s and connectors grouped
into one `GroupShape` — not pictures of shapes. A partner can drag a box, retype
a label, or recolour a phase in the meeting. Every node carries its evidence id
in the shape name (`diagram:node:ev:risk:r_014`), which is also how targeted
regeneration finds it later.

Layout arithmetic is deterministic and lives in Python: a timeline places markers
by real date interpolation, a tree by tidy-tree positioning, a matrix by scored
coordinates with seeded jitter so a snapshot test does not flap.
"""
from __future__ import annotations

import logging
from dataclasses import dataclass
from datetime import date
from html import escape
from pathlib import Path
from typing import Optional, Sequence

from app.templates.brand_system import BrandSystem
from app.visualizations.specs import DiagramNode, DiagramSpec

log = logging.getLogger("pmi.visualizations.diagrams")

#: Prefix for shape names, so a renderer or a critic can find a node again.
NODE_PREFIX = "diagram:node:"


@dataclass(frozen=True)
class Box:
    """A rectangle in inches, relative to the page."""

    left: float
    top: float
    width: float
    height: float

    def inset(self, dx: float, dy: float) -> "Box":
        return Box(self.left + dx, self.top + dy,
                   max(self.width - 2 * dx, 0.1), max(self.height - 2 * dy, 0.1))


# ============================================================ layout arithmetic
def layout_nodes(spec: DiagramSpec, box: Box) -> dict[str, Box]:
    """Where each node sits. Pure geometry — no drawing, no library."""
    kind = spec.diagram_type
    if kind in ("process_flow", "phase_diagram"):
        return _layout_row(spec.nodes, box)
    if kind in ("timeline", "milestone_track"):
        return _layout_timeline(spec.nodes, box)
    if kind in ("risk_matrix", "two_by_two"):
        return _layout_matrix(spec.nodes, box)
    if kind == "value_driver_tree":
        return _layout_tree(spec.nodes, box)
    if kind == "swimlane":
        return _layout_swimlane(spec.nodes, box)
    if kind == "dependency_map":
        return _layout_columns(spec.nodes, box)
    return _layout_row(spec.nodes, box)


def _layout_row(nodes: Sequence[DiagramNode], box: Box) -> dict[str, Box]:
    if not nodes:
        return {}
    gap = min(0.18, box.width * 0.02)
    width = (box.width - gap * (len(nodes) - 1)) / len(nodes)
    height = min(1.2, box.height * 0.45)
    top = box.top + (box.height - height) / 2
    return {node.node_id: Box(box.left + (width + gap) * index, top, width, height)
            for index, node in enumerate(nodes)}


def _layout_timeline(nodes: Sequence[DiagramNode], box: Box) -> dict[str, Box]:
    """Markers placed by real date interpolation, not evenly spaced.

    Evenly spacing dated events is the most common way a timeline lies: three
    milestones in one week and one six months later read as a steady cadence.
    """
    dated = [(n, _parse(n.at)) for n in nodes]
    known = [(n, d) for n, d in dated if d is not None]
    if len(known) < 2:
        return _layout_row([n for n, _ in dated], box)

    earliest = min(d for _, d in known)
    latest = max(d for _, d in known)
    span = max((latest - earliest).days, 1)

    marker_w = min(1.9, box.width / max(len(nodes), 1) * 1.4)
    marker_h = 0.62
    axis_y = box.top + box.height / 2
    out: dict[str, Box] = {}

    for index, (node, when) in enumerate(dated):
        if when is None:
            # Undated items are listed under the axis rather than placed on it.
            out[node.node_id] = Box(box.left, box.top + box.height - marker_h,
                                    marker_w, marker_h)
            continue
        fraction = (when - earliest).days / span
        left = box.left + fraction * (box.width - marker_w)
        # Alternate above and below so adjacent labels do not collide.
        top = axis_y - marker_h - 0.16 if index % 2 == 0 else axis_y + 0.16
        out[node.node_id] = Box(left, top, marker_w, marker_h)
    return out


def _layout_matrix(nodes: Sequence[DiagramNode], box: Box) -> dict[str, Box]:
    """Scored items in a 5x5 grid, with deterministic jitter for collisions."""
    size = 0.52
    grid = box.inset(0.6, 0.2)
    cell_w, cell_h = grid.width / 5, grid.height / 5
    placed: dict[str, Box] = {}
    occupancy: dict[tuple[int, int], int] = {}

    for node in nodes:
        column = _clamp_axis(node.column)
        row = _clamp_axis(node.row)
        if column is None or row is None:
            continue
        key = (column, row)
        seen = occupancy.get(key, 0)
        occupancy[key] = seen + 1
        # Seeded by position in the cell, so the same input always lays out the
        # same way and a snapshot test does not flap.
        offset_x = (seen % 3 - 1) * size * 0.55
        offset_y = (seen // 3 - 1) * size * 0.55
        left = grid.left + cell_w * (column - 1) + (cell_w - size) / 2 + offset_x
        # A matrix reads with low values at the bottom.
        top = grid.top + cell_h * (5 - row) + (cell_h - size) / 2 + offset_y
        placed[node.node_id] = Box(left, top, size, size)
    return placed


def _clamp_axis(value: Optional[int]) -> Optional[int]:
    if value is None:
        return None
    return max(1, min(5, int(value)))


def _layout_tree(nodes: Sequence[DiagramNode], box: Box) -> dict[str, Box]:
    """A tidy tree: children centred under their parent, levels evenly spaced."""
    children: dict[str, list[DiagramNode]] = {}
    roots: list[DiagramNode] = []
    for node in nodes:
        if node.parent_id:
            children.setdefault(node.parent_id, []).append(node)
        else:
            roots.append(node)

    depth_of: dict[str, int] = {}

    def depth(node: DiagramNode, level: int = 0) -> int:
        depth_of[node.node_id] = level
        kids = children.get(node.node_id, [])
        return max([depth(kid, level + 1) for kid in kids], default=level)

    levels = max([depth(root) for root in roots], default=0) + 1
    node_h = min(0.72, box.height / max(levels, 1) * 0.6)
    level_gap = (box.height - node_h) / max(levels - 1, 1)

    # Leaves take consecutive slots in reading order; a parent then sits at the
    # mean of its children's slots, which is what makes the tree read as tidy.
    leaves: list[str] = []

    def place_leaves(node: DiagramNode) -> None:
        kids = children.get(node.node_id, [])
        if not kids:
            leaves.append(node.node_id)
            return
        for kid in kids:
            place_leaves(kid)

    for root in roots:
        place_leaves(root)

    centre_of: dict[str, float] = {
        node_id: float(index) for index, node_id in enumerate(leaves)}
    for node in sorted(nodes, key=lambda n: -depth_of.get(n.node_id, 0)):
        kids = children.get(node.node_id, [])
        if kids:
            centre_of[node.node_id] = sum(
                centre_of.get(kid.node_id, 0.0) for kid in kids) / len(kids)

    slots = max(len(leaves), 1)
    node_w = min(2.4, box.width / slots * 0.9)

    out: dict[str, Box] = {}
    for node in nodes:
        centre = centre_of.get(node.node_id, 0.0)
        left = box.left + (box.width / slots) * (centre + 0.5) - node_w / 2
        top = box.top + level_gap * depth_of.get(node.node_id, 0)
        out[node.node_id] = Box(max(left, box.left), top, node_w, node_h)
    return out


def _layout_swimlane(nodes: Sequence[DiagramNode], box: Box) -> dict[str, Box]:
    lanes: dict[str, list[DiagramNode]] = {}
    for node in nodes:
        lanes.setdefault(node.parent_id or "", []).append(node)
    lane_h = box.height / max(len(lanes), 1)
    out: dict[str, Box] = {}
    for lane_index, (_lane, members) in enumerate(lanes.items()):
        inner = Box(box.left + 1.2, box.top + lane_h * lane_index + 0.08,
                    box.width - 1.2, lane_h - 0.16)
        out.update(_layout_row(members, inner))
    return out


def _layout_columns(nodes: Sequence[DiagramNode], box: Box) -> dict[str, Box]:
    columns = 3
    width = box.width / columns * 0.86
    height = min(0.62, box.height / max(len(nodes) / columns, 1) * 0.7)
    out: dict[str, Box] = {}
    for index, node in enumerate(nodes):
        column, row = index % columns, index // columns
        out[node.node_id] = Box(
            box.left + (box.width / columns) * column,
            box.top + (height + 0.14) * row, width, height)
    return out


def _parse(value: Optional[str]) -> Optional[date]:
    if not value:
        return None
    try:
        return date.fromisoformat(value[:10])
    except ValueError:
        return None


# ============================================================ pptx rendering
def to_pptx(spec: DiagramSpec, brand: BrandSystem, slide, box: Box):
    """Draw as real shapes, grouped so the whole diagram moves as one."""
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.util import Inches, Pt

    positions = layout_nodes(spec, box)
    shapes: list = []
    handles: dict[str, object] = {}

    if spec.diagram_type in ("timeline", "milestone_track"):
        shapes.append(_pptx_axis(slide, brand, box))
    if spec.diagram_type in ("risk_matrix", "two_by_two"):
        shapes.extend(_pptx_matrix_grid(slide, brand, box, spec))

    shape_kind = {
        "process_flow": MSO_SHAPE.CHEVRON,
        "phase_diagram": MSO_SHAPE.PENTAGON,
        "risk_matrix": MSO_SHAPE.OVAL,
        "two_by_two": MSO_SHAPE.OVAL,
        "timeline": MSO_SHAPE.ROUNDED_RECTANGLE,
        "milestone_track": MSO_SHAPE.ROUNDED_RECTANGLE,
    }.get(spec.diagram_type, MSO_SHAPE.ROUNDED_RECTANGLE)

    for node in spec.nodes:
        position = positions.get(node.node_id)
        if position is None:
            continue
        shape = slide.shapes.add_shape(
            shape_kind, Inches(position.left), Inches(position.top),
            Inches(position.width), Inches(position.height))
        shape.name = f"{NODE_PREFIX}{node.evidence_id or node.node_id}"

        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = brand.pptx_rgb(_node_color(node))
        shape.line.color.rgb = brand.pptx_rgb("surface")
        shape.line.width = Pt(1.0)
        shape.shadow.inherit = False

        frame = shape.text_frame
        frame.word_wrap = True
        frame.text = node.label
        paragraph = frame.paragraphs[0]
        paragraph.font.size = Pt(_label_size(spec, position))
        paragraph.font.bold = True
        paragraph.font.color.rgb = brand.pptx_rgb(_on_color(node))
        if node.sublabel and position.height > 0.5:
            sub = frame.add_paragraph()
            sub.text = node.sublabel
            sub.font.size = Pt(max(_label_size(spec, position) - 2, 7))
            sub.font.color.rgb = brand.pptx_rgb(_on_color(node))

        shapes.append(shape)
        handles[node.node_id] = shape

    # Skip edges where either endpoint is excluded
    active_node_ids = {node.node_id for node in spec.nodes}
    for edge in spec.edges:
        if edge.from_id not in active_node_ids or edge.to_id not in active_node_ids:
            continue
        start, end = positions.get(edge.from_id), positions.get(edge.to_id)
        if start is None or end is None:
            continue
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.ELBOW,
            Inches(start.left + start.width / 2), Inches(start.top + start.height),
            Inches(end.left + end.width / 2), Inches(end.top))
        connector.line.color.rgb = brand.pptx_rgb("muted")
        connector.line.width = Pt(1.25)
        shapes.append(connector)

    if len(shapes) > 1:
        try:
            return slide.shapes.add_group_shape(shapes)
        except (ValueError, AttributeError) as exc:            # noqa: BLE001
            log.debug("could not group diagram %s (%s)", spec.spec_id, exc)
    return shapes[0] if shapes else None


def _pptx_axis(slide, brand: BrandSystem, box: Box):
    from pptx.util import Inches, Pt

    y = box.top + box.height / 2
    line = slide.shapes.add_shape(
        _rect(), Inches(box.left), Inches(y - 0.02), Inches(box.width),
        Inches(0.04))
    line.name = "diagram:axis"
    line.fill.solid()
    line.fill.fore_color.rgb = brand.pptx_rgb("primary")
    line.line.fill.background()
    line.shadow.inherit = False
    return line


def _rect():
    from pptx.enum.shapes import MSO_SHAPE

    return MSO_SHAPE.RECTANGLE


def _pptx_matrix_grid(slide, brand: BrandSystem, box: Box,
                      spec: DiagramSpec) -> list:
    """A 5x5 grid shaded from the sequential palette, low-left to high-right."""
    from pptx.util import Inches, Pt

    grid = box.inset(0.6, 0.2)
    cell_w, cell_h = grid.width / 5, grid.height / 5
    palette = brand.sequential or [brand.semantic["surface_alt"]]
    cells = []
    for column in range(5):
        for row in range(5):
            severity = (column + row) / 8.0
            colour = palette[min(int(severity * (len(palette) - 1)),
                                 len(palette) - 1)]
            cell = slide.shapes.add_shape(
                _rect(), Inches(grid.left + cell_w * column),
                Inches(grid.top + cell_h * (4 - row)),
                Inches(cell_w), Inches(cell_h))
            cell.name = f"diagram:cell:{column + 1}x{row + 1}"
            cell.fill.solid()
            cell.fill.fore_color.rgb = brand.pptx_rgb(colour)
            cell.line.color.rgb = brand.pptx_rgb("surface")
            cell.line.width = Pt(0.75)
            cell.shadow.inherit = False
            cells.append(cell)

    if spec.axes:
        for label, left, top, width, height in (
                (spec.axes.x_label, grid.left, grid.top + grid.height + 0.02,
                 grid.width, 0.24),
                (spec.axes.y_label, box.left, grid.top, 0.55, grid.height)):
            if not label:
                continue
            textbox = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height))
            textbox.name = "diagram:axis-label"
            textbox.text_frame.text = label
            paragraph = textbox.text_frame.paragraphs[0]
            paragraph.font.size = Pt(9)
            paragraph.font.color.rgb = brand.pptx_rgb("muted")
            cells.append(textbox)
    return cells


def _node_color(node: DiagramNode) -> str:
    return {"good": "rag_green", "warn": "rag_amber", "bad": "rag_red",
            "muted": "surface_alt"}.get(node.status, "primary")


def _on_color(node: DiagramNode) -> str:
    return "text" if node.status in ("muted",) else "text_inverse"


def _label_size(spec: DiagramSpec, position: Box) -> float:
    """Smaller type in smaller boxes. The overflow critic checks the result."""
    if position.width < 0.7:
        return 7.0
    if position.width < 1.3:
        return 8.0
    return 10.0 if spec.diagram_type != "process_flow" else 11.0


# ============================================================= svg and png
def to_svg(spec: DiagramSpec, brand: BrandSystem, *, width: int = 720,
           height: int = 320) -> str:
    """The same layout, in SVG, for the HTML dashboard."""
    box = Box(0.2, 0.2, width / 96 - 0.4, height / 96 - 0.4)
    positions = layout_nodes(spec, box)
    scale = 96.0

    parts = [f'<svg viewBox="0 0 {width} {height}" role="img" '
             f'class="pmi-diagram" data-diagram-id="{escape(spec.spec_id)}">',
             f"<title>{escape(spec.title or spec.spec_id)}</title>",
             f"<desc>{escape(spec.alt_text)}</desc>"]

    if spec.diagram_type in ("timeline", "milestone_track"):
        y = (box.top + box.height / 2) * scale
        parts.append(f'<line x1="{box.left * scale:.0f}" y1="{y:.0f}" '
                     f'x2="{(box.left + box.width) * scale:.0f}" y2="{y:.0f}" '
                     f'stroke="{brand.semantic["primary"]}" stroke-width="3"/>')

    # Skip edges where either endpoint is excluded
    active_node_ids = {node.node_id for node in spec.nodes}
    for edge in spec.edges:
        if edge.from_id not in active_node_ids or edge.to_id not in active_node_ids:
            continue
        start, end = positions.get(edge.from_id), positions.get(edge.to_id)
        if start is None or end is None:
            continue
        parts.append(
            f'<line x1="{(start.left + start.width / 2) * scale:.0f}" '
            f'y1="{(start.top + start.height) * scale:.0f}" '
            f'x2="{(end.left + end.width / 2) * scale:.0f}" '
            f'y2="{end.top * scale:.0f}" stroke="{brand.semantic["muted"]}" '
            f'stroke-width="1.5"'
            + (' stroke-dasharray="4 3"' if edge.style != "solid" else "")
            + "/>")

    for node in spec.nodes:
        position = positions.get(node.node_id)
        if position is None:
            continue
        colour = brand.color(_node_color(node))
        parts.append(
            f'<g class="pmi-node" data-evidence-id="{escape(node.evidence_id)}">'
            f'<rect x="{position.left * scale:.0f}" y="{position.top * scale:.0f}" '
            f'width="{position.width * scale:.0f}" '
            f'height="{position.height * scale:.0f}" rx="4" fill="{colour}"/>'
            f'<text x="{(position.left + position.width / 2) * scale:.0f}" '
            f'y="{(position.top + position.height / 2) * scale + 4:.0f}" '
            f'text-anchor="middle" font-size="11" font-weight="600" '
            f'fill="{brand.color(_on_color(node))}">'
            f'{escape(_clip(node.label))}</text>'
            f"<title>{escape(node.label + (f' — {node.sublabel}' if node.sublabel else ''))}</title>"
            f"</g>")

    parts.append("</svg>")
    return "\n".join(parts)


def to_png(spec: DiagramSpec, brand: BrandSystem, out_dir: Path, *,
           size_in: tuple[float, float] = (9.0, 4.0), dpi: int = 200) -> Path:
    """Rasterise for Word and PDF, from the same layout as the deck."""
    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.patches as patches
    import matplotlib.pyplot as plt

    out_dir.mkdir(parents=True, exist_ok=True)
    box = Box(0.05, 0.05, 0.9, 0.9)
    positions = layout_nodes(spec, box)

    figure, axes = plt.subplots(figsize=size_in, dpi=dpi)
    axes.set_xlim(0, 1)
    axes.set_ylim(1, 0)                    # y grows downward, as in the deck
    axes.axis("off")
    figure.patch.set_facecolor(brand.chart.surface)

    if spec.diagram_type in ("timeline", "milestone_track"):
        axes.plot([box.left, box.left + box.width],
                  [box.top + box.height / 2] * 2,
                  color=brand.semantic["primary"], linewidth=3)

    # Skip edges where either endpoint is excluded
    active_node_ids = {node.node_id for node in spec.nodes}
    for edge in spec.edges:
        if edge.from_id not in active_node_ids or edge.to_id not in active_node_ids:
            continue
        start, end = positions.get(edge.from_id), positions.get(edge.to_id)
        if start is None or end is None:
            continue
        axes.annotate("", xy=(end.left + end.width / 2, end.top),
                      xytext=(start.left + start.width / 2,
                              start.top + start.height),
                      arrowprops={"arrowstyle": "-|>",
                                  "color": brand.semantic["muted"],
                                  "linewidth": 1.2})

    for node in spec.nodes:
        position = positions.get(node.node_id)
        if position is None:
            continue
        axes.add_patch(patches.FancyBboxPatch(
            (position.left, position.top), position.width, position.height,
            boxstyle="round,pad=0.002", linewidth=0,
            facecolor=brand.color(_node_color(node))))
        axes.text(position.left + position.width / 2,
                  position.top + position.height / 2, _clip(node.label),
                  ha="center", va="center", fontsize=9, fontweight="600",
                  color=brand.color(_on_color(node)))

    if spec.title:
        axes.set_title(spec.title, loc="left", fontsize=brand.chart.title_pt,
                       color=brand.semantic["text"])
    path = out_dir / f"{_safe(spec.spec_id)}.png"
    figure.savefig(path, bbox_inches="tight", facecolor=figure.get_facecolor())
    plt.close(figure)
    return path


def _clip(text: str, limit: int = 28) -> str:
    return text if len(text) <= limit else text[:limit - 1].rstrip() + "…"


def _safe(name: str) -> str:
    import re

    return re.sub(r"[^A-Za-z0-9_.-]+", "-", name).strip("-") or "diagram"
