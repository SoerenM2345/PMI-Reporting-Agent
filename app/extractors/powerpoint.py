"""PowerPoint extractor (spec §5.2): SteerCo decks, workstream updates, governance packs.

Reads tables, text frames, grouped shapes, speaker notes and native chart data.
Speaker notes matter more than they look: the caveat a workstream lead did not want
on the slide ("ERP date is not confirmed with the vendor") is routinely in the notes,
and it is exactly the kind of thing a Steering Committee should hear about.
"""
from __future__ import annotations

import logging
from pathlib import Path
from typing import Iterator

from pptx import Presentation

from app.extractors.base import (
    classify_table,
    extract_actions_from_text,
    find_progress_mentions,
    make_source,
    rows_to_records,
)
from app.models.pmi import ExtractionMethod, SourceFormat

log = logging.getLogger("pmi.extract.pptx")

suffixes: tuple[str, ...] = (".pptx",)
format: SourceFormat = SourceFormat.POWERPOINT


def extract(path: Path) -> list[dict]:
    records: list[dict] = []
    presentation = Presentation(str(path))

    for slide_number, slide in enumerate(presentation.slides, start=1):
        title = _title_of(slide)
        table_source = make_source(
            path.name, format,
            slide_number=slide_number,
            section_name=title or None,
            extraction_method=ExtractionMethod.TABLE_PARSE,
        )
        text_source = make_source(
            path.name, format,
            slide_number=slide_number,
            section_name=title or None,
            extraction_method=ExtractionMethod.TEXT_REGEX,
        )

        text_parts: list[str] = []

        for shape in _walk(slide.shapes):
            if getattr(shape, "has_table", False):
                grid = [
                    [cell.text.strip() for cell in row.cells]
                    for row in shape.table.rows
                ]
                if len(grid) >= 2:
                    records.extend(rows_to_records(
                        grid[0], grid[1:],
                        classify_table(grid[0], context=title),
                        table_source,
                    ))

            if getattr(shape, "has_text_frame", False):
                text_parts.append(shape.text_frame.text)

            if getattr(shape, "has_chart", False):
                records.extend(_chart(shape, title, table_source))

        # Speaker notes — where the honest caveats live.
        notes = _notes_of(slide)
        if notes:
            text_parts.append(notes)
            records.append({
                "type": "note",
                "text": f"Speaker notes (slide {slide_number}): {notes[:1500]}",
                "source": text_source,
            })

        text = "\n".join(p for p in text_parts if p)
        for pct in find_progress_mentions(text):
            records.append({"type": "kpi", "name": "Overall Progress", "value": pct,
                            "unit": "%", "source": text_source})
        for item in extract_actions_from_text(text):
            records.append({**item, "source": text_source})
        if text.strip():
            records.append({"type": "note", "text": text.strip()[:2000],
                            "source": text_source})

    return records


# ------------------------------------------------------------------- internals
def _walk(shapes) -> Iterator:
    """Yield every shape, descending into groups.

    The original code iterated `slide.shapes` only, so a table inside a grouped
    shape — which is how most consultants build a slide — was invisible.
    """
    for shape in shapes:
        yield shape
        if getattr(shape, "shape_type", None) is not None and hasattr(shape, "shapes"):
            yield from _walk(shape.shapes)


def _title_of(slide) -> str:
    try:
        if slide.shapes.title is not None:
            return slide.shapes.title.text.strip()
    except (AttributeError, ValueError):
        pass
    return ""


def _notes_of(slide) -> str:
    try:
        if slide.has_notes_slide:
            return slide.notes_slide.notes_text_frame.text.strip()
    except (AttributeError, ValueError):
        pass
    return ""


def _chart(shape, title: str, source) -> list[dict]:
    """Read a native pptx chart's series data (§5.2: 'Chart data where accessible')."""
    records: list[dict] = []
    try:
        chart = shape.chart
        categories = [str(c) for c in chart.plots[0].categories]
        for series in chart.plots[0].series:
            for category, value in zip(categories, series.values):
                if value is None:
                    continue
                records.append({
                    "type": "kpi",
                    "name": f"{series.name or title or 'Chart'} — {category}",
                    "value": value,
                    "source": source,
                })
    except Exception as exc:
        log.debug("chart on a slide was not readable: %s", exc)
    return records
