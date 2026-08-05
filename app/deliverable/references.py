"""Resolve and apply user-requested reuse of uploaded source material.

The report request may name an uploaded file as a layout, table or verbatim text
source.  Resolution happens before planning and is deterministic: filenames and
locators come from the user's words, while file content is read by the native
format libraries.  The resulting constraints are stored on the Deliverable so
the preview, approval record and renderer all refer to the same source.
"""
from __future__ import annotations

import csv
import hashlib
import logging
import re
from copy import deepcopy
from pathlib import Path
from typing import Iterable, Optional

from pydantic import BaseModel, Field

from app.context.schemas import GenerationContext, SourceUseConstraint
from app.deliverable.model import (
    Deliverable,
    PageDesign,
    TableElement,
    TextElement,
)
from app.report.content import Cell, Column
from app.visualizations.specs import TableSpec

log = logging.getLogger("pmi.deliverable.references")


class ReferenceResolution(BaseModel):
    constraints: list[SourceUseConstraint] = Field(default_factory=list)
    issues: list[str] = Field(default_factory=list)


_LAYOUT_WORDS = re.compile(r"\b(layout|template|design|theme|master)\b", re.I)
_TABLE_WORDS = re.compile(r"\b(table|tabular|sheet|range)\b", re.I)
_TEXT_WORDS = re.compile(
    r"\b(exact(?:\s+wording|\s+text)?|verbatim|word[- ]for[- ]word|"
    r"copy\s+(?:the\s+)?text|use\s+(?:the\s+)?text)\b", re.I)
_REFERENCE_WORDS = re.compile(
    r"\b(?:use|using|reuse|copy|take|follow|based on|from)\b.{0,80}"
    r"\b(?:file|layout|template|table|text|wording|sheet|slide|page)\b",
    re.I | re.S,
)


def resolve(session_id: str, request_text: str,
            available_files: Optional[Iterable[str]] = None) -> ReferenceResolution:
    """Resolve source-use instructions against this session's uploaded files.

    A source filename or its unambiguous stem must occur in the request.  This
    prevents a phrase such as "use the finance table" from silently picking the
    first spreadsheet when several are present.
    """
    from app.storage import json_store

    text = request_text or ""
    names = list(available_files or _uploaded_names(session_id))
    matched: list[tuple[str, int, int]] = []
    lowered = text.casefold()
    for name in sorted(names, key=len, reverse=True):
        candidates = [name]
        stem = Path(name).stem
        if stem and stem.casefold() != name.casefold():
            candidates.append(stem)
        for candidate in candidates:
            start = lowered.find(candidate.casefold())
            if start >= 0:
                matched.append((name, start, start + len(candidate)))
                break

    constraints: list[SourceUseConstraint] = []
    for name, start, end in matched:
        clause = _reference_fragment(text, start, end)
        path = json_store.uploads_dir(session_id) / name
        checksum = _checksum(path)
        kinds: list[str] = []
        if _LAYOUT_WORDS.search(clause):
            kinds.append("layout")
        if _TABLE_WORDS.search(clause):
            kinds.append("table")
        if _TEXT_WORDS.search(clause):
            kinds.append("exact_text")

        # "Use X exactly" with no noun is most safely treated as content, not
        # as permission to clone an unknown template or table.
        if not kinds and re.search(r"\b(?:exactly|verbatim)\b", clause, re.I):
            kinds.append("exact_text")

        for kind in kinds:
            selector = _selector(clause, kind)
            constraint = SourceUseConstraint(
                kind=kind, source_file=name, selector=selector,
                mode=("adapted" if kind == "layout" and path.suffix.lower() != ".pptx"
                      else "exact"),
                scope="whole_report" if kind == "layout" else "included_content",
                source_location=selector,
                checksum=checksum,
            )
            if constraint not in constraints:
                constraints.append(constraint)

    issues: list[str] = []
    if _REFERENCE_WORDS.search(text) and not constraints:
        if not matched:
            issues.append(
                "I can use that source, but I could not match the referenced "
                "filename to an uploaded file. Please name the file exactly."
            )
        else:
            issues.append(
                "I found the file, but could not tell whether you want its "
                "layout, a table, or exact text. Please say which one."
            )

    return ReferenceResolution(constraints=constraints, issues=issues)


def apply_to_context(session_id: str, context: GenerationContext,
                     constraints: list[SourceUseConstraint]) -> list[str]:
    """Attach constraints and, for a PPTX layout source, load its real master."""
    from app.storage import json_store

    warnings: list[str] = []
    refreshed = refresh(session_id, constraints)
    context.source_use_constraints = refreshed
    layout = next((item for item in refreshed if item.kind == "layout"), None)
    if layout is None:
        return warnings

    path = json_store.uploads_dir(session_id) / layout.source_file
    if path.suffix.lower() != ".pptx":
        warnings.append(
            f"{layout.source_file} is not a PowerPoint template; its layout "
            "can only be adapted, not reused exactly."
        )
        return warnings
    try:
        from app.templates import template_registry

        reference = template_registry.load(path)
        context.template_reference = reference
        context.brand_system = reference.brand
    except Exception as exc:  # noqa: BLE001
        warnings.append(
            f"The layout in {layout.source_file} could not be loaded "
            f"({type(exc).__name__}: {exc}); the default layout was used."
        )
    return warnings


def refresh(session_id: str,
            constraints: list[SourceUseConstraint]) -> list[SourceUseConstraint]:
    """Return constraints with checksums recomputed from the current uploads."""
    from app.storage import json_store

    current: list[SourceUseConstraint] = []
    for item in constraints:
        copy = item.model_copy(deep=True)
        copy.checksum = _checksum(json_store.uploads_dir(session_id)
                                  / copy.source_file)
        current.append(copy)
    return current


def apply_content(session_id: str, deliverable: Deliverable,
                  constraints: list[SourceUseConstraint]) -> list[str]:
    """Materialise referenced tables and exact passages into normal page specs."""
    from app.storage import json_store

    warnings: list[str] = []
    deliverable.source_use_constraints = refresh(session_id, constraints)
    for item in deliverable.source_use_constraints:
        path = json_store.uploads_dir(session_id) / item.source_file
        if item.kind == "table":
            try:
                columns, rows, location = _read_table(path, item.selector)
                _append_table_pages(deliverable, item, columns, rows, location)
            except Exception as exc:  # noqa: BLE001
                warnings.append(
                    f"The requested table from {item.source_file} could not be "
                    f"included ({type(exc).__name__}: {exc})."
                )
        elif item.kind == "exact_text":
            try:
                text, location = _read_text(path, item.selector)
                if not text.strip():
                    raise ValueError("the selected location contains no readable text")
                _append_text_pages(deliverable, item, text, location)
            except Exception as exc:  # noqa: BLE001
                warnings.append(
                    f"The requested exact text from {item.source_file} could not "
                    f"be included ({type(exc).__name__}: {exc})."
                )
    deliverable.renumber()
    return warnings


def _uploaded_names(session_id: str) -> list[str]:
    from app.storage import json_store

    meta = json_store.load_meta(session_id) or {}
    return [entry if isinstance(entry, str) else entry.get("name", "")
            for entry in meta.get("files", []) if entry]


def _checksum(path: Path) -> str:
    try:
        return hashlib.sha1(path.read_bytes()).hexdigest()[:16]
    except OSError:
        return ""


def _reference_fragment(text: str, start: int, end: int) -> str:
    """The instruction surrounding one filename in a multi-source sentence."""
    delimiters = list(re.finditer(r"(?:[;,.\n]|\band\b)", text, re.I))
    left = max((match.end() for match in delimiters if match.end() <= start),
               default=0)
    right = min((match.start() for match in delimiters if match.start() >= end),
                default=len(text))
    fragment = text[left:right]
    # If the nearest fragment begins directly with the filename ("use the
    # layout from X" was separated by a comma before X), include a small amount
    # of preceding context so the governing noun is not lost.
    if not (_LAYOUT_WORDS.search(fragment) or _TABLE_WORDS.search(fragment)
            or _TEXT_WORDS.search(fragment)):
        fragment = text[max(0, start - 100):right]
    return fragment


def _selector(clause: str, kind: str) -> str:
    if kind == "table":
        sheet = re.search(r"\bsheet\s+[\"'“”]?([^,;.!]+?)[\"'“”]?\s*(?:,|range|$)",
                          clause, re.I)
        cell_range = re.search(r"\b([A-Z]{1,3}\d+:[A-Z]{1,3}\d+)\b", clause)
        parts = []
        if sheet:
            parts.append("sheet:" + " ".join(sheet.group(1).split()).strip(" '\"“”"))
        if cell_range:
            parts.append("range:" + cell_range.group(1).upper())
        location = re.search(r"\b(page|slide)\s+(\d+)\b", clause, re.I)
        if location:
            parts.append(
                f"{location.group(1).lower()}:{location.group(2)}"
            )
        return ", ".join(parts)
    location = re.search(r"\b(page|slide)\s+(\d+)\b", clause, re.I)
    return f"{location.group(1).lower()}:{location.group(2)}" if location else ""


def _read_table(path: Path, selector: str):
    suffix = path.suffix.lower()
    if suffix in (".xlsx", ".xlsm"):
        from openpyxl import load_workbook
        from openpyxl.utils.cell import range_boundaries

        book = load_workbook(path, data_only=True, read_only=True)
        sheet_name = _selector_value(selector, "sheet")
        sheet = book[sheet_name] if sheet_name in book.sheetnames else book[book.sheetnames[0]]
        cell_range = _selector_value(selector, "range")
        if cell_range:
            min_col, min_row, max_col, max_row = range_boundaries(cell_range)
            values = [[sheet.cell(row=r, column=c).value
                       for c in range(min_col, max_col + 1)]
                      for r in range(min_row, max_row + 1)]
            location = f"sheet '{sheet.title}'!{cell_range}"
        else:
            values = [list(row) for row in sheet.iter_rows(values_only=True)]
            values = _trim_matrix(values)
            location = f"sheet '{sheet.title}'"
        return _matrix_table(values, location)

    if suffix == ".csv":
        with path.open(newline="", encoding="utf-8-sig") as handle:
            values = list(csv.reader(handle))
        return _matrix_table(values, "CSV table")

    if suffix == ".docx":
        from docx import Document

        document = Document(path)
        if not document.tables:
            raise ValueError("the document contains no table")
        values = [[cell.text for cell in row.cells]
                  for row in document.tables[0].rows]
        return _matrix_table(values, "table 1")

    if suffix == ".pptx":
        from pptx import Presentation

        deck = Presentation(path)
        slide_number = int(_selector_value(selector, "slide") or 0)
        slides = ([deck.slides[slide_number - 1]] if slide_number else deck.slides)
        for slide_index, slide in enumerate(slides, start=slide_number or 1):
            for shape in slide.shapes:
                if getattr(shape, "has_table", False):
                    values = [[cell.text for cell in row.cells]
                              for row in shape.table.rows]
                    return _matrix_table(values, f"slide {slide_index}")
        raise ValueError("the presentation contains no table")

    if suffix in (".html", ".htm"):
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(path.read_text(encoding="utf-8"), "html.parser")
        table = soup.find("table")
        if table is None:
            raise ValueError("the page contains no table")
        values = [[cell.get_text(" ", strip=True) for cell in row.find_all(["th", "td"])]
                  for row in table.find_all("tr")]
        return _matrix_table(values, "table 1")

    if suffix == ".pdf":
        import pdfplumber

        page_number = int(_selector_value(selector, "page") or 1)
        with pdfplumber.open(path) as document:
            tables = document.pages[page_number - 1].extract_tables()
        if not tables:
            raise ValueError(f"page {page_number} contains no detectable table")
        return _matrix_table(tables[0], f"page {page_number}")
    raise ValueError(f"tables are not supported for {suffix or 'this file type'}")


def _matrix_table(values, location: str):
    values = _trim_matrix(values)
    if not values:
        raise ValueError("the selected table is empty")
    width = max(len(row) for row in values)
    header = list(values[0]) + [None] * (width - len(values[0]))
    columns = [Column(header=_display(value) or f"Column {index + 1}")
               for index, value in enumerate(header)]
    rows = [[Cell(text=_display(value), value=value)
             for value in (list(row) + [None] * (width - len(row)))]
            for row in values[1:]]
    return columns, rows, location


def _trim_matrix(values):
    rows = [list(row) for row in values]
    while rows and not any(value not in (None, "") for value in rows[-1]):
        rows.pop()
    if not rows:
        return []
    width = max((index + 1 for row in rows for index, value in enumerate(row)
                 if value not in (None, "")), default=0)
    return [row[:width] for row in rows]


def _display(value) -> str:
    if value is None:
        return ""
    if hasattr(value, "isoformat"):
        return value.isoformat()
    return str(value)


def _read_text(path: Path, selector: str) -> tuple[str, str]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        import fitz

        document = fitz.open(path)
        page_number = int(_selector_value(selector, "page") or 0)
        pages = [document[page_number - 1]] if page_number else list(document)
        text = "\n\n".join(page.get_text("text").rstrip() for page in pages)
        return text, f"page {page_number}" if page_number else "all pages"
    if suffix == ".pptx":
        from pptx import Presentation

        deck = Presentation(path)
        slide_number = int(_selector_value(selector, "slide") or 0)
        slides = [deck.slides[slide_number - 1]] if slide_number else list(deck.slides)
        blocks = []
        for slide in slides:
            blocks.extend(shape.text for shape in slide.shapes
                          if getattr(shape, "has_text_frame", False) and shape.text)
        return "\n\n".join(blocks), (f"slide {slide_number}" if slide_number
                                       else "all slides")
    if suffix == ".docx":
        from docx import Document

        document = Document(path)
        return "\n".join(p.text for p in document.paragraphs), "document text"
    if suffix in (".xlsx", ".xlsm"):
        columns, rows, location = _read_table(path, selector)
        lines = ["\t".join(column.header for column in columns)]
        lines.extend("\t".join(cell.text for cell in row) for row in rows)
        return "\n".join(lines), location
    if suffix == ".csv":
        return path.read_text(encoding="utf-8-sig"), "CSV text"
    if suffix in (".html", ".htm"):
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(path.read_text(encoding="utf-8"), "html.parser")
        return soup.get_text("\n", strip=True), "page text"
    return path.read_text(encoding="utf-8"), "file text"


def _selector_value(selector: str, key: str) -> str:
    match = re.search(rf"(?:^|,\s*){re.escape(key)}:([^,]+)", selector or "", re.I)
    return match.group(1).strip() if match else ""


def _append_table_pages(deliverable: Deliverable, item: SourceUseConstraint,
                        columns: list[Column], rows: list[list[Cell]],
                        location: str) -> None:
    per_page = 14
    chunks = [rows[index:index + per_page]
              for index in range(0, len(rows), per_page)] or [[]]
    for part, chunk in enumerate(chunks, start=1):
        page_id = _unique_id(deliverable, f"reference-table-{Path(item.source_file).stem}")
        spec_id = page_id + "-table"
        title = f"Referenced table — {Path(item.source_file).stem}"
        if len(chunks) > 1:
            title += f" ({part}/{len(chunks)})"
        spec = TableSpec(
            spec_id=spec_id, title=title, columns=deepcopy(columns), rows=chunk,
            total_rows=len(rows), caption=f"Exact table from {item.source_file}",
            source_note=f"Source: {item.source_file} ({location})",
        )
        deliverable.specs.tables[spec_id] = spec
        page = _reference_page(deliverable, page_id, title,
                               f"Exact source table · {location}")
        page.elements = [TableElement(
            element_id=page_id + ".table", slot=_content_slot(page),
            spec_id=spec_id, caption=spec.caption, authored_by="python",
            prominence="primary",
        )]
        page.source_note = spec.source_note
        _insert_before_closing(deliverable, page)


def _append_text_pages(deliverable: Deliverable, item: SourceUseConstraint,
                       text: str, location: str) -> None:
    chunks = _split_exact(text, 1600)
    for part, chunk in enumerate(chunks, start=1):
        page_id = _unique_id(deliverable, f"reference-text-{Path(item.source_file).stem}")
        title = f"Exact text — {Path(item.source_file).stem}"
        if len(chunks) > 1:
            title += f" ({part}/{len(chunks)})"
        page = _reference_page(deliverable, page_id, title,
                               f"Verbatim source text · {location}")
        page.elements = [TextElement(
            element_id=page_id + ".body", slot=_content_slot(page), role="body",
            text=chunk, authored_by="python", prominence="primary",
        )]
        page.source_note = f"Source: {item.source_file} ({location}) · verbatim"
        _insert_before_closing(deliverable, page)


def _reference_page(deliverable: Deliverable, page_id: str, title: str,
                    subtitle: str) -> PageDesign:
    model = next((page for page in deliverable.pages
                  if page.purpose not in ("cover", "divider", "closing")), None)
    return PageDesign(
        page_id=page_id, index=len(deliverable.pages), section_id="source_reuse",
        purpose="appendix", composition="table_full" if "table" in page_id else "single",
        layout_id=model.layout_id if model else "",
        layout_name=model.layout_name if model else "",
        title=title, subtitle=subtitle, planned_by="user",
    )


def _content_slot(page: PageDesign) -> str:
    return "col1"


def _insert_before_closing(deliverable: Deliverable, page: PageDesign) -> None:
    index = next((i for i, existing in enumerate(deliverable.pages)
                  if existing.purpose == "closing"), len(deliverable.pages))
    deliverable.pages.insert(index, page)


def _unique_id(deliverable: Deliverable, base: str) -> str:
    slug = re.sub(r"[^a-z0-9]+", "-", base.casefold()).strip("-") or "reference"
    known = {page.page_id for page in deliverable.pages}
    candidate = slug
    number = 2
    while candidate in known:
        candidate = f"{slug}-{number}"
        number += 1
    return candidate


def _split_exact(text: str, limit: int) -> list[str]:
    """Split without rewriting, trimming or dropping any source characters."""
    if len(text) <= limit:
        return [text]
    parts: list[str] = []
    start = 0
    while start < len(text):
        end = min(start + limit, len(text))
        if end < len(text):
            boundary = text.rfind("\n", start, end)
            if boundary > start:
                end = boundary + 1
        parts.append(text[start:end])
        start = end
    return parts
