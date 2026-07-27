"""Phase 4 — export a saved draft to real files.

Covers: the normalized-document parser (Markdown → structure), every exporter
produces a file that opens, and — the point of the whole separation — an export
reflects the *saved draft* including user edits, and does not re-plan a different
narrative (Scenario 6). Structure is asserted (slide/section titles), never bytes,
matching the existing generator-test discipline.
"""
from __future__ import annotations

import shutil
from pathlib import Path

import pytest

from app.project import docmodel, drafting, exporting
from app.project import files as files_mod
from app.project.json_repositories import Repositories
from app.project.rebuild import rebuild

SAMPLES = Path(__file__).resolve().parents[1] / "data" / "samples"
PROJECT = "proj_p4"


@pytest.fixture
def draft(tmp_path, monkeypatch):
    monkeypatch.setattr("app.config.settings.storage_dir", tmp_path / "storage_data")
    repos = Repositories()
    for name in ("milestone_tracker.csv", "integration_tracker.xlsx"):
        src = tmp_path / name
        shutil.copy2(SAMPLES / name, src)
        files_mod.ingest_file(PROJECT, src, repos=repos)
    rebuild(PROJECT, repos=repos, trigger="upload")
    d = drafting.create_draft(PROJECT, audience="executive", repos=repos)
    return repos, d


# ------------------------------------------------------------- normalized doc
def test_markdown_parses_into_structure(draft):
    _, d = draft
    doc = docmodel.to_document(d)
    assert doc.title == d.title
    assert len(doc.sections) == len(d.sections)
    # The milestones section should parse into a table with headers.
    milestones = next(s for s in doc.sections if "milestone" in s.section_id)
    table = next(b for b in milestones.blocks if b.type == "table")
    assert "Milestone" in table.columns and table.rows


# ------------------------------------------------------- every format opens
@pytest.mark.parametrize("fmt, suffix, opener", [
    ("markdown", ".md", None),
    ("html", ".html", None),
    ("word", ".docx", "docx"),
    ("powerpoint", ".pptx", "pptx"),
    ("excel", ".xlsx", "xlsx"),
    ("pdf", ".pdf", "pdf"),
])
def test_export_produces_a_readable_file(draft, fmt, suffix, opener):
    repos, d = draft
    path = exporting.export_draft(PROJECT, d.draft_id, fmt, repos=repos)
    assert path.suffix == suffix and path.stat().st_size > 0

    if opener == "docx":
        from docx import Document
        Document(str(path))
    elif opener == "pptx":
        from pptx import Presentation
        assert len(Presentation(str(path)).slides) >= 1
    elif opener == "xlsx":
        from openpyxl import load_workbook
        load_workbook(str(path)).close()
    elif opener == "pdf":
        import fitz
        with fitz.open(str(path)) as doc:
            assert doc.page_count >= 1


# ----------------------------------------------- fidelity to the saved draft (S6)
def test_pptx_titles_match_the_draft_sections(draft):
    repos, d = draft
    from pptx import Presentation

    path = exporting.export_draft(PROJECT, d.draft_id, "powerpoint", repos=repos)
    prs = Presentation(str(path))
    titles = {slide.shapes.title.text for slide in prs.slides
              if slide.shapes.title and slide.shapes.title.text}
    # The deck's section slides are the draft's section headings — not a re-plan.
    draft_headings = {s.heading for s in d.sections if s.heading}
    assert draft_headings & titles == draft_headings or draft_headings <= titles


def test_export_reflects_a_user_edit(draft):
    repos, d = draft
    marker = "UNIQUE-EDIT-TOKEN-42 is the agreed status."
    section_id = d.sections[0].section_id
    drafting.edit_section(PROJECT, d.draft_id, section_id, marker, repos=repos)

    md = exporting.export_draft(PROJECT, d.draft_id, "markdown", repos=repos)
    assert marker in md.read_text(encoding="utf-8")

    html = exporting.export_draft(PROJECT, d.draft_id, "html", repos=repos)
    assert "UNIQUE-EDIT-TOKEN-42" in html.read_text(encoding="utf-8")


def test_unknown_format_and_draft_raise(draft):
    repos, d = draft
    with pytest.raises(exporting.ExportError):
        exporting.export_draft(PROJECT, d.draft_id, "keynote", repos=repos)
    with pytest.raises(exporting.ExportError):
        exporting.export_draft(PROJECT, "no_such_draft", "pdf", repos=repos)


# ------------------------------------------------------ export via orchestrator
def test_orchestrator_export_returns_download(draft):
    from app.project import orchestrator

    repos, d = draft
    resp = orchestrator.respond(PROJECT, "export it as PowerPoint",
                                active_draft_id=d.draft_id, repos=repos)
    assert resp.intent == "export"
    assert any(a["type"] == "download" and a["file"].endswith(".pptx")
               for a in resp.actions)
