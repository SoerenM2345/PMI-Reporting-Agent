"""The preview loop over HTTP: plan → read as text → render (§4).

The point of these routes is that a user sees what a report will say *before*
paying to generate it, and that what they approved is what gets rendered. So the
assertions are about agreement between the preview and the artefact, and about
refusing to render a plan the analysis has outgrown.
"""
from __future__ import annotations

import pytest
from fastapi.testclient import TestClient

from app.main import app


@pytest.fixture
def client():
    return TestClient(app)


@pytest.fixture
def analyzed(client, sample_files):
    """A session carrying the spec's 82-vs-75 conflict, already analyzed."""
    sid = client.post("/api/session").json()["session_id"]
    for name in ("integration_tracker.xlsx", "weekly_update.pptx"):
        with open(sample_files / name, "rb") as handle:
            client.post(f"/api/upload?session_id={sid}",
                        files={"files": (name, handle, "application/octet-stream")})

    body = client.post("/api/analyze", json={
        "session_id": sid,
        "request_text": "Create a SteerCo presentation for the current PMI status.",
        "audience": "Executive",
    }).json()
    assert not body.get("needs_audience")
    return sid


# ===================================================================== basics
def test_asking_for_a_plan_that_does_not_exist_says_how_to_make_one(client, analyzed):
    body = client.get(f"/api/content/{analyzed}").json()
    assert body["detail"]["error"] == "no_content"
    assert "POST" in body["detail"]["message"]


def test_planning_returns_a_readable_report_before_anything_is_generated(client, analyzed):
    body = client.post(f"/api/content/{analyzed}").json()

    assert body["version"] == 1
    assert body["stale"] is False
    assert body["markdown"].startswith("# ")
    # It reads as a report, not as a dump of the data model.
    assert "## " in body["markdown"]

    # Section ids are the planner's, not a fixed table's — what is asserted is
    # that the document opens by saying something and closes by disclosing what
    # it could not say, in whatever order this request called for.
    assert body["governing_message"]
    sections = body["sections"]
    assert sections and all(s["headline"] for s in sections)
    assert body["governing_message"] in body["markdown"]
    assert any("limitation" in s["headline"].lower()
               or "not reported" in s["headline"].lower()
               or "data quality" in s["headline"].lower()
               for s in sections), \
        f"nothing discloses the report's limits: {[s['headline'] for s in sections]}"


def test_the_plan_is_versioned_and_the_history_is_listed(client, analyzed):
    client.post(f"/api/content/{analyzed}")
    second = client.post(f"/api/content/{analyzed}").json()
    assert second["version"] == 2

    versions = client.get(f"/api/content/{analyzed}/versions").json()["versions"]
    assert [v["version"] for v in versions] == [2, 1]
    assert versions[0]["is_head"] is True


def test_reverting_appends_a_version_rather_than_erasing_one(client, analyzed):
    client.post(f"/api/content/{analyzed}")
    client.post(f"/api/content/{analyzed}")

    restored = client.post(f"/api/content/{analyzed}/revert?version=1").json()
    assert restored["version"] == 3

    # Nothing was destroyed on the way.
    assert client.get(f"/api/content/{analyzed}?version=2").status_code == 200


def test_a_revision_becomes_a_new_version_and_the_old_one_stays(client, analyzed):
    """"Remove the X page" over HTTP. The previous version stays on disk, which
    is what makes a revision safe to try."""
    planned = client.post(f"/api/content/{analyzed}").json()
    target = planned["sections"][-1]

    body = client.post(f"/api/content/{analyzed}/revise",
                       json={"instruction": f"remove the {target['headline']} page"}
                       ).json()

    if not body["changed"]:
        # Refusing is a real outcome and must come back readable, not as a 4xx.
        assert body["rejected"] or body["warnings"]
        assert body["version"] == planned["version"]
        return

    assert body["version"] == planned["version"] + 1
    assert target["section_id"] not in [s["section_id"] for s in body["sections"]]
    assert client.get(f"/api/content/{analyzed}?version=1").status_code == 200


def test_an_uninterpretable_revision_says_so_rather_than_erroring(client, analyzed):
    client.post(f"/api/content/{analyzed}")

    response = client.post(f"/api/content/{analyzed}/revise",
                           json={"instruction": "make it pop"})
    body = response.json()

    assert response.status_code == 200, "a misunderstanding is not a server error"
    assert body["changed"] is False
    assert body["markdown"], "the unchanged report still comes back"


def test_revising_before_planning_says_to_plan_first(client, analyzed):
    response = client.post(f"/api/content/{analyzed}/revise",
                           json={"instruction": "put risks first"})
    assert response.status_code == 404
    assert response.json()["detail"]["error"] == "no_content"


def test_reverting_to_a_version_that_never_existed_is_a_404(client, analyzed):
    client.post(f"/api/content/{analyzed}")
    response = client.post(f"/api/content/{analyzed}/revert?version=99")
    assert response.status_code == 404


# ================================================================== staleness
def test_resolving_a_conflict_marks_the_stored_plan_stale(client, analyzed):
    """The dangerous sequence. The plan states the figure one source reported;
    the user then decides the other one is right. Serving that plan unflagged
    would hand over a deck the user has already corrected."""
    client.post(f"/api/content/{analyzed}")
    assert client.get(f"/api/content/{analyzed}").json()["stale"] is False

    conflicts = client.get(f"/api/conflicts/{analyzed}").json()["unresolved"]
    if not conflicts:
        pytest.skip("sample data produced no unresolved conflict to decide")

    client.post(f"/api/conflicts/{analyzed}/resolve",
                json={"choices": {conflicts[0]["conflict_id"]: {"value": "80"}}})

    assert client.get(f"/api/content/{analyzed}").json()["stale"] is True


def test_generation_refuses_a_stale_plan_and_uses_a_fresh_one(client, analyzed):
    """A stale plan must never reach a renderer — it is silently wrong, which
    is the worst failure mode this system has."""
    client.post(f"/api/content/{analyzed}")

    conflicts = client.get(f"/api/conflicts/{analyzed}").json()["unresolved"]
    for conflict in conflicts:
        client.post(f"/api/conflicts/{analyzed}/resolve",
                    json={"choices": {conflict["conflict_id"]: {"value": "80"}}})

    body = client.post("/api/generate",
                       json={"session_id": analyzed, "force": True}).json()

    assert body["outputs"], "generation produced nothing"
    # Re-planned rather than reusing the stale version.
    runs = client.get(f"/api/session/{analyzed}").json().get("runs", [])
    assert runs and runs[-1].get("content_version") is None


def _summary_text(content: dict) -> str:
    """The executive-summary bullets from a `/api/content` payload, as one string."""
    for section in content.get("blocks", []):
        if section["section_id"] != "summary.executive":
            continue
        return "\n".join(
            item["text"]
            for block in section["blocks"]
            for item in block.get("items", [])
        )
    return ""


def test_resolving_a_conflict_reaches_the_summary_narrative(client, analyzed):
    """The executive summary is authored free text — the one place a figure the
    report states is *asserted* rather than derived. It used to be reused verbatim
    across every re-plan with no staleness check, so after the user settled the
    82-vs-75 conflict the deterministic sections self-corrected while the summary
    kept calling it unresolved. Re-planning must re-roll it: the losing narrative
    goes, and the value the user chose appears."""
    before = _summary_text(client.post(f"/api/content/{analyzed}").json())
    if "UNRESOLVED" not in before.upper():
        pytest.skip("this run's summary did not flag an unresolved conflict")

    conflicts = client.get(f"/api/conflicts/{analyzed}").json()["unresolved"]
    for conflict in conflicts:
        client.post(f"/api/conflicts/{analyzed}/resolve",
                    json={"choices": {conflict["conflict_id"]: {"value": "80"}}})

    after = _summary_text(client.post(f"/api/content/{analyzed}").json())
    assert "UNRESOLVED" not in after.upper(), \
        "the re-planned summary still calls a resolved conflict unresolved"
    assert "80" in after, "the value the user chose never reached the summary"


# ============================================== the preview matches the deck
def test_what_the_preview_says_is_what_the_deck_says(client, analyzed):
    """The whole promise of the preview. If these can disagree, approving the
    text means nothing."""
    from pptx import Presentation

    from app.config import get_settings

    preview = client.post(f"/api/content/{analyzed}").json()
    body = client.post("/api/generate",
                       json={"session_id": analyzed, "force": True}).json()

    deck_name = next(name for name in body["outputs"] if name.endswith(".pptx"))
    path = get_settings().output_dir / analyzed / deck_name

    slide_text = []
    for slide in Presentation(str(path)).slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text.append(shape.text_frame.text)
    deck = "\n".join(slide_text)

    # Every section headline the user approved appears on a slide.
    for section in preview["sections"]:
        if section["section_id"] == "summary.executive":
            continue
        assert section["headline"] in deck, \
            f"preview promised {section['headline']!r}, deck does not say it"


# ================================================ editing the preview (item 3b)
def _planned(session_id):
    """The stored `Deliverable` — the one object the preview and deck share."""
    from app.deliverable import session as session_plan

    return session_plan.load(session_id)


def _cells(session_id):
    """Every `(block_id, row, column, cell)` the preview exposes in a table.

    A table's cells live on its `TableSpec`, reached through the element's
    `spec_id`; the id the UI sends back is the *element's*, which is why the
    walk goes element -> spec rather than assuming the two are the same.
    """
    from app.deliverable.model import TableElement

    deliverable = _planned(session_id)
    for page in deliverable.pages:
        for element in page.elements:
            if not isinstance(element, TableElement):
                continue
            spec = deliverable.specs.tables.get(element.spec_id)
            if spec is None:
                continue
            for r, row in enumerate(spec.rows):
                for c, cell in enumerate(row):
                    yield element.element_id, r, c, cell


def _first_editable_cell(session_id):
    """A `(block_id, row, column, cell)` the preview would let a user edit."""
    for block_id, row, column, cell in _cells(session_id):
        if cell.ref is not None and cell.ref.field:
            return block_id, row, column, cell
    return None


def test_a_preview_edit_is_written_into_the_data_model(client, analyzed):
    """A preview edit must change the *data*, not just the stored plan.

    Patching the stored `Deliverable` alone is one line and updates the preview
    immediately — and produces a deck, a workbook and a document that disagree
    with the preview about a figure the user personally corrected. Everything is
    planned from `PMIDataModel`, so that is where the value has to land.
    """
    from app.storage import json_store

    client.post(f"/api/content/{analyzed}")
    found = _first_editable_cell(analyzed)
    if found is None:
        pytest.skip("no editable cell in this plan")
    block_id, row, column, cell = found

    body = client.post(f"/api/content/{analyzed}/cell", json={
        "block_id": block_id, "row": row, "column": column,
        "value": "Anna Schmidt" if cell.ref.field.endswith(("owner", "lead"))
                 else "12-08-2026",
    }).json()

    if not body["applied"]:
        pytest.skip(f"this cell rejected the value: {body['message']}")

    model = json_store.load_analysis(analyzed).data_model
    collection, id_attr = _collection(cell.ref.entity_type)
    entity = next(e for e in getattr(model, collection)
                  if getattr(e, id_attr) == cell.ref.entity_id)
    assert getattr(entity, cell.ref.field) is not None

    # And the report was re-planned, so the preview reflects it too.
    assert body["version"] > 1


def test_a_preview_edit_reaches_the_excel_output_too(client, analyzed):
    """The formats must not be able to disagree about a corrected figure.

    The workbook used to re-walk `PMIDataModel` instead of rendering the
    approved content, so it was the one format that could contradict the preview
    — and nothing in the system compared them, so nobody would have found out
    until the meeting.
    """
    from openpyxl import load_workbook

    from app.config import get_settings

    client.post(f"/api/content/{analyzed}")
    found = _first_editable_cell(analyzed)
    if found is None:
        pytest.skip("no editable cell in this plan")
    block_id, row, column, cell = found

    marker = "Wilhelmina Ravensworth"      # unmistakably not from any file
    if not cell.ref.field.endswith(("owner", "lead", "title", "name")):
        pytest.skip("the first editable cell does not hold free text")

    body = client.post(f"/api/content/{analyzed}/cell", json={
        "block_id": block_id, "row": row, "column": column, "value": marker,
    }).json()
    assert body["applied"], body["message"]
    assert marker in body["markdown"], "the preview does not show the edit"

    generated = client.post("/api/generate", json={
        "session_id": analyzed, "format": "excel", "force": True,
    }).json()
    name = next(n for n in generated["outputs"] if n.endswith(".xlsx"))
    workbook = load_workbook(str(get_settings().output_dir / analyzed / name))

    text = "\n".join(
        str(c.value)
        for sheet in workbook.worksheets
        for r in sheet.iter_rows() for c in r if c.value
    )
    assert marker in text, "the workbook does not carry the corrected value"


def test_a_computed_cell_says_why_it_cannot_be_edited(client, analyzed):
    """A risk score and a budget variance are ours, computed in Python (§11).

    Accepting an edit to one would silently overwrite arithmetic the system
    guarantees it does itself — so the cell carries no `ref` and the endpoint
    says what to correct instead.
    """
    client.post(f"/api/content/{analyzed}")

    target = next(((block_id, row, column)
                   for block_id, row, column, cell in _cells(analyzed)
                   if cell.ref is None or not cell.ref.field), None)
    if target is None:
        pytest.skip("every cell in this plan is editable")

    body = client.post(f"/api/content/{analyzed}/cell", json={
        "block_id": target[0], "row": target[1], "column": target[2],
        "value": "99",
    }).json()

    assert body["applied"] is False
    assert "computed" in body["message"]


def _collection(entity_type: str):
    from app.agent.corrections import _COLLECTIONS

    return _COLLECTIONS[entity_type]


def test_generation_still_works_for_a_caller_that_never_opens_the_preview(client, analyzed):
    """The preview is optional. Every pre-existing client must keep working."""
    body = client.post("/api/generate",
                       json={"session_id": analyzed, "force": True}).json()
    assert body["outputs"]
    assert body["summary"] is not None


# ================================================== re-rendering into formats
@pytest.mark.parametrize("fmt,suffix", [
    ("word", ".docx"), ("pdf", ".pdf"), ("html", ".html"),
])
def test_the_approved_content_can_be_rendered_into_any_format(
    client, analyzed, fmt, suffix
):
    """§17. The user reads the text once, then picks a format — re-rendering
    costs no LLM call and cannot change a word of what they approved."""
    client.post(f"/api/content/{analyzed}")

    body = client.post("/api/generate", json={
        "session_id": analyzed, "force": True, "format": fmt,
    }).json()

    assert any(name.endswith(suffix) for name in body["outputs"]), body["outputs"]
    assert not body["errors"], body["errors"]


def test_a_rendered_file_can_actually_be_downloaded(client, analyzed):
    client.post(f"/api/content/{analyzed}")
    body = client.post("/api/generate", json={
        "session_id": analyzed, "force": True, "format": "pdf",
    }).json()

    name = next(n for n in body["outputs"] if n.endswith(".pdf"))
    response = client.get(f"/api/download/{analyzed}/{name}")

    assert response.status_code == 200
    assert response.content[:4] == b"%PDF"
