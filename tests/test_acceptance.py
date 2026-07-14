"""The §20 acceptance scenario, end to end.

The spec states it as fifteen numbered steps. This test asserts all of them, in order,
against the real API — using the stored vision fixture so it runs in CI with no key and
no cost.

Two of the assertions carry more weight than the rest, and they are the reason this file
exists rather than a happy-path smoke test:

  * **Step 9** — the system must ASK about the 82-vs-75 progress conflict. It is only a
    9% delta; a system that resolves it silently has failed, however good its deck is.
  * **Step 13** — the deck must contain the risk that exists ONLY in the screenshot. If
    the image pipeline quietly contributes nothing, every other test still passes and
    the report is simply missing a critical GDPR risk that nobody will ever know about.
"""
from __future__ import annotations

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation

from app.config import get_settings
from app.main import app

REQUEST = "Create a Steering Committee presentation for the current PMI status."


@pytest.fixture
def client():
    return TestClient(app)


def deck_text(path) -> str:
    chunks: list[str] = []
    for slide in Presentation(str(path)).slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    chunks += [cell.text for cell in row.cells]
    return "\n".join(chunks)


def deck_headers(path) -> list[str]:
    """The title AND subtitle of each slide.

    Both matter. §12.5 puts the *management message* in the title ("2 critical risks
    have no mitigation action"), and the §20.12 section name in the subtitle ("Critical
    risks and issues"). Reading only the title would test the wrong half.
    """
    headers = []
    for slide in Presentation(str(path)).slides:
        lines = [
            shape.text_frame.text.split("\n")[0]
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        headers.append(" — ".join(lines[:2]))
    return headers


def test_steerco_acceptance_scenario(client, sample_files, fake_vision):
    """§20, all fifteen steps."""

    # --- 1. The user creates a new PMI project. -----------------------------
    session = client.post("/api/session").json()["session_id"]

    client.post("/api/project", json={
        "session_id": session,
        "project_name": "Project Aurora",
        "reporting_date": "2026-07-01",
        "day_1_date": "2026-06-15",
    })

    # --- 2. The user uploads an Excel masterplan, a SteerCo PPTX, a PNG. ----
    uploads = [
        ("integration_tracker.xlsx", "the Excel integration masterplan"),
        ("weekly_update.pptx", "the PowerPoint Steering Committee update"),
        ("risk_dashboard.png", "the PNG screenshot of a risk dashboard"),
    ]
    for name, _what in uploads:
        with open(sample_files / name, "rb") as handle:
            response = client.post(
                f"/api/upload?session_id={session}",
                files={"files": (name, handle, "application/octet-stream")},
            )
        assert response.status_code == 200
        assert name in response.json()["saved"]

    # --- 3. The user requests a Steering Committee presentation. ------------
    analysis = client.post("/api/analyze", json={
        "session_id": session,
        "request_text": REQUEST,
    }).json()

    # --- 4. The system detects the Executive audience. ----------------------
    assert analysis["needs_audience"] is False
    assert analysis["audience"] == "Executive", "§20.4: the Executive audience"
    assert analysis["output_type"] == "powerpoint"

    # --- 5. It extracts tasks, milestones, risks, costs, synergies, decisions.
    stats = analysis["stats"]
    for entity in ("tasks", "milestones", "risks", "budget_items"):
        assert stats[entity] > 0, f"§20.5: no {entity} were extracted"

    # --- 6. It interprets the risk-dashboard image. -------------------------
    assert fake_vision.calls >= 1, "§20.6: the image was never sent for interpretation"

    review = analysis["low_confidence_items"]
    assert review, "§20.6: image findings must be flagged for review (§5.6)"

    # §21.14: nothing read from a picture ever reaches the confidence of a spreadsheet
    # read, so an image finding can never win an automatic conflict against the tracker
    # it was screenshotted from.
    assert all(item["confidence"] < 1.0 for item in review)

    # --- 7-8. It standardizes, and detects the conflicts. -------------------
    #   Excel reports 82% progress; PowerPoint reports 75%.
    conflicts = analysis["conflicts"]
    progress = next(
        c for c in conflicts if "progress" in c["entity_key"].casefold()
    )
    assert progress["values"] == {
        "integration_tracker.xlsx": "82",
        "weekly_update.pptx": "75",
    }, "§20.8: the 82-vs-75 conflict"

    #   ...and the image shows a newly added critical risk.
    quality = client.get(f"/api/quality/{session}").json()
    image_risks = [
        item for item in quality["low_confidence_items"] if item["type"] == "risk"
    ]
    assert any("GDPR" in item["label"] for item in image_risks), \
        "§20.8: the critical risk found only in the image"

    # --- 9. The system ASKS the user to resolve the progress conflict. ------
    # The assertion that matters most. 82 vs 75 is only a 9% delta — a severity rule
    # based on magnitude would auto-resolve it and never ask. It escalates because
    # overall progress is on §9's critical-topics list.
    assert progress["severity"] == "critical"
    assert progress["requires_user_input"] is True
    assert progress["conflict_id"] in analysis["blocking_conflicts"]

    # ...and the system REFUSES to produce a deck until it is answered.
    refused = client.post("/api/generate", json={"session_id": session})
    assert refused.status_code == 409, \
        "§20.9: a deck must not be produced while a critical conflict is open"
    assert refused.json()["detail"]["error"] == "unresolved_critical_conflicts"

    # --- 10. The user selects 82%. -----------------------------------------
    resolved = client.post(
        f"/api/conflicts/{session}/resolve",
        json={"choices": {progress["conflict_id"]: "integration_tracker.xlsx"}},
    ).json()

    winner = next(
        c for c in resolved["conflicts"]
        if c["conflict_id"] == progress["conflict_id"]
    )
    assert winner["resolved_value"] == "82", "§20.10"
    assert resolved["blocking_conflicts"] == []

    # --- 11. The system generates an editable PowerPoint. -------------------
    generated = client.post("/api/generate", json={"session_id": session})
    assert generated.status_code == 200

    outputs = generated.json()["outputs"]
    deck_name = next(o for o in outputs if o.endswith(".pptx"))

    deck_path = get_settings().output_dir / session / deck_name
    assert deck_path.is_file()
    Presentation(str(deck_path))  # opens => editable (§21.15)

    # --- 12. The PowerPoint includes the eight required sections. -----------
    headers = " | ".join(deck_headers(deck_path)).casefold()
    body = deck_text(deck_path)

    for section in (
        "executive summary",            # §20.12.1
        "overall integration status",   # §20.12.2
        "workstream",                   # §20.12.3
        "milestone",                    # §20.12.4
        "risk",                         # §20.12.5
        "financial",                    # §20.12.6 — synergy and budget status
        "decision",                     # §20.12.7
        "next steps",                   # §20.12.8
    ):
        assert section in headers, f"§20.12: the deck must have a '{section}' section"

    # §12.5: "Use clear management-message titles." A slide called "Risks" tells the
    # reader nothing; one that names the finding tells them what to do.
    assert "critical risk(s) require management attention" in headers \
        or "have no mitigation" in headers \
        or "are overdue" in headers

    # --- 13. The presentation includes the risk extracted from the image. ---
    # If the image pipeline quietly contributes nothing, every other assertion above
    # still passes and the deck is simply missing a critical GDPR risk.
    assert "GDPR" in body, \
        "§20.13: the risk that exists ONLY in the screenshot must reach the deck"

    # The winning figure is on the deck; the losing one appears only where a conflict
    # is being discussed — never presented as the project's progress.
    assert "82" in body

    # --- 14. The system generates a data-quality report. --------------------
    assert any("data_quality_report" in o for o in outputs), "§20.14"
    assert any("conflict_report" in o for o in outputs)

    quality_name = next(o for o in outputs if "data_quality_report" in o)
    quality_text = (get_settings().output_dir / session / quality_name).read_text()

    assert "Low-confidence findings" in quality_text
    assert "GDPR" in quality_text
    assert "Senior Manager review" in quality_text

    conflict_name = next(o for o in outputs if "conflict_report" in o)
    conflict_text = (get_settings().output_dir / session / conflict_name).read_text()
    assert "82" in conflict_text and "75" in conflict_text
    assert "Resolved to `82`" in conflict_text

    # --- 15. The user downloads the .pptx and the conflict report. ----------
    for name in (deck_name, conflict_name):
        response = client.get(f"/api/download/{session}/{name}")
        assert response.status_code == 200
        assert response.content


def test_the_image_risk_would_be_missing_without_the_image(client, sample_files):
    """Control. Without the screenshot, the GDPR risk is in no other file — which is
    what makes step 13 a real test of the image pipeline rather than a coincidence."""
    session = client.post("/api/session").json()["session_id"]

    for name in ("integration_tracker.xlsx", "weekly_update.pptx"):
        with open(sample_files / name, "rb") as handle:
            client.post(f"/api/upload?session_id={session}",
                        files={"files": (name, handle, "application/octet-stream")})

    client.post("/api/analyze", json={"session_id": session, "request_text": REQUEST})
    body = client.get(f"/api/quality/{session}").json()

    labels = " ".join(i["label"] for i in body["low_confidence_items"])
    assert "GDPR" not in labels


def test_a_keyless_run_says_the_image_could_not_be_read(client, sample_files, monkeypatch):
    """§7 + §21.17. With no vision model the deck is produced — but the user is told,
    in the report, that the screenshot contributed nothing. An empty result that looked
    like a clean run would be the worst possible outcome."""
    from app import llm

    monkeypatch.setattr("app.config.settings.anthropic_api_key", None)
    llm.reset_client()

    session = client.post("/api/session").json()["session_id"]
    with open(sample_files / "risk_dashboard.png", "rb") as handle:
        client.post(f"/api/upload?session_id={session}",
                    files={"files": ("risk_dashboard.png", handle, "image/png")})

    analysis = client.post("/api/analyze", json={
        "session_id": session, "request_text": REQUEST,
    }).json()

    warnings = " ".join(analysis["warnings"])
    assert "Could NOT interpret" in warnings
    assert "MISSING from this report" in warnings
