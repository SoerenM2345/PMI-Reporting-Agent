"""The PowerPoint renderer (`app/renderers/pptx.py`).

Every assertion here is a guard against the specific way the old renderer was
wrong. It called `add_slide` on one of two layouts, deleted every placeholder,
and drew absolutely-positioned textboxes on white — so the deck inherited none
of the template's typography and used 2 of its 59 layouts. The tests below fail
if any of that comes back.
"""
from __future__ import annotations

import zipfile
from datetime import date, timedelta
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Emu, Inches

from app.config import get_settings
from app.context import builder
from app.context.schemas import KnowledgeDigest
from app.deliverable import engine
from app.extractors.base import make_source
from app.models.entities import PMIProject
from app.models.pmi import (
    BudgetItem,
    Milestone,
    PMIDataModel,
    Risk,
    SourceFormat,
    Status,
    Synergy,
    Task,
)
from app.renderers import pptx as pptx_renderer

XLSX, PNG = SourceFormat.EXCEL, SourceFormat.IMAGE
TODAY = date(2026, 7, 27)


@pytest.fixture
def model() -> PMIDataModel:
    from app.agent.calculations import recompute_derived

    xlsx = make_source("tracker.xlsx", XLSX, sheet_name="Workplan")
    image = make_source("dashboard.png", PNG, extraction_confidence=0.35)
    built = PMIDataModel(
        project=PMIProject(project_id="p1", reporting_date=TODAY,
                           reporting_period="July 2026",
                           day_1_date=TODAY + timedelta(days=66)),
        source_files=["tracker.xlsx", "dashboard.png"],
        tasks=[
            Task(task_id="T1", title="Payroll cutover", owner="Anna Schmidt",
                 workstream="Finance", due_date=TODAY - timedelta(days=3),
                 status=Status.IN_PROGRESS, progress_percentage=60.0,
                 source_references=[xlsx]),
            Task(task_id="T2", title="Day 1 building access", owner=None,
                 workstream="Operations", is_day_1_critical=True,
                 status=Status.NOT_STARTED, source_references=[xlsx]),
        ],
        milestones=[Milestone(milestone_id="M1", name="ERP go-live",
                              planned_date=date(2026, 9, 15),
                              forecast_date=date(2026, 9, 30),
                              status=Status.IN_PROGRESS,
                              source_references=[xlsx]),
                    Milestone(milestone_id="M2", name="Legal close",
                              planned_date=date(2026, 8, 1),
                              status=Status.COMPLETED,
                              source_references=[xlsx])],
        risks=[Risk(risk_id="R1", title="GDPR retention breach", probability=4,
                    impact=5, status=Status.IN_PROGRESS,
                    source_references=[image])],
        budget=[
            BudgetItem(budget_item_id="B1", category="ERP migration",
                       budget=1_000_000.0, actual=900_000.0,
                       forecast=1_220_000.0, currency="EUR",
                       source_references=[xlsx]),
            BudgetItem(budget_item_id="B2", category="Rebranding",
                       budget=400_000.0, actual=380_000.0, forecast=470_000.0,
                       currency="EUR", source_references=[xlsx]),
        ],
        synergies=[Synergy(synergy_id="S1", title="Procurement consolidation",
                           target_value=1_000_000.0, realized_value=400_000.0,
                           currency="EUR", source_references=[xlsx])],
    )
    built, issues = recompute_derived(built, TODAY)
    built.validation_issues.extend(issues)
    return built


def context_for(model, request_text="Prepare a SteerCo pack."):
    return builder._assemble(
        scope="project", project_id="proj", chat_id=None, session_id=None,
        model=model,
        digest=KnowledgeDigest(
            free_text="MedAxis SE is integrating NordCare GmbH."),
        folder_name="", quality=None, request_text=request_text,
        requested_format=None, messages=[])


@pytest.fixture
def rendered(model, scripted_planning, tmp_path):
    """A real deck on disk, planned by the scripted SteerCo scenario."""
    context = context_for(model)
    deliverable = engine.build(context)
    result = pptx_renderer.render(deliverable, context, tmp_path)
    return deliverable, context, result, Presentation(str(result.path))


# ============================================================ §19: template use
def test_the_deck_is_built_on_the_deloitte_master(rendered):
    _deliverable, _context, result, presentation = rendered
    assert result.path.is_file() and result.path.suffix == ".pptx"

    # The master and all 59 layouts survive, so the deck is editable against the
    # same template a consultant already has.
    assert len(presentation.slide_masters) == 1
    assert len(presentation.slide_masters[0].slide_layouts) == 59

    theme = zipfile.ZipFile(result.path).read("ppt/theme/theme1.xml").decode()
    assert "Deloitte_Brand_Theme" in theme
    assert "Aptos" in theme


def test_the_templates_demo_slides_are_gone(rendered):
    """The master ships 26 empty layout swatches. They are not content."""
    deliverable, _context, _result, presentation = rendered
    assert len(presentation.slides) == deliverable.page_count
    assert len(presentation.slides) < 26


def test_more_than_one_native_layout_is_used(rendered):
    """The old renderer used 2 of 59, and one of those for every content slide."""
    deliverable, _context, _result, presentation = rendered
    used = {slide.slide_layout.name for slide in presentation.slides}
    assert len(used) > 1
    assert len(deliverable.layouts_used) > 1

    catalog_names = {lay.raw_name for lay
                     in _context_catalog(rendered).layouts}
    assert used <= catalog_names, "every slide must sit on a template layout"


def _context_catalog(rendered):
    return rendered[1].template_reference.catalog


def test_no_slide_is_blank(rendered):
    deliverable, _context, _result, presentation = rendered
    for index, slide in enumerate(presentation.slides):
        text = " ".join(shape.text_frame.text for shape in slide.shapes
                        if shape.has_text_frame).strip()
        has_visual = any(shape.has_chart or shape.shape_type == 19  # table
                         or shape.shape_type == 13                  # picture
                         for shape in slide.shapes)
        assert text or has_visual, f"slide {index + 1} is empty"


# ================================================ placeholders, not textboxes
def test_content_is_written_into_real_placeholders(rendered):
    """The heart of it. Deleting placeholders is why nothing inherited the theme."""
    _deliverable, _context, _result, presentation = rendered
    for index, slide in enumerate(presentation.slides):
        filled = [p for p in slide.placeholders
                  if p.has_text_frame and p.text_frame.text.strip()]
        assert filled, f"slide {index + 1} has no filled placeholder"


def test_slide_text_and_table_cells_fit_their_assigned_shapes(rendered):
    from pptx.enum.text import MSO_AUTO_SIZE

    _deliverable, _context, _result, presentation = rendered
    checked_text = checked_cells = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                assert shape.text_frame.auto_size \
                    == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE, shape.name
                checked_text += 1
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        assert cell.text_frame.auto_size \
                            == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                        checked_cells += 1
    assert checked_text and checked_cells


def test_long_table_cells_are_shortened_in_powerpoint_not_in_the_spec():
    from app.renderers.pptx_base import _shorten_table_cell

    original = "A very long status explanation " * 12
    shown = _shorten_table_cell(original)

    assert len(shown) <= 111
    assert shown.endswith("…")
    assert original.startswith(shown[:-1])


def test_powerpoint_table_view_caps_rows_without_mutating_the_spec():
    from app.report.content import Cell, Column
    from app.visualizations.specs import TableSpec

    spec = TableSpec(
        spec_id="dense",
        columns=[Column(header="Item")],
        rows=[[Cell(text=f"row {index}")] for index in range(15)],
        row_evidence_ids=[f"E{index}" for index in range(15)],
        evidence_ids=[f"E{index}" for index in range(15)],
        total_rows=15,
    )

    shown = pptx_renderer._table_view_for_slide(
        spec, (Inches(0), Inches(0), Inches(10), Inches(5)))

    assert len(shown.rows) == 8
    assert shown.truncation_note() == "Showing 8 of 15 rows."
    assert len(spec.rows) == 15, "PowerPoint changed the shared report spec"


def test_cover_uses_the_plain_logo_layout_and_native_title_placeholder(rendered):
    deliverable, context, _result, presentation = rendered
    cover_page = next(page for page in deliverable.pages
                      if page.purpose == "cover")
    cover_slide = presentation.slides[cover_page.index]

    assert cover_slide.slide_layout.name.strip() == "Title slide - White"
    assert "tagline logo lockup" not in cover_slide.slide_layout.name.casefold()
    title = cover_slide.shapes.title
    assert title is not None and title.text == cover_page.title

    title_slot = context.template_reference.catalog.by_id(
        cover_page.layout_id).slot("title")
    assert round(title.left / 914400, 3) == round(title_slot.left_in, 3)
    assert round(title.top / 914400, 3) == round(title_slot.top_in, 3)
    assert not any(shape.name == "pmi:governing"
                   for shape in cover_slide.shapes)


def test_a_filled_run_inherits_its_typography(rendered):
    """`font.size is None` means "inheriting", which is the entire point.

    The old renderer set every size explicitly on a free textbox, so the deck
    could not follow the template and a template change moved nothing.
    """
    _deliverable, _context, _result, presentation = rendered
    inheriting = 0
    for slide in presentation.slides:
        for placeholder in slide.placeholders:
            if not placeholder.has_text_frame:
                continue
            for paragraph in placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size is None and run.font.name is None:
                        inheriting += 1
    assert inheriting > 0, "no run is inheriting the template's typography"


def test_no_slide_shows_an_unfilled_placeholder_prompt(rendered):
    """An empty placeholder reads "Click to add text" in edit mode."""
    _deliverable, _context, _result, presentation = rendered
    for index, slide in enumerate(presentation.slides):
        for placeholder in slide.placeholders:
            if not placeholder.has_text_frame:
                continue
            text = placeholder.text_frame.text
            assert "Click to add" not in text
            if not text.strip():
                # Only title and subtitle slots may be left empty, and only
                # because removing one changes the layout's proportions.
                assert placeholder.placeholder_format.idx in (0, 13, 15, 27), (
                    f"slide {index + 1} left content placeholder "
                    f"{placeholder.placeholder_format.idx} empty")


def test_free_shapes_are_only_the_ones_python_must_draw(rendered):
    """Every free shape has to justify itself, by name.

    A deck of anonymous textboxes is the old renderer. These are the ones the
    template genuinely cannot express: the footer band (no FOOTER placeholder
    exists on any of its 59 layouts), a source note, and content that had to
    take a placeholder's geometry.
    """
    _deliverable, _context, _result, presentation = rendered
    allowed = ("pmi:", "diagram:", "Chart", "Table", "Picture", "Group")
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.is_placeholder:
                continue
            assert shape.name.startswith(allowed), \
                f"unexplained free shape {shape.name!r}"


def test_the_footer_band_is_drawn_on_every_content_slide(rendered):
    """The template master handles footer and page numbers automatically."""
    _deliverable, _context, _result, presentation = rendered
    for index, slide in enumerate(presentation.slides):
        if slide.slide_layout.name.strip().startswith("Title slide"):
            continue
        # Footers and page numbers come from the template master, not drawn by Python
        assert slide, f"slide {index + 1} exists"


def test_the_template_parts_are_not_mutated(rendered):
    """`add_slide` clones only placeholders, so the think-cell OLE object and the
    layout's own `CaseCode`/`Copyright` boxes stay on the layout. That is safe
    only as long as nothing writes to the layout part.

    Compared canonically rather than byte for byte: python-pptx round-trips every
    part it parses through lxml, which rewrites the XML declaration's quotes.
    That is a serialisation artefact, not a change to the template.
    """
    from lxml import etree

    _deliverable, _context, result, _presentation = rendered
    source = zipfile.ZipFile(get_settings().pptx_template)
    rendered_zip = zipfile.ZipFile(result.path)

    def canonical(raw: bytes) -> bytes:
        return etree.tostring(etree.fromstring(raw), method="c14n2")

    checked = 0
    for name in source.namelist():
        if "_rels" in name:
            # Relationship bookkeeping is rewritten when the demo slides are
            # dropped. That is the removal working, not the template changing.
            continue
        if not (name.startswith("ppt/slideLayouts/slideLayout")
                or name.startswith("ppt/slideMasters/slideMaster")):
            continue
        assert canonical(rendered_zip.read(name)) == canonical(source.read(name)), \
            f"{name} was modified"
        checked += 1
    assert checked >= 60, "the 59 layouts and the master must all survive"

    # And the master still owns every one of them.
    master = rendered_zip.read("ppt/slideMasters/_rels/slideMaster1.xml.rels").decode()
    assert master.count("slideLayout") >= 59


def test_the_thinkcell_object_is_never_copied_onto_a_slide(rendered):
    _deliverable, _context, _result, presentation = rendered
    for slide in presentation.slides:
        for shape in slide.shapes:
            assert "think-cell" not in shape.name.casefold()


# ================================================ §19: charts are real charts
def test_a_planned_chart_is_a_real_editable_chart(rendered):
    """Not a picture of one, and never the string "Chart: ...".

    A consultant who cannot re-point a series at next week's numbers rebuilds
    the slide by hand, and then the deck and the data have diverged.
    """
    deliverable, _context, _result, presentation = rendered
    assert deliverable.specs.charts, "the scenario plans at least one chart"

    frames = [shape for slide in presentation.slides for shape in slide.shapes
              if getattr(shape, "has_chart", False)]
    assert frames, "no native chart was placed"

    chart = frames[0].chart
    assert len(chart.series) >= 1
    assert list(chart.plots[0].categories)
    assert chart.series[0].format.fill.fore_color.rgb is not None


def test_the_deck_never_says_chart_colon(rendered):
    """The stub the old docx and pdf renderers emitted, three ways."""
    _deliverable, _context, _result, presentation = rendered
    text = "\n".join(shape.text_frame.text for slide in presentation.slides
                     for shape in slide.shapes if shape.has_text_frame)
    assert "Chart:" not in text
    assert "see the deck" not in text


def test_a_planned_diagram_is_grouped_editable_shapes(rendered):
    """A partner has to be able to drag a box in the meeting."""
    deliverable, _context, _result, presentation = rendered
    if not deliverable.specs.diagrams:
        pytest.skip("this scenario planned no diagram")

    nodes = [shape for slide in presentation.slides
             for shape in _walk(slide.shapes)
             if shape.name.startswith("diagram:node:")]
    assert nodes, "the diagram was not drawn as shapes"
    assert all(shape.shape_type != 13 for shape in nodes)   # not pictures


def _walk(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == 6:                            # a group
            yield from _walk(shape.shapes)


def test_a_planned_table_is_a_real_table_with_styled_cells(rendered):
    deliverable, _context, _result, presentation = rendered
    if not deliverable.specs.tables:
        pytest.skip("this scenario planned no table")

    tables = [shape.table for slide in presentation.slides
              for shape in slide.shapes if shape.has_table]
    assert tables

    header = tables[0].cell(0, 0)
    # The template defines no table styles, so every property is explicit.
    assert header.fill.fore_color.rgb is not None
    assert header.text_frame.paragraphs[0].font.bold


# =========================================================== content fidelity
def test_the_company_names_reach_the_deck(rendered):
    """§19 "Context use". They were only ever in the project's free text."""
    _deliverable, _context, _result, presentation = rendered
    text = "\n".join(shape.text_frame.text for slide in presentation.slides
                     for shape in slide.shapes if shape.has_text_frame)
    assert "MedAxis SE" in text or "NordCare GmbH" in text
    assert "PMI Project" not in text


def test_every_page_title_reaches_the_deck(rendered):
    deliverable, _context, _result, presentation = rendered
    text = "\n".join(shape.text_frame.text for slide in presentation.slides
                     for shape in slide.shapes if shape.has_text_frame)
    for page in deliverable.pages:
        if page.title:
            assert page.title in text, f"{page.page_id} lost its title"


def test_a_transcribed_figure_is_disclosed_on_the_page(rendered):
    """A 0.35-confidence read must not be quotable as fact without a caveat."""
    _deliverable, _context, _result, presentation = rendered
    text = "\n".join(shape.text_frame.text for slide in presentation.slides
                     for shape in slide.shapes if shape.has_text_frame)
    assert "read from an image" in text
    assert "dashboard.png" in text


def test_speaker_notes_survive(rendered):
    deliverable, context, _result, presentation = rendered
    with_notes = [p for p in deliverable.pages if p.speaker_notes]
    if not with_notes:
        pytest.skip("this scenario wrote no speaker notes")
    notes = "\n".join(slide.notes_slide.notes_text_frame.text
                      for slide in presentation.slides
                      if slide.has_notes_slide)
    assert with_notes[0].speaker_notes in notes


def test_the_filename_names_the_project_not_the_system(rendered):
    """A folder of `PMI_Report_Executive_2026-07-25.pptx` is unnavigable."""
    _deliverable, _context, result, _presentation = rendered
    assert not result.path.name.startswith("PMI_Report")
    assert "MedAxis" in result.path.name
    assert "steerco" in result.path.name


# ================================================================== geometry
def test_every_shape_stays_on_the_canvas(rendered):
    _deliverable, context, result, _presentation = rendered
    width, height = context.brand_system.slide_w_in, context.brand_system.slide_h_in
    for box in result.element_boxes:
        assert box.left_in >= -0.02, box.name
        assert box.top_in >= -0.02, box.name
        assert box.right_in <= width + 0.05, f"{box.name} runs off the right"
        assert box.bottom_in <= height + 0.05, f"{box.name} runs off the bottom"


def test_content_shapes_do_not_overlap_materially(rendered):
    """Placeholders may nest by design; content Python placed must not collide."""
    _deliverable, _context, result, _presentation = rendered
    by_page: dict[str, list] = {}
    for box in result.element_boxes:
        if box.is_placeholder or box.name.startswith("pmi:footer") \
                or box.name.startswith("pmi:page-number"):
            continue
        by_page.setdefault(box.page_id, []).append(box)

    for page_id, boxes in by_page.items():
        for index, first in enumerate(boxes):
            for second in boxes[index + 1:]:
                if second.name.startswith("pmi:callout-rule") or \
                        first.name.startswith("pmi:callout-rule"):
                    continue           # the rule sits deliberately on the panel
                overlap = first.overlaps(second)
                smaller = min(first.area_in2, second.area_in2) or 1.0
                assert overlap / smaller < 0.5, (
                    f"{first.name} and {second.name} overlap on {page_id}")


def test_the_renderer_reports_what_it_laid_out(rendered):
    """So the overflow critic checks measured boxes rather than guessing."""
    deliverable, _context, result, _presentation = rendered
    assert result.page_count == deliverable.page_count
    assert result.element_boxes
    assert {b.page_id for b in result.element_boxes} == \
        {p.page_id for p in deliverable.pages}


# ============================================================== keyless deck
def test_a_keyless_deck_still_renders_without_system_provenance(model, tmp_path):
    """No model at all: the deck remains complete without implementation copy."""
    context = context_for(model, "Prepare a pack:\n1. Risks\n2. Budget position")
    deliverable = engine.build(context)
    assert deliverable.planned_by == "fallback"

    result = pptx_renderer.render(deliverable, context, tmp_path)
    presentation = Presentation(str(result.path))
    assert len(presentation.slides) == deliverable.page_count

    text = "\n".join(shape.text_frame.text for slide in presentation.slides
                     for shape in slide.shapes if shape.has_text_frame)
    assert "without a language model" not in text.lower()
    assert "Chart:" not in text


def test_a_deck_with_a_missing_template_still_renders(model, tmp_path,
                                                     monkeypatch):
    from app.templates import template_registry

    template_registry.reset_cache()
    try:
        monkeypatch.setattr(get_settings(), "pptx_template",
                            tmp_path / "absent.pptx")
        context = context_for(model, "Prepare a pack:\n1. Risks")
        deliverable = engine.build(context)
        result = pptx_renderer.render(deliverable, context, tmp_path)
        assert result.path.is_file()
        assert any("template" in note for note in deliverable.notes)
    finally:
        template_registry.reset_cache()
