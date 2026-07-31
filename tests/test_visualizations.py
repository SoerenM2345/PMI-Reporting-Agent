"""Charts, diagrams, tables and on-demand calculations.

The rule this layer exists to enforce: a chart either has validated data behind
every point, or it is not a chart and the caller shows a table instead. There is
no path to a caption over empty space — which the old renderers reached three
different ways.
"""
from __future__ import annotations

from datetime import date

import pytest

from app.evidence.model import EvidenceIndex, EvidenceItem
from app.extractors.base import make_source
from app.generation import calculator
from app.models.pmi import SourceFormat
from app.templates import template_registry
from app.visualizations import builder, charts, diagrams, scales, tables, validator
from app.visualizations.specs import (
    ChartRequest,
    ChartSeries,
    ChartSpec,
    DataPoint,
    DiagramNodeRequest,
    DiagramRequest,
    SeriesRequest,
)

XLSX = SourceFormat.EXCEL


@pytest.fixture(scope="module")
def brand():
    return template_registry.default().brand


@pytest.fixture
def evidence() -> EvidenceIndex:
    """Three budget lines and two synergies, with one figure not reported."""
    index = EvidenceIndex()
    source = make_source("tracker.xlsx", XLSX, sheet_name="Budget")

    for order, (name, budget, actual, forecast) in enumerate([
            ("ERP migration", 1_000_000.0, 900_000.0, 1_220_000.0),
            ("Rebranding", 400_000.0, 380_000.0, 470_000.0),
            ("Advisory", 250_000.0, 240_000.0, 250_000.0)], start=1):
        index.add(EvidenceItem(
            evidence_id=f"ev:budget:B{order}", kind="budget",
            origin="normalized_value", label=name,
            statement=f"{name} budget line.", value=budget - forecast,
            display="", currency="EUR", entity_type="budget",
            entity_id=f"B{order}", sources=[source],
            payload={"category": name, "budget": budget, "actual": actual,
                     "forecast": forecast, "variance": budget - forecast},
            search_text=name))

    index.add(EvidenceItem(
        evidence_id="ev:synergy:S1", kind="synergy", origin="normalized_value",
        label="Procurement", statement="Procurement synergy.", currency="EUR",
        entity_type="synergy", entity_id="S1", sources=[source],
        payload={"target_value": 1_000_000.0, "realized_value": 400_000.0},
        search_text="procurement synergy"))
    index.add(EvidenceItem(
        evidence_id="ev:synergy:S2", kind="synergy", origin="normalized_value",
        label="Headcount", statement="Headcount synergy.", currency="EUR",
        entity_type="synergy", entity_id="S2", sources=[source],
        # No realised figure: nobody reported it.
        payload={"target_value": 500_000.0, "realized_value": None},
        search_text="headcount synergy"))
    return index


def budget_request(**kwargs) -> ChartRequest:
    defaults = dict(
        spec_id="c1", chart_type="column", title="Forecast exceeds budget",
        insight="Forecast is above budget on two of three lines.",
        evidence_ids=["ev:budget:B1", "ev:budget:B2", "ev:budget:B3"],
        series=[SeriesRequest(name="Budget", value_field="budget"),
                SeriesRequest(name="Forecast", value_field="forecast")],
        caption="Budget against current forecast, by cost line.",
    )
    defaults.update(kwargs)
    return ChartRequest(**defaults)


# ======================================================== values come from evidence
def test_every_point_carries_the_evidence_it_came_from(evidence):
    spec = builder.build_chart(budget_request(), evidence)
    for series in spec.series:
        for point in series.points:
            assert point.evidence_id.startswith("ev:budget:")
            assert evidence.get(point.evidence_id) is not None


def test_values_are_read_from_the_named_field(evidence):
    spec = builder.build_chart(budget_request(), evidence)
    budget, forecast = spec.series
    assert budget.points[0].value == 1_000_000.0
    assert forecast.points[0].value == 1_220_000.0
    assert budget.points[0].display == "EUR 1,000,000"


def test_an_unknown_field_is_dropped_with_a_warning(evidence):
    spec = builder.build_chart(budget_request(series=[
        SeriesRequest(name="Budget", value_field="budget"),
        SeriesRequest(name="Nonsense", value_field="__import__"),
    ]), evidence)
    assert [s.name for s in spec.series] == ["Budget"]
    assert any("unknown field" in w for w in spec.warnings)


def test_a_missing_figure_stays_missing(evidence):
    """Never a zero. "Nobody reported this" and "this is nil" are different
    findings and the second is usually worse news."""
    spec = builder.build_chart(ChartRequest(
        spec_id="c2", chart_type="column", title="Synergy realisation",
        evidence_ids=["ev:synergy:S1", "ev:synergy:S2"],
        series=[SeriesRequest(name="Realised", value_field="realized_value")],
        caption="Realised synergy by initiative."), evidence)

    point = spec.series[0].points[1]
    assert point.value is None
    assert point.display == "Not Reported"
    assert point.note == "not reported"


def test_sorting_puts_missing_values_last_not_at_zero(evidence):
    spec = builder.build_chart(ChartRequest(
        spec_id="c3", chart_type="bar", title="Realisation",
        evidence_ids=["ev:synergy:S2", "ev:synergy:S1"],
        series=[SeriesRequest(name="Realised", value_field="realized_value")],
        sort="value_desc", caption="x"), evidence)
    assert [p.value for p in spec.series[0].points] == [400_000.0, None]


def test_the_chart_discloses_a_transcribed_figure(evidence):
    faint = make_source("photo.png", SourceFormat.IMAGE,
                        extraction_confidence=0.3)
    evidence.get("ev:budget:B1").sources = [faint]
    spec = builder.build_chart(budget_request(), evidence)
    assert "read from an image" in spec.caption
    assert "photo.png" in spec.source_note


def test_the_chart_discloses_a_disputed_figure(evidence):
    evidence.get("ev:budget:B1").conflict_ids.append("cf_001")
    spec = builder.build_chart(budget_request(), evidence)
    assert "disputed" in spec.caption
    assert spec.series[0].points[0].emphasis == "warn"


# ================================================================= validation
def test_a_valid_chart_passes(evidence):
    spec = builder.build_chart(budget_request(), evidence)
    assert validator.validate_chart(spec, evidence).ok


def test_a_point_that_disagrees_with_its_evidence_is_rejected(evidence):
    spec = builder.build_chart(budget_request(), evidence)
    spec.series[0].points[0].value = 42.0
    result = validator.validate_chart(spec, evidence)
    assert not result.ok
    assert "not in" in result.summary


def test_a_point_citing_nothing_is_rejected(evidence):
    spec = builder.build_chart(budget_request(), evidence)
    spec.series[0].points[0].evidence_id = ""
    assert not validator.validate_chart(spec, evidence).ok


def test_mixed_currencies_on_one_axis_are_rejected(evidence):
    spec = builder.build_chart(budget_request(), evidence)
    spec.series[0].currency, spec.series[1].currency = "EUR", "USD"
    result = validator.validate_chart(spec, evidence)
    assert not result.ok
    assert "currencies" in result.summary
    assert "not comparable" in result.summary


def test_mixed_units_on_one_axis_are_rejected(evidence):
    spec = builder.build_chart(budget_request(), evidence)
    spec.series[0].unit, spec.series[1].unit = "%", "EUR"
    assert "mixes units" in validator.validate_chart(spec, evidence).summary


def test_a_pie_containing_a_missing_value_is_rejected(evidence):
    """You cannot honestly stack an unknown: the reader takes the visible
    segments as the whole."""
    spec = builder.build_chart(ChartRequest(
        spec_id="c4", chart_type="pie", title="Realisation",
        evidence_ids=["ev:synergy:S1", "ev:synergy:S2"],
        series=[SeriesRequest(name="Realised", value_field="realized_value")],
        caption="x"), evidence)
    result = validator.validate_chart(spec, evidence)
    assert not result.ok
    assert "presents its parts as a whole" in result.summary


def test_a_column_chart_may_show_a_gap(evidence):
    spec = builder.build_chart(ChartRequest(
        spec_id="c5", chart_type="column", title="Realisation",
        evidence_ids=["ev:synergy:S1", "ev:synergy:S2"],
        series=[SeriesRequest(name="Realised", value_field="realized_value")],
        caption="x"), evidence)
    result = validator.validate_chart(spec, evidence)
    assert result.ok
    assert any("not reported" in w for w in result.warnings)


def test_a_pie_with_too_many_slices_is_rejected(evidence):
    spec = ChartSpec(spec_id="c6", chart_type="pie", caption="x", alt_text="x",
                     series=[ChartSeries(name="s", points=[
                         DataPoint(label=f"L{n}", value=float(n),
                                   evidence_id="ev:budget:B1")
                         for n in range(1, 10)])])
    assert "unreadable" in validator.validate_chart(spec, evidence).summary


def test_a_percentage_axis_outside_zero_to_one_hundred_is_rejected(evidence):
    spec = ChartSpec(spec_id="c7", chart_type="column", caption="x",
                     alt_text="x", series=[ChartSeries(
                         name="s", unit="%",
                         points=[DataPoint(label="A", value=140.0,
                                           evidence_id="ev:budget:B1",
                                           display="140%")])])
    spec.value_axis.is_percentage = True
    result = validator.validate_chart(spec, evidence)
    assert "percentage axis" in result.summary


def test_a_waterfall_that_does_not_reconcile_is_rejected(evidence):
    """A bridge that does not add up is worse than no bridge."""
    evidence.add(EvidenceItem(
        evidence_id="ev:calc:erp_increase", kind="calculation",
        origin="computed_value", label="ERP increase", value=220_000.0))
    evidence.add(EvidenceItem(
        evidence_id="ev:calc:overstated", kind="calculation",
        origin="computed_value", label="Overstated close", value=1_500_000.0))

    def bridge(end_evidence: str, end: float) -> ChartSpec:
        return ChartSpec(
            spec_id="c8", chart_type="waterfall", caption="x", alt_text="x",
            series=[ChartSeries(name="s", points=[
                DataPoint(label="Open", value=1_000_000.0,
                          evidence_id="ev:budget:B1"),
                DataPoint(label="ERP", value=220_000.0,
                          evidence_id="ev:calc:erp_increase"),
                DataPoint(label="Close", value=end,
                          evidence_id=end_evidence)])])

    assert "does not reconcile" in validator.validate_chart(
        bridge("ev:calc:overstated", 1_500_000.0), evidence).summary
    # And the honest one passes: 1,000,000 + 220,000 = 1,220,000.
    assert validator.validate_chart(
        bridge("ev:budget:B1", 1_220_000.0), evidence).ok


def test_a_chart_with_no_caption_is_rejected(evidence):
    spec = builder.build_chart(budget_request(caption=""), evidence)
    spec.caption = ""
    assert "no caption" in validator.validate_chart(spec, evidence).summary


# ====================================================== §19: real charts exist
def test_a_chart_renders_to_a_real_png(evidence, brand, tmp_path):
    spec = builder.build_chart(budget_request(), evidence)
    path = charts.to_png(spec, brand, tmp_path)
    assert path.is_file() and path.stat().st_size > 2000
    assert path.suffix == ".png"


def test_a_chart_renders_to_real_svg_with_traceable_marks(evidence, brand):
    spec = builder.build_chart(budget_request(), evidence)
    svg = charts.to_svg(spec, brand)

    assert svg.startswith("<svg") and svg.rstrip().endswith("</svg>")
    assert svg.count('class="pmi-mark"') == 6      # 3 categories x 2 series
    assert 'data-evidence-id="ev:budget:B1"' in svg
    assert "<title>" in svg and "<desc>" in svg
    # Self-contained: nothing to fetch.
    assert "http" not in svg and "<image" not in svg


def test_a_missing_value_is_labelled_in_the_svg_not_drawn_as_zero(evidence, brand):
    spec = builder.build_chart(ChartRequest(
        spec_id="c9", chart_type="column", title="Realisation",
        evidence_ids=["ev:synergy:S1", "ev:synergy:S2"],
        series=[SeriesRequest(name="Realised", value_field="realized_value")],
        caption="x"), evidence)
    svg = charts.to_svg(spec, brand)
    assert "not reported" in svg
    assert svg.count('class="pmi-mark"') == 1      # one bar, not two


def test_a_native_pptx_chart_is_a_real_editable_chart(evidence, brand, tmp_path):
    """A PNG of a chart is not a chart: a consultant who cannot re-point a
    series will rebuild the slide by hand."""
    from pptx import Presentation
    from pptx.util import Inches

    spec = builder.build_chart(budget_request(), evidence)
    assert spec.is_native_pptx

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    frame = charts.to_pptx_native(spec, brand, slide,
                                  (Inches(1), Inches(1), Inches(8), Inches(4.5)))

    assert frame.has_chart
    chart = frame.chart
    assert len(chart.series) == 2
    assert list(chart.plots[0].categories) == ["ERP migration", "Rebranding",
                                               "Advisory"]
    # Styled explicitly, because the template defines no chartStyles.xml.
    assert chart.series[0].format.fill.fore_color.rgb is not None


def test_a_non_native_type_is_flagged_for_image_rendering(evidence):
    spec = builder.build_chart(budget_request(chart_type="heatmap"), evidence)
    assert spec.render_as_image and not spec.is_native_pptx


# ===================================================================== scales
def test_axes_are_rounded_to_readable_bounds():
    """An axis topping out at 1,237,000 tells a reader the maximum is precise."""
    scale = scales.nice_scale([0.0, 1_237_000.0])
    assert scale.maximum >= 1_237_000.0
    assert scale.ticks[0] == 0.0 and scale.ticks[-1] == scale.maximum
    # Evenly spaced on a 1-2-2.5-5 step, so every gridline is a round figure.
    steps = {round(b - a, 6) for a, b in zip(scale.ticks, scale.ticks[1:])}
    assert len(steps) == 1
    step = steps.pop()
    assert str(step / 10 ** len(str(int(step)).rstrip("0"))).rstrip("0.") in \
        ("", "1", "2", "25", "5")


def test_a_missing_value_does_not_drag_the_axis_to_zero():
    """A phantom zero makes every other bar look bigger."""
    with_gap = scales.nice_scale([500.0, None, 700.0], include_zero=False)
    assert with_gap.minimum > 0


def test_a_percentage_axis_runs_zero_to_one_hundred():
    scale = scales.nice_scale([12.0, 48.0], is_percentage=True)
    assert (scale.minimum, scale.maximum) == (0.0, 100.0)


def test_a_stacked_column_with_a_gap_has_no_total():
    """Summing the parts that happen to be known understates the column while
    looking authoritative."""
    totals = scales.stacked_totals([[10.0, 20.0], [5.0, None]])
    assert totals == [15.0, None]


def test_long_axis_labels_are_abbreviated():
    assert scales.tick_label(1_200_000, currency="EUR") == "EUR 1.2m"
    assert scales.tick_label(45, is_percentage=True) == "45%"


# =================================================================== diagrams
def test_a_timeline_places_markers_by_real_dates(evidence, brand):
    """Evenly spacing dated events is the commonest way a timeline lies."""
    for order, when in enumerate(["2026-08-01", "2026-08-08", "2027-02-01"],
                                 start=1):
        evidence.add(EvidenceItem(
            evidence_id=f"ev:milestone:M{order}", kind="milestone",
            origin="normalized_value", label=f"Gate {order}",
            due=date.fromisoformat(when), search_text=f"gate {order}"))

    spec = builder.build_diagram(DiagramRequest(
        spec_id="d1", diagram_type="timeline", title="Plan",
        evidence_ids=["ev:milestone:M1", "ev:milestone:M2", "ev:milestone:M3"],
        caption="Milestone plan."), evidence)
    assert validator.validate_diagram(spec, evidence).ok

    box = diagrams.Box(0.0, 0.0, 12.0, 3.0)
    placed = diagrams.layout_nodes(spec, box)
    lefts = [placed[n.node_id].left for n in spec.nodes]
    # The first two are a week apart, the third six months later.
    assert lefts[1] - lefts[0] < (lefts[2] - lefts[1]) / 4


def test_a_timeline_with_one_date_is_rejected(evidence):
    evidence.add(EvidenceItem(evidence_id="ev:milestone:M9", kind="milestone",
                              origin="normalized_value", label="Only one",
                              due=date(2026, 8, 1)))
    spec = builder.build_diagram(DiagramRequest(
        spec_id="d2", diagram_type="timeline", caption="x",
        evidence_ids=["ev:milestone:M9"]), evidence)
    assert "needs at least two" in validator.validate_diagram(
        spec, evidence).summary


def test_a_risk_matrix_reads_its_coordinates_from_evidence(evidence):
    for order, (probability, impact) in enumerate([(4, 5), (2, 2)], start=1):
        evidence.add(EvidenceItem(
            evidence_id=f"ev:risk:R{order}", kind="risk",
            origin="normalized_value", label=f"Risk {order}",
            payload={"probability": probability, "impact": impact},
            search_text=f"risk {order}"))

    spec = builder.build_diagram(DiagramRequest(
        spec_id="d3", diagram_type="risk_matrix", caption="Risk grid.",
        evidence_ids=["ev:risk:R1", "ev:risk:R2"],
        x_axis_label="Probability", y_axis_label="Impact"), evidence)

    assert validator.validate_diagram(spec, evidence).ok
    assert (spec.nodes[0].column, spec.nodes[0].row) == (4, 5)

    placed = diagrams.layout_nodes(spec, diagrams.Box(0, 0, 10, 5))
    high, low = placed["n1"], placed["n2"]
    assert high.left > low.left            # higher probability, further right
    assert high.top < low.top              # higher impact, further up


def test_a_value_driver_tree_needs_one_root(evidence):
    nodes = [DiagramNodeRequest(node_id="a", label="Total"),
             DiagramNodeRequest(node_id="b", label="Part", parent_id="a")]
    spec = builder.build_diagram(DiagramRequest(
        spec_id="d4", diagram_type="value_driver_tree", caption="x",
        nodes=nodes), evidence)
    assert validator.validate_diagram(spec, evidence).ok

    spec.nodes[1].parent_id = ""
    assert "exactly one root" in validator.validate_diagram(
        spec, evidence).summary


def test_a_tree_centres_parents_over_their_children(evidence):
    spec = builder.build_diagram(DiagramRequest(
        spec_id="d5", diagram_type="value_driver_tree", caption="x", nodes=[
            DiagramNodeRequest(node_id="root", label="Total"),
            DiagramNodeRequest(node_id="l", label="Left", parent_id="root"),
            DiagramNodeRequest(node_id="r", label="Right", parent_id="root"),
        ]), evidence)
    placed = diagrams.layout_nodes(spec, diagrams.Box(0, 0, 12, 4))
    root_centre = placed["root"].left + placed["root"].width / 2
    left_centre = placed["l"].left + placed["l"].width / 2
    right_centre = placed["r"].left + placed["r"].width / 2
    assert left_centre < root_centre < right_centre
    assert root_centre == pytest.approx((left_centre + right_centre) / 2, abs=0.05)
    assert placed["root"].top < placed["l"].top


def test_an_edge_naming_a_missing_node_is_rejected(evidence):
    from app.visualizations.specs import DiagramEdgeRequest

    spec = builder.build_diagram(DiagramRequest(
        spec_id="d6", diagram_type="process_flow", caption="x",
        nodes=[DiagramNodeRequest(node_id="a", label="A")],
        edges=[DiagramEdgeRequest(from_id="a", to_id="ghost")]), evidence)
    assert "not a node" in validator.validate_diagram(spec, evidence).summary


def test_a_diagram_renders_as_editable_shapes_not_a_picture(evidence, brand):
    """A partner has to be able to drag a box in the meeting."""
    from pptx import Presentation

    spec = builder.build_diagram(DiagramRequest(
        spec_id="d7", diagram_type="process_flow", title="Cutover",
        caption="Cutover sequence.", nodes=[
            DiagramNodeRequest(node_id=f"n{n}", label=f"Step {n}")
            for n in range(1, 5)]), evidence)

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    group = diagrams.to_pptx(spec, brand, slide, diagrams.Box(0.5, 2.0, 9.0, 2.0))

    assert group is not None
    assert not any(shape.shape_type == 13 for shape in slide.shapes)  # no picture
    names = [s.name for s in group.shapes]
    assert len(names) == 4
    assert all(name.startswith(diagrams.NODE_PREFIX) for name in names)


def test_a_diagram_renders_to_svg_and_png(evidence, brand, tmp_path):
    spec = builder.build_diagram(DiagramRequest(
        spec_id="d8", diagram_type="process_flow", caption="x", nodes=[
            DiagramNodeRequest(node_id="a", label="Plan"),
            DiagramNodeRequest(node_id="b", label="Build")]), evidence)

    svg = diagrams.to_svg(spec, brand)
    assert svg.startswith("<svg") and 'class="pmi-node"' in svg
    assert diagrams.to_png(spec, brand, tmp_path).is_file()


def test_matrix_collisions_are_deterministic(evidence):
    for order in range(1, 4):
        evidence.add(EvidenceItem(
            evidence_id=f"ev:risk:X{order}", kind="risk",
            origin="normalized_value", label=f"Same cell {order}",
            payload={"probability": 3, "impact": 3}))
    spec = builder.build_diagram(DiagramRequest(
        spec_id="d9", diagram_type="risk_matrix", caption="x",
        evidence_ids=[f"ev:risk:X{n}" for n in range(1, 4)]), evidence)

    first = diagrams.layout_nodes(spec, diagrams.Box(0, 0, 10, 5))
    second = diagrams.layout_nodes(spec, diagrams.Box(0, 0, 10, 5))
    assert first == second, "a snapshot test must not flap"
    assert len({(b.left, b.top) for b in first.values()}) == 3


# ===================================================================== tables
def test_a_table_gets_columns_from_the_evidence_kind(evidence):
    spec = tables.build_for("t1", evidence,
                            ["ev:budget:B1", "ev:budget:B2", "ev:budget:B3"])
    assert [c.header for c in spec.columns] == ["Cost line", "Budget", "Actual",
                                                "Forecast", "Variance"]
    assert spec.rows[0][0].text == "ERP migration"
    assert spec.rows[0][1].text == "EUR 1,000,000"
    assert validator.validate_table(spec, evidence).ok


def test_a_negative_variance_is_marked_bad(evidence):
    spec = tables.build_for("t2", evidence, ["ev:budget:B1"])
    variance = spec.rows[0][4]
    assert variance.value == -220_000.0
    assert variance.emphasis == "bad"


def test_a_missing_field_that_matters_is_flagged_in_the_cell(evidence):
    """An unowned critical risk should be visible while scanning the column."""
    evidence.add(EvidenceItem(
        evidence_id="ev:risk:R7", kind="risk", origin="normalized_value",
        label="GDPR breach", entity_type="risk", entity_id="R7",
        severity="critical", status="in_progress",
        payload={"risk_score": 20, "mitigation_action": None},
        sources=[make_source("tracker.xlsx", XLSX)]))
    spec = tables.build_for("t3", evidence, ["ev:risk:R7"])
    text = " ".join(cell.text for cell in spec.rows[0])
    assert "⚠ NO MITIGATION" in text
    assert "⚠ NO OWNER" in text


def test_a_table_cell_stays_editable(evidence):
    """The write-back handle that makes inline correction possible."""
    spec = tables.build_for("t4", evidence, ["ev:budget:B1"])
    reference = spec.rows[0][1].ref
    assert reference is not None
    assert (reference.entity_type, reference.entity_id, reference.field) == \
        ("budget", "B1", "budget")


def test_a_computed_row_offers_no_edit_handle(evidence):
    """Offering to edit a derived figure would be a lie."""
    evidence.add(EvidenceItem(
        evidence_id="ev:fact:budget.variance", kind="fact",
        origin="computed_value", label="Budget variance", value=-220_000.0,
        display="EUR -220,000"))
    spec = tables.build_for("t5", evidence, ["ev:fact:budget.variance"])
    assert all(cell.ref is None for cell in spec.rows[0])


def test_a_truncated_table_says_how_much_it_is_showing(evidence):
    for order in range(10, 30):
        evidence.add(EvidenceItem(
            evidence_id=f"ev:budget:B{order}", kind="budget",
            origin="normalized_value", label=f"Line {order}",
            payload={"category": f"Line {order}", "budget": 1000.0}))
    spec = tables.build_table("t6", list(evidence.of_kind("budget")),
                              row_limit=12)
    assert spec.is_truncated
    assert "Showing 12 of" in spec.truncation_note()
    assert validator.validate_table(spec, evidence).ok


def test_a_table_showing_a_subset_silently_is_rejected(evidence):
    spec = tables.build_for("t7", evidence, ["ev:budget:B1"])
    spec.total_rows = 50
    spec.row_limit = None
    object.__setattr__(spec, "total_rows", 50)
    result = validator.validate_table(spec, evidence)
    assert result.warnings, "a truncated table must disclose it"


# ============================================================== calculations
def test_a_requested_sum_is_computed_and_becomes_evidence(evidence):
    results, warnings = calculator.execute([calculator.CalculationRequest(
        calculation_id="total", op="sum", value_field="budget",
        input_evidence_ids=["ev:budget:B1", "ev:budget:B2", "ev:budget:B3"],
        label="Total approved budget", reason="The reader asked for the total.")],
        evidence)

    result = results[0]
    assert not result.refused
    assert result.value == 1_650_000.0
    assert result.display == "EUR 1,650,000"
    assert result.derivation.formula.startswith("sum of budget")
    assert warnings == []

    added = calculator.as_evidence(results, evidence)
    assert added[0].evidence_id == "ev:calc:total"
    assert "1650000" in evidence.numeric_corpus()


def test_a_computed_figure_cites_what_it_was_computed_from(evidence):
    results, _ = calculator.execute([calculator.CalculationRequest(
        calculation_id="total", op="sum", value_field="budget",
        input_evidence_ids=["ev:budget:B1"], label="Total")], evidence)
    added = calculator.as_evidence(results, evidence)
    assert added[0].source_files == ["tracker.xlsx"]


def test_mixed_currencies_are_refused_with_a_readable_reason(evidence):
    evidence.get("ev:budget:B2").currency = "USD"
    results, warnings = calculator.execute([calculator.CalculationRequest(
        calculation_id="mix", op="sum", value_field="budget",
        input_evidence_ids=["ev:budget:B1", "ev:budget:B2"],
        label="Total")], evidence)

    assert results[0].refused
    assert "different currencies" in results[0].refusal_reason
    assert "Convert them" in results[0].refusal_reason
    assert warnings, "the refusal must reach the page, not just a log"


def test_a_total_over_partly_reported_records_is_refused(evidence):
    """A sum over a partly-reported set is not a total, and presenting it as one
    is the commonest way an integration report understates a number."""
    results, _ = calculator.execute([calculator.CalculationRequest(
        calculation_id="realised", op="sum", value_field="realized_value",
        input_evidence_ids=["ev:synergy:S1", "ev:synergy:S2"],
        label="Total realised")], evidence)
    assert results[0].refused
    assert "incomplete while looking complete" in results[0].refusal_reason


def test_a_zero_denominator_is_refused(evidence):
    evidence.add(EvidenceItem(evidence_id="ev:zero", kind="fact",
                              origin="computed_value", label="Zero", value=0.0))
    results, _ = calculator.execute([calculator.CalculationRequest(
        calculation_id="r", op="ratio",
        input_evidence_ids=["ev:budget:B1", "ev:zero"], label="Ratio")],
        evidence)
    assert results[0].refused


def test_an_unsupported_operation_is_refused_not_executed(evidence):
    """`getattr` over a model-supplied operation name is arbitrary code."""
    results, _ = calculator.execute([calculator.CalculationRequest(
        calculation_id="x", op="__import__",
        input_evidence_ids=["ev:budget:B1"], label="Nope")], evidence)
    assert results[0].refused
    assert "not a supported operation" in results[0].refusal_reason


def test_a_percentage_result_is_labelled_as_one(evidence):
    results, _ = calculator.execute([calculator.CalculationRequest(
        calculation_id="pct", op="percent_of", value_field="realized_value",
        input_evidence_ids=["ev:synergy:S1", "ev:synergy:S1"],
        label="Realisation")], evidence)
    assert results[0].unit == "%"
    assert results[0].display == "100%"


def test_too_many_calculations_are_capped(evidence):
    requests = [calculator.CalculationRequest(
        calculation_id=f"c{n}", op="sum", value_field="budget",
        input_evidence_ids=["ev:budget:B1"], label=f"c{n}") for n in range(40)]
    results, warnings = calculator.execute(requests, evidence)
    assert len(results) == calculator.MAX_REQUESTS
    assert any("only the first" in w for w in warnings)


def test_a_refused_calculation_never_becomes_evidence(evidence):
    results, _ = calculator.execute([calculator.CalculationRequest(
        calculation_id="bad", op="ratio", input_evidence_ids=["ev:budget:B1"],
        label="Bad")], evidence)
    assert calculator.as_evidence(results, evidence) == []
    assert evidence.get("ev:calc:bad") is None
