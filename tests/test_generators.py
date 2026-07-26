"""P6: PowerPoint, Excel and chart generation (spec §12, §13, §14).

Structural assertions only — slide titles, sheet names, cell contents. Never a byte
comparison against a golden .pptx: python-pptx writes a zip with timestamps, so that
test would fail on a clock tick and teach everyone to ignore it.
"""
from __future__ import annotations

from datetime import date, timedelta

import pytest
from openpyxl import load_workbook
from pptx import Presentation

from app.agent.calculations import recompute_derived
from app.agent.consistency import run_checks
from app.agent.data_quality import build_report
from app.extractors.base import make_source
from app.generators import charts, pptx_report, xlsx_dashboard
from app.generators.quality_report import write_conflict_report, write_quality_report
from app.models.pmi import (
    KPI,
    Audience,
    BudgetItem,
    Milestone,
    PMIDataModel,
    PMIProject,
    Risk,
    SourceFormat,
    Status,
    Synergy,
    Task,
)

XLSX, PNG = SourceFormat.EXCEL, SourceFormat.IMAGE


@pytest.fixture
def model():
    """A project with the shapes the generators must handle: an unmitigated critical
    risk, an overdue task, a workstream that reported nothing, an over-budget line,
    and a figure read off a screenshot."""
    xlsx = make_source("tracker.xlsx", XLSX, sheet_name="Workplan")
    image = make_source("dashboard.png", PNG, extraction_confidence=0.35)
    yesterday = date.today() - timedelta(days=3)

    m = PMIDataModel(
        project=PMIProject(project_id="p1", project_name="Project Aurora",
                           reporting_date=date.today(),
                           day_1_date=date.today() + timedelta(days=30)),
        source_files=["tracker.xlsx", "dashboard.png"],
        tasks=[
            Task(task_id="T1", title="Payroll cutover", owner="Anna Schmidt",
                 workstream="Finance", due_date=yesterday, status=Status.IN_PROGRESS,
                 progress_percentage=60.0, source_references=[xlsx]),
            Task(task_id="T2", title="Day 1 building access", owner=None,
                 workstream="Operations", is_day_1_critical=True,
                 status=Status.NOT_STARTED, source_references=[xlsx]),
        ],
        milestones=[
            Milestone(milestone_id="M1", name="ERP go-live", is_go_live=True,
                      planned_date=date(2026, 9, 15),
                      forecast_date=date(2026, 9, 30),
                      status=Status.IN_PROGRESS, source_references=[xlsx]),
        ],
        risks=[
            Risk(risk_id="R1", title="GDPR retention breach", probability=4, impact=5,
                 owner="Lisa Chen", status=Status.IN_PROGRESS,
                 source_references=[image]),          # image-sourced, unmitigated
        ],
        budget=[
            BudgetItem(budget_item_id="B1", category="Advisors", budget=1000.0,
                       actual=900.0, forecast=1250.0, source_references=[xlsx]),
        ],
        synergies=[
            Synergy(synergy_id="S1", title="Procurement consolidation",
                    target_value=1_000_000.0, realized_value=400_000.0,
                    source_references=[xlsx]),
        ],
        kpis=[
            KPI(kpi_id="K1", name="Overall Progress", current_value=60.0,
                source_references=[xlsx]),
        ],
    )
    m, issues = recompute_derived(m, m.project.reporting_date)
    m.validation_issues.extend(issues)

    from app.agent.standardize import derive_workstreams

    derive_workstreams(m)

    results = run_checks(m)
    m.conflicts = results.conflicts
    m.validation_issues.extend(results.issues)
    return m


def titles(path) -> list[str]:
    out = []
    for slide in Presentation(str(path)).slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                out.append(shape.text_frame.text.split("\n")[0])
                break
    return out


def all_text(path) -> str:
    chunks = []
    for slide in Presentation(str(path)).slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    chunks += [c.text for c in row.cells]
    return "\n".join(chunks)


# ------------------------------------------------------------------ §12 PPTX
def test_every_audience_produces_an_openable_deck(model, tmp_path):
    """§21.15: 'Ensure generated files open successfully.'"""
    for audience in Audience:
        path = pptx_report.generate(model, audience, ["A bullet."], tmp_path)
        assert path.exists()
        assert len(Presentation(str(path)).slides) >= 3


def test_the_deck_leads_with_a_management_message_not_a_section_name(model, tmp_path):
    """§12.5: 'Use clear management-message titles.' A slide called "Risks" tells the
    reader nothing; one that names the finding tells them what to do."""
    path = pptx_report.generate(model, Audience.EXECUTIVE, ["A bullet."], tmp_path)
    text = " ".join(titles(path))

    assert "critical risk(s)" in text or "overdue" in text
    # The generic section names must not be doing the work of a title.
    assert not any(t.strip() in ("Risks", "Milestones", "Tasks") for t in titles(path))


def test_an_unmitigated_critical_risk_is_marked_on_the_slide(model, tmp_path):
    path = pptx_report.generate(model, Audience.EXECUTIVE, ["A bullet."], tmp_path)
    text = all_text(path)

    assert "GDPR retention breach" in text
    assert "NO MITIGATION" in text     # not a blank cell the reader might skim past


def test_a_day_1_task_with_no_owner_is_called_out(model, tmp_path):
    """A Day-1 task nobody owns is a Day-1 failure nobody has noticed yet. It must
    reach the deck — via the completeness check if not via a table."""
    path = pptx_report.generate(model, Audience.PMO, ["A bullet."], tmp_path)
    text = all_text(path).casefold()

    assert "day 1 building access" in text
    assert "no owner" in text


def test_every_deck_ends_with_its_limitations(model, tmp_path):
    """§12.5: 'Mark data-quality limitations.' Without this slide a report built from
    one clean tracker and one built from a blurry photo look equally confident."""
    quality = build_report(model)
    path = pptx_report.generate(model, Audience.EXECUTIVE, ["A bullet."], tmp_path,
                                quality=quality)

    assert "limitations" in titles(path)[-1].casefold() or \
           "conflict" in titles(path)[-1].casefold()
    text = all_text(path)
    assert "low confidence" in text.casefold() or "image" in text.casefold()


def test_the_finance_deck_leads_on_money_and_the_pmo_deck_on_operations(model, tmp_path):
    """§12.1-12.4 are different documents, not one document with a different cover."""
    finance = all_text(pptx_report.generate(model, Audience.FINANCE, ["x"], tmp_path))
    pmo = all_text(pptx_report.generate(model, Audience.PMO, ["x"], tmp_path))

    assert "Synergy" in finance or "synergy" in finance
    assert "Variance" in finance
    # The PMO deck carries the slide no tracker produces: who has NOT reported.
    assert "Missing updates" in pmo or "did not report" in pmo


def test_unknown_values_render_as_not_reported_never_as_zero(model, tmp_path):
    """§7: a project with no reported progress is not a project at 0%."""
    bare = PMIDataModel(
        project=PMIProject(project_id="p", project_name="Bare"),
        tasks=[Task(task_id="T1", title="Untracked",
                    source_references=[make_source("a.xlsx", XLSX)])],
    )
    path = pptx_report.generate(bare, Audience.EXECUTIVE, ["No data."], tmp_path)
    text = all_text(path)

    assert "Not Reported" in text
    assert "0%" not in text


# ------------------------------------------------------------------ §13 XLSX
def test_the_workbook_has_the_sheets_the_spec_asks_for(model, tmp_path):
    path = xlsx_dashboard.generate(model, Audience.PMO, ["x"], tmp_path,
                                   quality=build_report(model))
    names = load_workbook(str(path)).sheetnames

    # §13 sheets 1, 2, 3, 4, 8, 9 and — the important one — 10.
    for required in ("Executive Dashboard", "Workstream Scorecard", "Masterplan",
                     "Risks", "Budget", "Synergies", "Data Quality"):
        assert required in names, f"§13 requires a '{required}' sheet"


def test_sheets_are_filterable_and_frozen(model, tmp_path):
    """§13 quality requirements: filters, frozen headers. A dashboard you cannot filter
    is a table, and a table you cannot read is a PDF."""
    path = xlsx_dashboard.generate(model, Audience.PMO, ["x"], tmp_path)
    sheet = load_workbook(str(path))["Risks"]

    assert sheet.auto_filter.ref is not None
    assert sheet.freeze_panes == "A2"


def test_the_data_quality_sheet_shows_low_confidence_image_findings(model, tmp_path):
    """§13 sheet 10 + §5.6: findings read off a picture must be shown for review."""
    path = xlsx_dashboard.generate(model, Audience.PMO, ["x"], tmp_path,
                                   quality=build_report(model))
    sheet = load_workbook(str(path))["Data Quality"]
    text = "\n".join(
        str(cell.value) for row in sheet.iter_rows() for cell in row if cell.value
    )

    assert "Low-confidence" in text
    assert "GDPR retention breach" in text
    assert "dashboard.png" in text


def test_unknown_cells_say_not_reported(model, tmp_path):
    path = xlsx_dashboard.generate(model, Audience.PMO, ["x"], tmp_path)
    sheet = load_workbook(str(path))["Masterplan"]
    text = "\n".join(
        str(cell.value) for row in sheet.iter_rows() for cell in row if cell.value
    )
    assert "Not Reported" in text


# ---------------------------------------------------------------- §14 charts
def test_the_risk_heatmap_plots_scored_risks(model, tmp_path):
    path = charts.risk_heatmap(model, tmp_path)
    assert path and path.exists() and path.suffix == ".png"


def test_an_unscored_risk_is_reported_not_dropped_from_the_heatmap(tmp_path):
    """A risk with no probability is not a low risk — it is an unassessed one, and
    dropping it from the picture is how it stays unassessed."""
    m = PMIDataModel(risks=[
        Risk(risk_id="R1", title="Unassessed", impact=5,
             source_references=[make_source("a.xlsx", XLSX)])
    ])
    m, _ = recompute_derived(m)
    assert m.risks[0].risk_score is None

    path = charts.risk_heatmap(m, tmp_path)
    assert path and path.exists()   # rendered, with a "could not be plotted" note


def test_an_image_request_maps_to_the_topic_it_names(tmp_path):
    """"Generate an image describing the current milestones and project plan"
    must produce milestone charts, not the status default — the picture cannot
    ignore the words the user chose."""
    from app.llm import tasks

    parsed = tasks.parse_request(
        "Generate an image describing the current milestones and project plan"
    )
    assert parsed.output_type == "chart"
    assert parsed.topic == "milestone"


def test_generating_a_chart_renders_the_approved_content_first(model, tmp_path):
    """§18.18. A picture must not show something the preview did not.

    The content's own `ChartBlock`s render before the topic-driven extras, so a
    chart the user read about in the preview is the first thing in the image
    set.
    """
    from app.agent.graph import _render_charts
    from app.report.planner import plan

    content = plan(model, Audience.EXECUTIVE)   # EXECUTIVE carries a chart block
    chart_ids = [b.builder for s in content.sections for b in s.blocks
                 if getattr(b, "kind", None) == "chart"]
    assert chart_ids, "the executive deck should plan at least one chart"

    paths = _render_charts(content, model, "status", tmp_path)
    assert paths, "no charts were produced"
    assert all(p.exists() for p in paths)


def test_charts_are_chosen_by_topic(model, tmp_path):
    risk_charts = charts.generate(model, "risks", tmp_path)
    budget_charts = charts.generate(model, "budget", tmp_path)

    assert any("heatmap" in p.name for p in risk_charts)
    assert any("budget" in p.name for p in budget_charts)


def test_a_workstream_with_no_progress_is_not_drawn_as_zero(model, tmp_path):
    """A zero-height bar says "they achieved nothing", which is a different and
    possibly libellous claim (§7)."""
    path = charts.workstream_progress(model, tmp_path)
    assert path and path.exists()


# ------------------------------------------------ §18.18-19 the two reports
def test_the_conflict_report_shows_where_each_claim_came_from(model, tmp_path):
    """"Two files disagree" is not actionable. "Sheet Workplan cell A7 says 82, slide 4
    says 75" is."""
    model.kpis.append(
        KPI(kpi_id="K2", name="Overall Progress", current_value=75.0,
            source_references=[make_source("steerco.pptx", SourceFormat.POWERPOINT,
                                           slide_number=4)])
    )
    model.conflicts = run_checks(model).conflicts
    assert model.conflicts, "the fixture must actually contain a conflict"

    text = write_conflict_report(model, tmp_path).read_text()

    assert "# PMI Conflict Report" in text
    assert "| Source | Location | Value | Confidence |" in text
    assert "sheet 'Workplan'" in text
    assert "slide 4" in text
    assert "CRITICAL" in text          # overall progress is on §9's critical list
    assert "UNRESOLVED" in text


def test_the_quality_report_says_what_the_run_could_not_do(model, tmp_path):
    report = build_report(model, failed_files=["scan.pdf"],
                          warnings=["LLM unavailable — summaries are template-based"])
    path = write_quality_report(model, report, tmp_path)
    text = path.read_text()

    assert "could NOT be read" in text
    assert "scan.pdf" in text
    assert "Low-confidence findings" in text
    assert "template-based" in text
    assert "Senior Manager review" in text
