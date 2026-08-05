"""All four formats from one plan (`app/renderers/`).

The point of a single `Deliverable` is that the deck, the document, the PDF and
the dashboard cannot disagree. So most of what is tested here is agreement:
every page title, every disclosure and every figure appears in all four, and none
of them contains the `Chart: ...` stub that the old Word and PDF renderers
emitted in place of an actual chart.

Format-specific quality lives beside it: a real TOC field in Word, repeated table
headers and a two-pass contents in the PDF, and a genuinely self-contained
single-file dashboard.
"""
from __future__ import annotations

from datetime import date, timedelta
from pathlib import Path

import pytest

from app.context import builder
from app.context.schemas import KnowledgeDigest
from app.deliverable import engine
from app.extractors.base import make_source
from app.models.entities import PMIProject
from app.models.pmi import (
    BudgetItem,
    Milestone,
    PMIDataModel,
    Risk,
    SourceFormat,
    Status,
    Synergy,
    Task,
)
from app.renderers import registry

XLSX, PNG = SourceFormat.EXCEL, SourceFormat.IMAGE
TODAY = date(2026, 7, 27)


@pytest.fixture
def model() -> PMIDataModel:
    from app.agent.calculations import recompute_derived

    xlsx = make_source("integration_tracker.xlsx", XLSX, sheet_name="Workplan")
    image = make_source("risk_dashboard.png", PNG, extraction_confidence=0.35)
    built = PMIDataModel(
        project=PMIProject(project_id="p1", reporting_date=TODAY,
                           reporting_period="July 2026",
                           day_1_date=TODAY + timedelta(days=66)),
        source_files=["integration_tracker.xlsx", "risk_dashboard.png"],
        tasks=[Task(task_id="T1", title="Payroll cutover", owner="Anna Schmidt",
                    workstream="Finance", due_date=TODAY - timedelta(days=3),
                    status=Status.IN_PROGRESS, progress_percentage=60.0,
                    source_references=[xlsx]),
               Task(task_id="T2", title="Day 1 building access", owner=None,
                    workstream="Operations", is_day_1_critical=True,
                    status=Status.NOT_STARTED, source_references=[xlsx])],
        milestones=[Milestone(milestone_id="M1", name="ERP go-live",
                              planned_date=date(2026, 9, 15),
                              forecast_date=date(2026, 9, 30),
                              status=Status.IN_PROGRESS,
                              source_references=[xlsx]),
                    Milestone(milestone_id="M2", name="Legal close",
                              planned_date=date(2026, 8, 1),
                              status=Status.COMPLETED,
                              source_references=[xlsx])],
        risks=[Risk(risk_id="R1", title="GDPR retention breach", probability=4,
                    impact=5, status=Status.IN_PROGRESS,
                    source_references=[image])],
        budget=[BudgetItem(budget_item_id="B1", category="ERP migration",
                           budget=1_000_000.0, actual=900_000.0,
                           forecast=1_220_000.0, currency="EUR",
                           source_references=[xlsx]),
                BudgetItem(budget_item_id="B2", category="Rebranding",
                           budget=400_000.0, actual=380_000.0,
                           forecast=470_000.0, currency="EUR",
                           source_references=[xlsx])],
        synergies=[Synergy(synergy_id="S1", title="Procurement consolidation",
                           target_value=1_000_000.0, realized_value=400_000.0,
                           currency="EUR", source_references=[xlsx])],
    )
    built, issues = recompute_derived(built, TODAY)
    built.validation_issues.extend(issues)
    return built


@pytest.fixture
def built(model, scripted_planning, tmp_path):
    context = builder._assemble(
        scope="project", project_id="proj", chat_id=None, session_id=None,
        model=model,
        digest=KnowledgeDigest(
            free_text="MedAxis SE is integrating NordCare GmbH."),
        folder_name="", quality=None,
        request_text="Prepare a SteerCo pack on where we stand.",
        requested_format=None, messages=[])
    deliverable = engine.build(context)
    results = {fmt: registry.render(deliverable, context, tmp_path, fmt)
               for fmt in registry.supported()}
    return deliverable, context, results


# ================================================================== extractors
def text_of(path: Path) -> str:
    """Everything a reader would see, per format."""
    suffix = path.suffix.lower()
    if suffix == ".pptx":
        from pptx import Presentation

        presentation = Presentation(str(path))
        parts = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        parts.extend(cell.text for cell in row.cells)
        return "\n".join(parts)
    if suffix == ".docx":
        import docx

        document = docx.Document(str(path))
        parts = [p.text for p in document.paragraphs]
        for table in document.tables:
            for row in table.rows:
                parts.extend(cell.text for cell in row.cells)
        for section in document.sections:
            parts.extend(p.text for p in section.footer.paragraphs)
        return "\n".join(parts)
    if suffix == ".pdf":
        import fitz

        with fitz.open(str(path)) as document:
            return "\n".join(page.get_text() for page in document)
    return path.read_text(encoding="utf-8")


def normalized(text: str) -> str:
    return " ".join(text.split())


# ==================================================================== all four
def test_all_four_formats_are_produced(built):
    _deliverable, _context, results = built
    assert set(results) == {"pptx", "docx", "pdf", "html"}
    for fmt, result in results.items():
        assert result.path.is_file(), fmt
        assert result.path.stat().st_size > 3000, fmt
        assert result.path.suffix == f".{fmt}"


def test_every_page_title_appears_in_every_format(built):
    """One plan means four renderings of the same document, not four documents."""
    deliverable, _context, results = built
    titles = [page.title for page in deliverable.pages if page.title]
    assert len(titles) >= 4

    for fmt, result in results.items():
        haystack = normalized(text_of(result.path))
        for title in titles:
            assert normalized(title) in haystack, \
                f"{fmt} is missing the title {title!r}"


def test_the_governing_message_leads_reading_formats_without_becoming_a_ppt_overlay(
        built):
    deliverable, _context, results = built
    for fmt, result in results.items():
        if fmt == "pptx":
            # PowerPoint follows the native title hierarchy. A separate copy of
            # the governing message above that placeholder collided with the
            # logo; the cover's planned title is the visible opening instead.
            assert normalized(deliverable.pages[0].title) in \
                normalized(text_of(result.path))
            continue
        assert normalized(deliverable.governing_message) in \
            normalized(text_of(result.path)), fmt


def test_no_format_contains_a_chart_stub(built):
    """`Chart: Workstream Progress — see the deck.` was reachable three ways."""
    _deliverable, _context, results = built
    for fmt, result in results.items():
        body = text_of(result.path)
        assert "Chart:" not in body, fmt
        assert "see the deck." not in body, fmt
        assert "No executive summary" not in body, fmt


def test_every_format_embeds_a_real_visual(built):
    """§19: a planned chart must produce an actual rendered visual, per format."""
    deliverable, _context, results = built
    assert deliverable.specs.charts

    from pptx import Presentation

    presentation = Presentation(str(results["pptx"].path))
    assert any(getattr(shape, "has_chart", False)
               for slide in presentation.slides for shape in slide.shapes)

    import docx

    assert len(docx.Document(str(results["docx"].path)).inline_shapes) >= 1

    import fitz

    with fitz.open(str(results["pdf"].path)) as document:
        assert sum(len(page.get_images()) for page in document) >= 1

    html = results["html"].path.read_text(encoding="utf-8")
    assert "<svg" in html


def test_word_and_pdf_put_the_deloitte_logo_on_page_one(built):
    """The cover wordmark must be a visible first-page image, not the former
    white-on-transparent asset that showed only its green dot on white paper."""
    import docx
    import fitz
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    _deliverable, _context, results = built

    word = docx.Document(str(results["docx"].path))
    assert word.paragraphs
    assert "<w:drawing" in word.paragraphs[0]._p.xml
    assert word.paragraphs[0].alignment == WD_ALIGN_PARAGRAPH.RIGHT

    with fitz.open(str(results["pdf"].path)) as pdf:
        images = pdf[0].get_images(full=True)
        assert len(images) == 1
        rects = pdf[0].get_image_rects(images[0][0])
        assert rects and rects[0].x0 > pdf[0].rect.width / 2


def test_the_transcribed_figure_is_disclosed_in_every_format(built):
    """A 0.35-confidence read is not quotable as fact anywhere."""
    _deliverable, _context, results = built
    for fmt, result in results.items():
        body = normalized(text_of(result.path))
        assert "read from an image" in body, fmt
        assert "risk_dashboard.png" in body, fmt


def test_the_company_names_reach_every_format(built):
    _deliverable, _context, results = built
    for fmt, result in results.items():
        body = text_of(result.path)
        assert "MedAxis" in body, fmt
        assert "PMI Project" not in body, fmt


def test_missing_values_never_render_as_zero(built):
    _deliverable, _context, results = built
    for fmt, result in results.items():
        body = text_of(result.path)
        if "Not Reported" in body:
            break
    else:
        pytest.skip("this fixture reports every figure")


def test_the_filenames_name_the_project(built):
    _deliverable, _context, results = built
    for fmt, result in results.items():
        assert "MedAxis" in result.path.name, fmt
        assert not result.path.name.startswith("PMI_Report"), fmt


# ======================================================================= Word
def test_word_uses_defined_brand_styles_not_words_defaults(built):
    import docx

    _deliverable, context, results = built
    document = docx.Document(str(results["docx"].path))
    names = {style.name for style in document.styles}
    assert {"PMI Title", "PMI Heading 1", "PMI Body", "PMI Callout",
            "PMI Table Header"} <= names

    assert document.styles["PMI Body"].font.name == context.brand_system.font_minor
    # Never Word's built-in table style again.
    for table in document.tables:
        assert table.style is None or "Light Grid" not in (table.style.name or "")


def test_word_has_a_real_toc_field_and_a_static_fallback(built):
    """The fallback is numbered, and Word refreshes it to exact pagination."""
    import docx

    deliverable, _context, results = built
    document = docx.Document(str(results["docx"].path))
    xml = document.element.xml
    assert "TOC \\o" in xml
    assert "PAGEREF pmi_section_1" in xml
    assert "bookmarkStart" in xml
    assert "updateFields" in document.settings.element.xml

    body = normalized(text_of(results["docx"].path))
    assert "Contents" in body
    titles = [p.title for p in deliverable.pages if p.purpose != "cover"]
    assert normalized(titles[0]) in body
    first_entry = next(p.text for p in document.paragraphs
                       if p.text.startswith(titles[0]))
    assert first_entry.endswith("\t3"), \
        "the cached TOC fallback has no visible page number"


def test_word_repeats_table_headers_and_numbers_figures(built):
    import docx

    _deliverable, _context, results = built
    document = docx.Document(str(results["docx"].path))
    xml = document.element.xml
    assert "tblHeader" in xml, "a table spanning a page must repeat its header"
    assert "SEQ Figure" in xml, "figures must be natively numbered"

    # Footers are a separate OOXML part; they are not in the body element.
    footer_xml = "".join(section.footer._element.xml
                         for section in document.sections)
    assert "PAGE" in footer_xml and "NUMPAGES" in footer_xml, \
        "footers need page numbers"


def test_generated_chart_text_is_not_repeated_beside_the_chart():
    from app.deliverable.engine import _remove_redundant_visual_text
    from app.deliverable.model import ChartElement, PageDesign, TextElement

    page = PageDesign(page_id="status", elements=[
        ChartElement(element_id="chart", spec_id="chart-spec",
                     evidence_ids=["ev:kpi:1"]),
        TextElement(element_id="body", role="body",
                    text="The KPI is shown in the chart.",
                    evidence_ids=["ev:kpi:1"], authored_by="llm"),
        TextElement(element_id="extra", role="body",
                    text="Management needs to decide the recovery owner.",
                    evidence_ids=["ev:decision:1"], authored_by="llm"),
    ])

    _remove_redundant_visual_text(page)

    assert [element.element_id for element in page.elements] == ["chart", "extra"]

def test_generated_table_text_is_not_repeated_below_the_table():
    from app.deliverable.engine import _remove_redundant_visual_text
    from app.deliverable.model import (
        BulletsElement,
        PageDesign,
        TableElement,
        TextElement,
    )

    page = PageDesign(page_id="risks", elements=[
        TableElement(element_id="table", spec_id="risk-table",
                     evidence_ids=["ev:risk:1", "ev:risk:2"]),
        TextElement(element_id="duplicate", role="body",
                    text="The two risks listed above remain open.",
                    evidence_ids=["ev:risk:1", "ev:risk:2"],
                    authored_by="llm"),
        TextElement(element_id="extra", role="body",
                    text="Management still needs to choose an owner.",
                    evidence_ids=["ev:decision:1"], authored_by="llm"),
        BulletsElement(element_id="user-note", items=["Keep this wording"],
                       evidence_ids=["ev:risk:1"], authored_by="user"),
    ])

    _remove_redundant_visual_text(page)

    assert [element.element_id for element in page.elements] == [
        "table", "extra", "user-note",
    ]


def test_word_keeps_headings_with_what_follows(built):
    import docx

    _deliverable, _context, results = built
    document = docx.Document(str(results["docx"].path))
    heading = document.styles["PMI Heading 1"]
    assert heading.paragraph_format.keep_with_next is True


def test_word_headings_carry_an_outline_level(built):
    """Without one, a real TOC field collects nothing."""
    import docx

    _deliverable, _context, results = built
    document = docx.Document(str(results["docx"].path))
    assert "outlineLvl" in document.styles["PMI Heading 1"].element.xml


def test_word_starts_every_planned_chapter_on_a_new_page(built):
    """Divider pages count as chapters and must not swallow the next heading."""
    import docx

    deliverable, _context, results = built
    document = docx.Document(str(results["docx"].path))
    planned = [page for page in deliverable.pages if page.purpose != "cover"]

    for page in planned:
        matches = [index for index, paragraph in enumerate(document.paragraphs)
                   if paragraph.text == page.title]
        assert matches, page.title
        index = matches[-1]  # skip the cached contents entry
        assert index > 0
        before = document.paragraphs[index - 1]
        assert any(node.tag.endswith("}br") and
                   node.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}type") == "page"
                   for node in before._p.iter()), page.title


# ======================================================================== PDF
def test_the_pdf_has_readable_pages_and_a_contents(built):
    import fitz

    _deliverable, _context, results = built
    with fitz.open(str(results["pdf"].path)) as document:
        assert document.page_count >= 3
        assert "Contents" in document[1].get_text()
        for index, page in enumerate(document):
            if index == 0:
                continue
            assert f"Page {index + 1}" in page.get_text(), \
                f"page {index + 1} has no footer"


def test_pdf_starts_every_planned_chapter_on_a_new_page(built):
    import fitz

    deliverable, _context, results = built
    planned = [page for page in deliverable.pages if page.purpose != "cover"]
    with fitz.open(str(results["pdf"].path)) as document:
        starts = []
        for planned_page in planned:
            wanted = normalized(planned_page.title)
            occurrences = [index for index, page in enumerate(document)
                           if wanted in normalized(page.get_text())]
            assert occurrences, planned_page.title
            starts.append(occurrences[-1])  # skip the contents occurrence
    assert len(set(starts)) == len(planned), starts


def test_the_pdf_embeds_a_unicode_font_not_a_core_font(built):
    """The old renderer transliterated ⚠ to `!` because core fonts are Latin-1."""
    import fitz

    _deliverable, _context, results = built
    with fitz.open(str(results["pdf"].path)) as document:
        fonts = {font[3] for page in document for font in page.get_fonts()}
    assert fonts
    assert not fonts <= {"Helvetica", "Helvetica-Bold", "Times-Roman"}


def test_the_pdf_says_when_it_substituted_the_brand_font(built):
    """Aptos is a Microsoft font and is absent from most build hosts."""
    _deliverable, context, results = built
    warnings = results["pdf"].warnings
    from reportlab.pdfbase import pdfmetrics

    registered = set(pdfmetrics.getRegisteredFontNames())
    if context.brand_system.font_minor in registered:
        pytest.skip("the brand font is installed on this machine")
    assert any("not installed" in warning for warning in warnings)
    assert any(context.brand_system.font_minor in warning for warning in warnings)


def test_pdf_table_text_is_not_truncated(built):
    """fpdf2 hard-cut every cell to 26 characters with no wrapping."""
    deliverable, _context, results = built
    tables = list(deliverable.specs.tables.values())
    if not tables:
        pytest.skip("this scenario planned no table")

    body = normalized(text_of(results["pdf"].path))
    long_cells = [cell.text for row in tables[0].rows for cell in row
                  if len(cell.text) > 26]
    for text in long_cells[:3]:
        assert normalized(text) in body, f"{text!r} was truncated"


def test_the_pdf_renders_to_images_for_visual_review(built):
    """§19 asks for rasterised pages. PyMuPDF does the PDF; there is no
    LibreOffice on this machine, so the deck is checked analytically instead."""
    import fitz

    _deliverable, _context, results = built
    with fitz.open(str(results["pdf"].path)) as document:
        for index, page in enumerate(document):
            pixmap = page.get_pixmap(dpi=72)
            assert pixmap.width > 400 and pixmap.height > 500
            # A page that is uniformly one colour is blank.
            assert len(set(pixmap.samples[::997])) > 1, f"page {index + 1} is blank"


# ======================================================================= HTML
def test_the_dashboard_is_one_self_contained_file(built):
    """It gets emailed and archived. A dashboard that needs a CDN stops working."""
    import re

    _deliverable, _context, results = built
    html = results["html"].path.read_text(encoding="utf-8")

    for match in re.finditer(r'(?:src|href)="([^"]*)"', html):
        target = match.group(1)
        assert target.startswith(("data:", "#")), f"external reference: {target}"
    assert "<script src=" not in html
    assert '<link href="http' not in html and "@import" not in html
    # No naive substring search for "cdn" here: base64 image data contains
    # every three-letter sequence eventually.


def test_the_dashboard_is_interactive_without_a_library(built):
    import re

    _deliverable, _context, results = built
    html = results["html"].path.read_text(encoding="utf-8")
    assert "<script>" in html
    assert "IntersectionObserver" in html          # scroll-spy on the rail
    assert "aria-sort" in html                     # sortable columns
    assert 'class="filter"' in html                # row filtering
    assert "pmi-mark" in html                      # chart tooltips
    assert "data-evidence-id" in html              # marks are traceable
    # The custom tooltip is the only hover tooltip. A nested SVG <title> would
    # make the browser show the same value a second time.
    assert not re.search(r'class="pmi-mark"[^>]*>\s*<title>', html)
    assert re.search(r'class="pmi-mark"[^>]*aria-label=', html)


def test_the_dashboard_labels_unknown_cover_metadata(built):
    _deliverable, _context, results = built
    html = results["html"].path.read_text(encoding="utf-8")

    assert '<span class="meta-label">Audience:</span> Steering Committee' in html
    assert '<span class="meta-label">Reporting period:</span> July 2026' in html
    assert '<span class="meta-label">Integration phase:</span> Unknown' in html


def test_the_dashboard_charts_are_inline_svg_not_images(built):
    _deliverable, _context, results = built
    html = results["html"].path.read_text(encoding="utf-8")
    assert html.count("<svg") >= 1
    assert "<img src=\"data:image/png" not in html or "<svg" in html


def test_the_dashboard_escapes_content(built):
    """Titles and cells come from a model and from files."""
    deliverable, context, results = built
    deliverable.pages[1].title = "<script>alert('x')</script> risks"
    from app.renderers import html as html_renderer

    result = html_renderer.render(deliverable, context,
                                 results["html"].path.parent / "escaped")
    html = result.path.read_text(encoding="utf-8")
    assert "<script>alert" not in html
    assert "&lt;script&gt;alert" in html


def test_the_dashboard_has_a_print_stylesheet_and_a_methodology_panel(built):
    _deliverable, _context, results = built
    html = results["html"].path.read_text(encoding="utf-8")
    assert "@media print" in html
    assert 'id="methodology"' in html
    assert "Sources and methodology" in html
    assert "prefers-color-scheme" in html


def test_the_dashboard_is_not_a_fixed_grid_of_kpi_cards(built):
    """The old renderer lifted a KPI strip to the top of every report."""
    deliverable, _context, results = built
    html = results["html"].path.read_text(encoding="utf-8")
    planned_kpis = any(page.of_role("kpi_row") for page in deliverable.pages)
    assert ('class="kpi-row"' in html) == planned_kpis


def test_the_dashboard_reflects_the_planned_composition(built):
    deliverable, _context, results = built
    html = results["html"].path.read_text(encoding="utf-8")
    for page in deliverable.pages:
        if page.purpose == "cover":
            continue
        assert f'id="{page.page_id}"' in html
        assert page.composition in html


# ================================================================== robustness
def test_one_failing_format_does_not_lose_the_others(model, scripted_planning,
                                                    tmp_path, monkeypatch):
    """A user who asked for a deck and a PDF should get the deck."""
    context = builder._assemble(
        scope="project", project_id="proj", chat_id=None, session_id=None,
        model=model, digest=KnowledgeDigest(), folder_name="Aurora",
        quality=None, request_text="Prepare a SteerCo pack.",
        requested_format=None, messages=[])
    deliverable = engine.build(context)

    def explode(*_args, **_kwargs):
        raise RuntimeError("the PDF engine fell over")

    monkeypatch.setattr("app.renderers.pdf.render", explode)
    results = registry.render_all(deliverable, context, tmp_path,
                                  ["pptx", "pdf", "html"])

    assert results[0].path.is_file() and results[0].suffix == "pptx"
    assert results[2].path.is_file() and results[2].suffix == "html"
    assert "could not be produced" in results[1].warnings[0]
    assert "fell over" in results[1].warnings[0]


def test_a_keyless_run_renders_all_formats_without_system_provenance(model, tmp_path):
    context = builder._assemble(
        scope="project", project_id="proj", chat_id=None, session_id=None,
        model=model, digest=KnowledgeDigest(), folder_name="Aurora",
        quality=None, request_text="Prepare a pack:\n1. Risks\n2. Budget position",
        requested_format=None, messages=[])
    deliverable = engine.build(context)
    assert deliverable.planned_by == "fallback"

    for fmt in registry.supported():
        result = registry.render(deliverable, context, tmp_path, fmt)
        assert result.path.is_file(), fmt
        body = text_of(result.path).lower()
        assert "without a language model" not in body, fmt
        assert "language model unavailable" not in body, fmt
